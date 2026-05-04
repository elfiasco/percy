"""PPTX inspection for Bridge onboarding diagnostics."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation

from percy.diagnostics.common import emu_to_inches, ensure_dir, enum_name, length_to_points, safe_get, write_json
from percy.diagnostics.inheritance import placeholder_info, resolve_text_shape


def inspect_pptx(pptx_path: str | Path, out_dir: str | Path | None = None) -> dict[str, Any]:
    """Return and optionally write a detailed PPTX inspection report."""

    path = Path(pptx_path)
    presentation = Presentation(str(path))
    report: dict[str, Any] = {
        "source_path": str(path),
        "slide_width": emu_to_inches(presentation.slide_width),
        "slide_height": emu_to_inches(presentation.slide_height),
        "slide_count": len(presentation.slides),
        "slides": [],
        "diagnostics": [],
    }

    for slide_index, slide in enumerate(presentation.slides, start=1):
        shapes = []
        for shape in slide.shapes:
            shapes.extend(_inspect_shape_tree(shape, slide_index))
        slide_report = {
            "slide_number": slide_index,
            "shape_count": len(shapes),
            "top_level_shape_count": len(slide.shapes),
            "shapes": shapes,
        }
        report["slides"].append(slide_report)

    report["diagnostics"] = _summarize_gaps(report)

    if out_dir is not None:
        output_dir = ensure_dir(out_dir)
        write_json(report, output_dir / "inspection.json")
        for slide in report["slides"]:
            write_json(slide, output_dir / f"slide-{slide['slide_number']:03d}.json")

    return report


def _inspect_shape_tree(shape: Any, slide_number: int, group_id: str | None = None, group_path: list[str] | None = None) -> list[dict[str, Any]]:
    if enum_name(safe_get(lambda: shape.shape_type)) == "GROUP" and safe_get(lambda: shape.shapes, None) is not None:
        current_group_id = group_id or f"slide-{slide_number}:group-{safe_get(lambda: shape.shape_id)}"
        current_path = [*(group_path or []), current_group_id]
        group_report = _inspect_shape(shape, slide_number, group_id, group_path or [])
        group_report["group"] = {
            "is_group_container": True,
            "group_id": current_group_id,
            "group_path": current_path,
            "child_count": len(shape.shapes),
        }
        children = []
        for child in shape.shapes:
            children.extend(_inspect_shape_tree(child, slide_number, current_group_id, current_path))
        return [group_report, *children]
    return [_inspect_shape(shape, slide_number, group_id, group_path or [])]


def _inspect_shape(shape: Any, slide_number: int, group_id: str | None = None, group_path: list[str] | None = None) -> dict[str, Any]:
    placeholder = placeholder_info(shape)
    shape_report: dict[str, Any] = {
        "slide_number": slide_number,
        "shape_id": safe_get(lambda: shape.shape_id),
        "shape_name": safe_get(lambda: shape.name),
        "shape_type": enum_name(safe_get(lambda: shape.shape_type)),
        "position": {
            "left": emu_to_inches(safe_get(lambda: shape.left)),
            "top": emu_to_inches(safe_get(lambda: shape.top)),
            "width": emu_to_inches(safe_get(lambda: shape.width)),
            "height": emu_to_inches(safe_get(lambda: shape.height)),
        },
        "rotation": safe_get(lambda: shape.rotation, 0.0),
        "has_text_frame": bool(safe_get(lambda: shape.has_text_frame, False)),
        "has_table": bool(safe_get(lambda: shape.has_table, False)),
        "has_chart": bool(safe_get(lambda: shape.has_chart, False)),
        "has_picture": _looks_like_picture(shape),
        "placeholder": placeholder,
        "group": {
            "is_group_container": False,
            "group_id": group_id,
            "group_path": group_path or [],
        },
        "semantic_role": _semantic_role(shape, placeholder),
        "fill": _inspect_fill(shape),
        "line": _inspect_line(shape),
        "raw_xml": safe_get(lambda: shape.element.xml),
    }

    if shape_report["has_text_frame"]:
        shape_report["text"] = _inspect_text(shape)
    if shape_report["has_table"]:
        shape_report["table"] = _inspect_table(shape)
    if shape_report["has_chart"]:
        shape_report["chart"] = _inspect_chart(shape)
    if shape_report["has_picture"]:
        shape_report["picture"] = _inspect_picture(shape)

    return shape_report


def _inspect_text(shape: Any) -> dict[str, Any]:
    text_frame = shape.text_frame
    return {
        "text": safe_get(lambda: shape.text),
        "word_wrap": safe_get(lambda: text_frame.word_wrap),
        "auto_size": enum_name(safe_get(lambda: text_frame.auto_size)),
        "vertical_anchor": enum_name(safe_get(lambda: text_frame.vertical_anchor)),
        "paragraphs": [_inspect_paragraph(paragraph, shape) for paragraph in text_frame.paragraphs],
    }


def _inspect_paragraph(paragraph: Any, shape: Any) -> dict[str, Any]:
    return {
        "alignment": enum_name(safe_get(lambda: paragraph.alignment)),
        "level": safe_get(lambda: paragraph.level),
        "line_spacing": safe_get(lambda: paragraph.line_spacing),
        "space_before": safe_get(lambda: paragraph.space_before),
        "space_after": safe_get(lambda: paragraph.space_after),
        "runs": [_inspect_run(run, paragraph, shape) for run in paragraph.runs],
    }


def _inspect_run(run: Any, paragraph: Any | None = None, shape: Any | None = None) -> dict[str, Any]:
    font = run.font
    resolved = resolve_text_shape(shape, paragraph, run) if shape is not None else {}
    return {
        "text": run.text,
        "font_name": safe_get(lambda: font.name),
        "font_size": length_to_points(safe_get(lambda: font.size)),
        "font_bold": safe_get(lambda: font.bold),
        "font_italic": safe_get(lambda: font.italic),
        "font_underline": safe_get(lambda: font.underline),
        "font_color": _inspect_color(safe_get(lambda: font.color)),
        "resolved": resolved,
    }


def _inspect_table(shape: Any) -> dict[str, Any]:
    table = shape.table
    rows = []
    for row in table.rows:
        rows.append([cell.text for cell in row.cells])
    return {
        "row_count": len(table.rows),
        "column_count": len(table.columns),
        "data": rows,
    }


def _inspect_chart(shape: Any) -> dict[str, Any]:
    chart = shape.chart
    series_count = safe_get(lambda: len(chart.series))
    return {
        "chart_type": enum_name(safe_get(lambda: chart.chart_type)),
        "has_title": safe_get(lambda: chart.has_title),
        "chart_style": safe_get(lambda: chart.chart_style),
        "series_count": series_count,
        "series_error": None if series_count is not None else "Could not iterate chart series with python-pptx.",
        "raw_chart_xml": safe_get(lambda: chart.part.blob.decode("utf-8", errors="replace")),
    }


def _inspect_picture(shape: Any) -> dict[str, Any]:
    image = safe_get(lambda: shape.image)
    return {
        "filename": safe_get(lambda: image.filename),
        "content_type": safe_get(lambda: image.content_type),
        "extension": safe_get(lambda: image.ext),
        "width_px": safe_get(lambda: image.size[0]),
        "height_px": safe_get(lambda: image.size[1]),
        "dpi": safe_get(lambda: image.dpi),
    }


def _inspect_fill(shape: Any) -> dict[str, Any]:
    fill = safe_get(lambda: shape.fill)
    if fill is None:
        return {}
    return {
        "type": enum_name(safe_get(lambda: fill.type)),
        "fore_color": _inspect_color(safe_get(lambda: fill.fore_color)),
        "transparency": safe_get(lambda: fill.transparency),
    }


def _inspect_line(shape: Any) -> dict[str, Any]:
    line = safe_get(lambda: shape.line)
    if line is None:
        return {}
    return {
        "color": _inspect_color(safe_get(lambda: line.color)),
        "width": safe_get(lambda: line.width),
        "dash_style": enum_name(safe_get(lambda: line.dash_style)),
    }


def _inspect_color(color: Any) -> str | None:
    if color is None:
        return None
    rgb = safe_get(lambda: color.rgb)
    if rgb is not None:
        return str(rgb)
    theme_color = safe_get(lambda: color.theme_color)
    if theme_color is not None:
        return f"theme:{enum_name(theme_color)}"
    return None


def _looks_like_picture(shape: Any) -> bool:
    return safe_get(lambda: shape.image, None) is not None


def _summarize_gaps(report: dict[str, Any]) -> list[dict[str, Any]]:
    gaps: list[dict[str, Any]] = []
    for slide in report["slides"]:
        for shape in slide["shapes"]:
            if shape["has_chart"]:
                gaps.append(_gap(shape, "chart_semantics", "Chart data/formatting is captured mostly as XML."))
            if shape.get("group", {}).get("is_group_container"):
                gaps.append(_gap(shape, "group_container", "Group container will not become a Bridge element; children are onboarded individually with group_id lineage."))
            if shape["shape_type"] not in {"TEXT_BOX", "PICTURE", "TABLE", "CHART"}:
                gaps.append(_gap(shape, "shape_type", f"Shape type may rebuild as a generic rectangle: {shape['shape_type']}"))
            if shape.get("rotation"):
                gaps.append(_gap(shape, "rotation", "Rotation is inspected but not fully rebuilt yet."))
            for paragraph in shape.get("text", {}).get("paragraphs", []):
                for run in paragraph.get("runs", []):
                    unresolved = run.get("resolved", {}).get("unresolved", [])
                    if unresolved:
                        gaps.append(
                            _gap(
                                shape,
                                "unresolved_inheritance",
                                f"Could not synthesize inherited text properties: {', '.join(unresolved)}",
                            )
                        )
    return gaps


def _gap(shape: dict[str, Any], code: str, message: str) -> dict[str, Any]:
    return {
        "slide_number": shape["slide_number"],
        "shape_id": shape["shape_id"],
        "shape_name": shape["shape_name"],
        "code": code,
        "message": message,
    }


def _semantic_role(shape: Any, placeholder: dict[str, Any]) -> str | None:
    name = (safe_get(lambda: shape.name, "") or "").lower()
    placeholder_type = placeholder.get("type")
    if placeholder_type == "SLIDE_NUMBER" or "slide number" in name or "pagenumber" in name or "page number" in name:
        return "page_number"
    if placeholder_type in {"FOOTER", "DATE"} or "footer" in name or "date" in name:
        return "footer"
    return None
