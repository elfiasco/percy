"""Bridge to PPTX best-effort rebuilding."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import BubbleChartData, CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_LINE, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.shapes.autoshape import AutoShapeType
from pptx.util import Inches, Pt

from percy.bridge import (
    BridgeChart,
    BridgeConnector,
    BridgeFreeform,
    BridgeGroup,
    BridgeImage,
    BridgeShape,
    BridgeTable,
    BridgeText,
    PercyDocument,
)
from percy.diagnostics.common import safe_get


_THEME_COLOR_SLOT_MAP = {
    "DARK_1":    "dk1",
    "DARK_2":    "dk2",
    "LIGHT_1":   "lt1",
    "LIGHT_2":   "lt2",
    "ACCENT_1":  "accent1",
    "ACCENT_2":  "accent2",
    "ACCENT_3":  "accent3",
    "ACCENT_4":  "accent4",
    "ACCENT_5":  "accent5",
    "ACCENT_6":  "accent6",
    "HYPERLINK": "hlink",
    "FOLLOWED_HYPERLINK": "folHlink",
}


def _patch_presentation_theme(presentation: Any, theme_colors: dict[str, str]) -> None:
    """Overwrite the theme color slots in the blank presentation's XML so that
    scheme color references resolve to the original document's palette."""
    if not theme_colors:
        return
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        theme_part = presentation.slide_masters[0].part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        theme_xml = etree.fromstring(theme_part.blob)
        clr_scheme = theme_xml.find(".//" + qn("a:clrScheme"))
        if clr_scheme is None:
            return
        for bridge_key, slot_name in _THEME_COLOR_SLOT_MAP.items():
            hex_val = (theme_colors.get(bridge_key) or "").lstrip("#")
            if len(hex_val) != 6:
                continue
            slot_el = clr_scheme.find(qn("a:" + slot_name))
            if slot_el is None:
                continue
            # Slot contains either <a:srgbClr> or <a:sysClr>; replace with srgbClr
            for child in list(slot_el):
                slot_el.remove(child)
            srgb = etree.SubElement(slot_el, qn("a:srgbClr"))
            srgb.set("val", hex_val.upper())
        theme_part._blob = etree.tostring(theme_xml, xml_declaration=True,
                                           encoding="UTF-8", standalone=True)
    except Exception:
        pass


def rebuild_pptx(document: PercyDocument, out_path: str | Path) -> dict[str, Any]:
    output_path = Path(out_path)
    presentation = Presentation()
    slide_width = document.metadata.slide_width or document.custom_properties.get("slide_width")
    slide_height = document.metadata.slide_height or document.custom_properties.get("slide_height")
    if slide_width:
        presentation.slide_width = Inches(float(slide_width))
    if slide_height:
        presentation.slide_height = Inches(float(slide_height))

    blank_layout = presentation.slide_layouts[6]
    diagnostics: list[dict[str, Any]] = []
    while len(presentation.slides) < len(document.slides):
        presentation.slides.add_slide(blank_layout)

    theme_colors = document.theme_colors or {}
    _patch_presentation_theme(presentation, theme_colors)

    for bridge_slide, pptx_slide in zip(document.slides, presentation.slides):
        _apply_slide_background(pptx_slide, bridge_slide, theme_colors)
        _apply_slide_notes(pptx_slide, bridge_slide)
        slide_default_text_color = getattr(bridge_slide, "default_text_color", None)
        for element in bridge_slide.elements:
            diagnostics.extend(_add_element(pptx_slide, element, theme_colors, slide_default_text_color))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output_path))
    return {"pptx_path": str(output_path), "diagnostics": diagnostics}


def _apply_slide_background(pptx_slide: Any, bridge_slide: Any, theme_colors: dict[str, str]) -> None:
    grad_stops = getattr(bridge_slide, "background_gradient_stops", None) or []
    grad_angle = getattr(bridge_slide, "background_gradient_angle", 0.0) or 0.0
    bg_color = bridge_slide.background_color

    if grad_stops and len(grad_stops) >= 2:
        try:
            fill = pptx_slide.background.fill
            _apply_gradient_fill(fill, grad_stops, grad_angle, theme_colors)
            return
        except Exception:
            pass

    if not bg_color:
        return
    try:
        fill = pptx_slide.background.fill
        fill.solid()
        _set_color(fill.fore_color, bg_color, theme_colors)
    except Exception:
        pass


def _apply_slide_notes(pptx_slide: Any, bridge_slide: Any) -> None:
    """Write speaker notes text to the pptx slide notes_slide if present."""
    cp = getattr(bridge_slide, "custom_properties", None) or {}
    notes_text: str = cp.get("notes_text", "")
    if not notes_text:
        return
    try:
        notes_slide = pptx_slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text
    except Exception:
        pass


def _add_element(slide: Any, element: Any, theme_colors: dict[str, str] | None = None, default_text_color: str | None = None) -> list[dict[str, Any]]:
    tc = theme_colors or {}
    if isinstance(element, BridgeText):
        _add_text(slide, element, tc, default_text_color)
        return []
    if isinstance(element, BridgeImage):
        return _add_image(slide, element, tc)
    if isinstance(element, BridgeTable):
        _add_table(slide, element, tc, default_text_color)
        return []
    if isinstance(element, BridgeShape):
        return _add_shape(slide, element, tc, default_text_color)
    if isinstance(element, BridgeConnector):
        return _add_connector(slide, element, tc)
    if isinstance(element, BridgeChart):
        return _add_chart(slide, element, tc)
    if isinstance(element, BridgeFreeform):
        return _add_freeform(slide, element, tc, default_text_color)
    if isinstance(element, BridgeGroup):
        group_diags: list[dict[str, Any]] = []
        for child in element.children:
            group_diags.extend(_add_element(slide, child, tc, default_text_color))
        return group_diags
    return [_diag(element, "unsupported_element", f"Unsupported element type: {type(element).__name__}")]


def _apply_outer_shadow(shape: Any, shadow: Any, theme_colors: dict[str, str] | None = None) -> None:
    """Inject an <a:outerShdw> element into shape's spPr effectLst."""
    if shadow is None or not shadow.has_shadow:
        return
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        sp_pr = shape.element.spPr
        effect_lst = sp_pr.find(qn("a:effectLst"))
        if effect_lst is None:
            effect_lst = etree.SubElement(sp_pr, qn("a:effectLst"))
        # Remove any existing outerShdw
        existing = effect_lst.find(qn("a:outerShdw"))
        if existing is not None:
            effect_lst.remove(existing)
        shdw = etree.SubElement(effect_lst, qn("a:outerShdw"))
        if shadow.blur is not None:
            shdw.set("blurRad", str(int(shadow.blur * 12700)))
        if shadow.distance is not None:
            shdw.set("dist", str(int(shadow.distance * 12700)))
        if shadow.direction is not None:
            shdw.set("dir", str(int(shadow.direction * 60000)))
        if shadow.align:
            shdw.set("algn", shadow.align)
        shdw.set("rotWithShape", "1" if shadow.rot_with_shape else "0")
        if shadow.color:
            color_hex = (shadow.color.value or "").lstrip("#")
            if len(color_hex) == 6:
                solid = etree.SubElement(shdw, qn("a:srgbClr"))
                solid.set("val", color_hex.upper())
                if shadow.color.alpha is not None:
                    alpha_el = etree.SubElement(solid, qn("a:alpha"))
                    alpha_el.set("val", str(shadow.color.alpha))
    except Exception:
        pass


def _add_text(slide: Any, element: BridgeText, theme_colors: dict[str, str] | None = None, default_text_color: str | None = None) -> None:
    box = slide.shapes.add_textbox(*_box(element))
    _apply_transform(box, element)
    _apply_text_frame(box.text_frame, element.paragraphs, theme_colors, default_text_color, slide_part=safe_get(lambda: slide.part))
    _apply_text_box_settings(box.text_frame, element.text_frame)
    fab = element.fill_and_border
    if fab.has_fill and fab.fill_color:
        _apply_fill_color(box.fill, fab.fill_color, theme_colors)
    if fab.has_border:
        _apply_line(box.line, fab, theme_colors)
    _apply_outer_shadow(box, getattr(element, "shadow", None), theme_colors)


_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_BULLET_TAGS = ("buClrTx", "buClr", "buFontTx", "buFont", "buSzTx", "buSzPct", "buSzPts", "buNone", "buChar", "buAutoNum", "buBlip")


def _apply_paragraph_bullet(paragraph: Any, bridge_paragraph: Any, slide_part: Any = None) -> None:
    """Apply bullet character and paragraph indentation to a rebuilt paragraph."""
    from lxml import etree
    bullet_type = getattr(bridge_paragraph, "bullet_type", "none") or "none"
    bullet_char = getattr(bridge_paragraph, "bullet_char", None)
    left_indent = getattr(bridge_paragraph, "left_indent", None)
    first_line_indent = getattr(bridge_paragraph, "first_line_indent", None)

    if bullet_type not in ("char", "image") and left_indent is None and first_line_indent is None:
        return

    pPr = paragraph._p.get_or_add_pPr()

    # For bullet paragraphs, skip marL/indent=0 — PowerPoint applies correct default hanging
    # indent. Writing explicit 0 overrides that default and eliminates the hanging indent.
    effective_marL = left_indent if (left_indent and left_indent != 0.0) else None
    effective_indent = first_line_indent if (first_line_indent and first_line_indent != 0.0) else None
    if bullet_type not in ("char", "image"):
        effective_marL = left_indent
        effective_indent = first_line_indent

    if effective_marL is not None:
        pPr.set("marL", str(int(round(effective_marL * 914400))))
    if effective_indent is not None:
        pPr.set("indent", str(int(round(effective_indent * 914400))))

    if bullet_type == "char" and bullet_char:
        bullet_font = getattr(bridge_paragraph, "bullet_font", None)
        for tag in _BULLET_TAGS:
            for el in pPr.findall(f"{{{_A_NS}}}{tag}"):
                pPr.remove(el)
        if bullet_font:
            buFont = etree.SubElement(pPr, f"{{{_A_NS}}}buFont")
            buFont.set("typeface", bullet_font)
        buChar = etree.SubElement(pPr, f"{{{_A_NS}}}buChar")
        buChar.set("char", bullet_char)

    elif bullet_type == "image":
        blip_bytes = getattr(bridge_paragraph, "bullet_blip_bytes", None)
        blip_ext = getattr(bridge_paragraph, "bullet_blip_ext", None) or "png"
        if blip_bytes and slide_part is not None:
            try:
                import io
                from pptx.opc.constants import RELATIONSHIP_TYPE as RT
                _R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
                # Add image to slide part
                ct_map = {"png": "image/png", "jpeg": "image/jpeg", "jpg": "image/jpeg"}
                content_type = ct_map.get(blip_ext, "image/png")
                from pptx.opc.package import Part
                from pptx.opc.packuri import PackURI
                img_part = Part(
                    partname=PackURI(f"/ppt/media/buBlip{id(bridge_paragraph)}.{blip_ext}"),
                    content_type=content_type,
                    package=slide_part.package,
                    blob=blip_bytes,
                )
                rId = slide_part.relate_to(img_part, RT.IMAGE)
                # Clear existing bullet tags and inject buBlip
                for tag in _BULLET_TAGS:
                    for el in pPr.findall(f"{{{_A_NS}}}{tag}"):
                        pPr.remove(el)
                buFontTx = etree.SubElement(pPr, f"{{{_A_NS}}}buFontTx")
                buBlip = etree.SubElement(pPr, f"{{{_A_NS}}}buBlip")
                blip_el = etree.SubElement(buBlip, f"{{{_A_NS}}}blip")
                blip_el.set(f"{{{_R_NS}}}embed", rId)
            except Exception:
                pass


def _apply_text_frame(text_frame: Any, paragraphs: list[Any], theme_colors: dict[str, str] | None = None, default_text_color: str | None = None, slide_part: Any = None) -> None:
    text_frame.clear()
    for paragraph_index, bridge_paragraph in enumerate(paragraphs or []):
        paragraph = text_frame.paragraphs[0] if paragraph_index == 0 else text_frame.add_paragraph()
        paragraph.level = bridge_paragraph.indent_level
        alignment = _enum_member(PP_ALIGN, bridge_paragraph.alignment)
        if alignment is not None:
            paragraph.alignment = alignment
        if bridge_paragraph.line_spacing is not None:
            ls = bridge_paragraph.line_spacing
            try:
                if 0 < ls < 10.0:
                    paragraph.line_spacing = ls          # proportional multiplier
                elif ls >= 10.0:
                    paragraph.line_spacing = Pt(ls)      # fixed in points
            except Exception:
                pass
        if bridge_paragraph.space_before is not None:
            paragraph.space_before = Pt(bridge_paragraph.space_before)
        if bridge_paragraph.space_after is not None:
            paragraph.space_after = Pt(bridge_paragraph.space_after)
        _apply_paragraph_bullet(paragraph, bridge_paragraph, slide_part)
        if not bridge_paragraph.runs:
            paragraph.text = ""
            end_sz = getattr(bridge_paragraph, "end_para_font_size", None)
            if end_sz is not None:
                try:
                    from lxml import etree
                    _A = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    end_pr = paragraph._p.find(f"{{{_A}}}endParaRPr")
                    if end_pr is None:
                        end_pr = etree.SubElement(paragraph._p, f"{{{_A}}}endParaRPr")
                    end_pr.set("sz", str(int(round(end_sz * 100))))
                except Exception:
                    pass
        for run_data in bridge_paragraph.runs:
            if getattr(run_data, "is_line_break", False):
                from lxml import etree
                br = etree.SubElement(paragraph._p, "{http://schemas.openxmlformats.org/drawingml/2006/main}br")
                continue
            run = paragraph.add_run()
            run.text = run_data.text
            run.font.name = run_data.font_name
            if run_data.font_size is not None:
                run.font.size = Pt(run_data.font_size)
            bold = run_data.font_bold
            if bold is None:
                fn = run_data.font_name or ""
                fn_lower = fn.lower().replace("-", " ").replace("_", " ")
                if any(w in fn_lower for w in ("bold", "black", "heavy", "extrabold", "demibold", "semibold")):
                    bold = True
            run.font.bold = bold
            run.font.italic = run_data.font_italic
            run.font.underline = run_data.font_underline
            if run_data.font_color is not None:
                _set_color(run.font.color, run_data.font_color, theme_colors)
            elif default_text_color:
                _set_color(run.font.color, default_text_color, theme_colors)
            if getattr(run_data, "font_caps", None):
                _set_run_caps(run, run_data.font_caps)
            if getattr(run_data, "char_spacing", None) is not None:
                _set_run_rpr_attr(run, "spc", str(int(run_data.char_spacing)))
            if getattr(run_data, "baseline_shift", None) is not None:
                _set_run_rpr_attr(run, "baseline", str(int(run_data.baseline_shift)))
            strike = getattr(run_data, "strikethrough", None)
            if strike and strike != "noStrike":
                _set_run_rpr_attr(run, "strike", strike)
            if getattr(run_data, "hyperlink", None):
                _set_run_hyperlink(run, run_data.hyperlink)


def _add_image(slide: Any, element: BridgeImage, theme_colors: dict[str, str] | None = None) -> list[dict[str, Any]]:
    if not element.image_data.image_bytes:
        _add_generic_shape(slide, element)
        return [_diag(element, "image_missing_bytes", "Image had no bytes; rebuilt as generic rectangle.")]
    picture = slide.shapes.add_picture(BytesIO(element.image_data.image_bytes), *_box(element))
    _apply_transform(picture, element)
    picture.crop_left = element.cropping.crop_left
    picture.crop_right = element.cropping.crop_right
    picture.crop_top = element.cropping.crop_top
    picture.crop_bottom = element.cropping.crop_bottom
    _apply_line(picture.line, element.border, theme_colors)
    _apply_picture_geometry(picture, element)
    _apply_outer_shadow(picture, getattr(element, "shadow", None), theme_colors)
    return _unsupported_color_diags(element, [element.border.border_color])


def _apply_picture_geometry(picture: Any, element: BridgeImage) -> None:
    """Apply prstGeom to a picture element (for rounded-rect and other clip shapes)."""
    prst = getattr(element, "shape_geometry", None)
    if not prst:
        return
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        # <p:pic> uses <p:spPr> not <p:sp>/<p:spPr>
        sp_pr = picture.element.find(qn("p:spPr"))
        if sp_pr is None:
            return
        # Remove existing prstGeom (default rect)
        existing = sp_pr.find(qn("a:prstGeom"))
        if existing is not None:
            sp_pr.remove(existing)
        # Build new prstGeom with adjustments
        new_geom = etree.Element(qn("a:prstGeom"))
        new_geom.set("prst", prst)
        adj = getattr(element, "shape_geometry_adj", {}) or {}
        if adj:
            av_lst = etree.SubElement(new_geom, qn("a:avLst"))
            for name, fmla in adj.items():
                gd = etree.SubElement(av_lst, qn("a:gd"))
                gd.set("name", name)
                gd.set("fmla", fmla)
        else:
            etree.SubElement(new_geom, qn("a:avLst"))
        # Insert after xfrm if present, otherwise at start
        xfrm = sp_pr.find(qn("a:xfrm"))
        insert_pos = list(sp_pr).index(xfrm) + 1 if xfrm is not None else 0
        sp_pr.insert(insert_pos, new_geom)
    except Exception:
        pass


def _add_table(slide: Any, element: BridgeTable, theme_colors: dict[str, str] | None = None, default_text_color: str | None = None) -> None:
    rows = max(len(element.data), 1)
    cols = max(len(element.data[0]) if element.data else 1, 1)
    table_shape = slide.shapes.add_table(rows, cols, *_box(element))
    table = table_shape.table
    _apply_table_dimensions(table, element)
    _apply_table_properties(table, element)
    _apply_table_merges(table, element)
    for row_index, row in enumerate(element.data):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell_format = safe_get(lambda: element.cell_formats[row_index][col_index])
            if cell_format is not None and cell_format.merge.is_spanned:
                continue
            if cell_format is not None and cell_format.paragraphs:
                _apply_text_frame(cell.text_frame, cell_format.paragraphs, theme_colors, default_text_color, slide_part=safe_get(lambda: slide.part))
            else:
                cell.text = "" if value is None else str(value)
            if cell_format is not None:
                _apply_table_cell_format(cell, cell_format, theme_colors)


def _apply_table_dimensions(table: Any, element: BridgeTable) -> None:
    for index, width in enumerate(element.dimensions.column_widths):
        if index < len(table.columns) and width:
            table.columns[index].width = Inches(width)
    for index, height in enumerate(element.dimensions.row_heights):
        if index < len(table.rows) and height:
            table.rows[index].height = Inches(height)


def _apply_table_properties(table: Any, element: BridgeTable) -> None:
    props = element.table_properties
    table.first_row = props.first_row_header
    table.first_col = props.first_col_header
    table.last_row = props.last_row_total
    table.last_col = props.last_col_total
    table.horz_banding = props.banded_rows
    table.vert_banding = props.banded_cols


def _apply_table_merges(table: Any, element: BridgeTable) -> None:
    for row in element.cell_formats:
        for cell_format in row:
            merge = cell_format.merge
            if not merge.is_merge_origin or merge.merge_span_rows <= 1 and merge.merge_span_cols <= 1:
                continue
            end_row = cell_format.grid_row + merge.merge_span_rows - 1
            end_col = cell_format.grid_col + merge.merge_span_cols - 1
            if end_row < len(table.rows) and end_col < len(table.columns):
                table.cell(cell_format.grid_row, cell_format.grid_col).merge(table.cell(end_row, end_col))


def _apply_table_cell_solid_fill(cell: Any, color: Any, theme_colors: dict[str, str] | None = None) -> None:
    """Write solidFill directly into tcPr XML, preserving alpha transparency."""
    from pptx.oxml.ns import qn
    from lxml import etree
    from percy.bridge.elements import ColorSpec
    try:
        tc = cell._tc
        tc_pr = tc.find(qn("a:tcPr"))
        if tc_pr is None:
            tc_pr = etree.SubElement(tc, qn("a:tcPr"))
        for fill_tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
            el = tc_pr.find(qn(fill_tag))
            if el is not None:
                tc_pr.remove(el)

        if isinstance(color, ColorSpec):
            # Resolve without alpha to get pure RGB — then write alpha as separate XML element.
            # color.resolve() pre-bakes alpha into the hex, which would double-apply transparency.
            color_no_alpha = ColorSpec(
                value=color.value, lum_mod=color.lum_mod, lum_off=color.lum_off,
                shade=color.shade, tint=color.tint, hue_mod=color.hue_mod, sat_mod=color.sat_mod,
            )
            hex_val = color_no_alpha.resolve(theme_colors or {}).lstrip("#")
        else:
            hex_val = str(color).lstrip("#")
        if len(hex_val) != 6:
            return

        solid_el = etree.SubElement(tc_pr, qn("a:solidFill"))
        clr_el = etree.SubElement(solid_el, qn("a:srgbClr"))
        clr_el.set("val", hex_val.upper())
        if isinstance(color, ColorSpec) and color.alpha is not None and color.alpha < 100000:
            etree.SubElement(clr_el, qn("a:alpha")).set("val", str(int(color.alpha)))
    except Exception:
        pass


def _apply_table_cell_format(cell: Any, cell_format: Any, theme_colors: dict[str, str] | None = None) -> None:
    # Apply borders first — OOXML schema requires lnL/lnR/lnT/lnB before fill elements in tcPr
    if cell_format.borders is not None:
        _apply_table_cell_borders(cell, cell_format.borders, theme_colors)
    if cell_format.fill_type == "noFill":
        try:
            from pptx.oxml.ns import qn
            from lxml import etree
            tc = cell._tc
            tc_pr = tc.find(qn("a:tcPr"))
            if tc_pr is None:
                tc_pr = etree.SubElement(tc, qn("a:tcPr"))
            for fill_tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
                el = tc_pr.find(qn(fill_tag))
                if el is not None:
                    tc_pr.remove(el)
            tc_pr.append(etree.Element(qn("a:noFill")))
        except Exception:
            pass
    elif cell_format.fill_color:
        _apply_table_cell_solid_fill(cell, cell_format.fill_color, theme_colors)
    if cell_format.margins.margin_left is not None:
        cell.margin_left = Inches(cell_format.margins.margin_left)
    if cell_format.margins.margin_right is not None:
        cell.margin_right = Inches(cell_format.margins.margin_right)
    if cell_format.margins.margin_top is not None:
        cell.margin_top = Inches(cell_format.margins.margin_top)
    if cell_format.margins.margin_bottom is not None:
        cell.margin_bottom = Inches(cell_format.margins.margin_bottom)
    vertical_anchor = _enum_member(MSO_VERTICAL_ANCHOR, cell_format.alignment.vertical_alignment)
    if vertical_anchor is not None:
        cell.vertical_anchor = vertical_anchor
    for paragraph in cell.text_frame.paragraphs:
        alignment = _enum_member(PP_ALIGN, cell_format.alignment.text_alignment)
        if alignment is not None:
            paragraph.alignment = alignment


def _apply_table_cell_borders(cell: Any, borders: Any, theme_colors: dict[str, str] | None = None) -> None:
    from pptx.oxml.ns import qn
    from lxml import etree

    tc = cell._tc
    tc_pr = tc.find(qn("a:tcPr"))
    if tc_pr is None:
        tc_pr = etree.SubElement(tc, qn("a:tcPr"))

    border_map = {
        "lnL": borders.border_left,
        "lnR": borders.border_right,
        "lnT": borders.border_top,
        "lnB": borders.border_bottom,
        "lnTlToBr": borders.diagonal_down,
        "lnBlToTr": borders.diagonal_up,
    }
    for tag_name, border in border_map.items():
        if border is None:
            continue
        full_tag = qn(f"a:{tag_name}")
        existing = tc_pr.find(full_tag)
        if existing is not None:
            tc_pr.remove(existing)
        ln = etree.SubElement(tc_pr, full_tag)
        if not border.visible:
            etree.SubElement(ln, qn("a:noFill"))
            continue
        if border.width is not None:
            ln.set("w", str(int(Pt(border.width))))
        if border.color:
            hex_val = _resolve_color_hex(border.color, theme_colors or {})
            if hex_val and len(hex_val) == 6:
                solid = etree.SubElement(ln, qn("a:solidFill"))
                etree.SubElement(solid, qn("a:srgbClr")).set("val", hex_val.upper())
        dash_val = (border.dash_style or border.style or "solid").lower()
        etree.SubElement(ln, qn("a:prstDash")).set("val", dash_val)


def _add_shape(slide: Any, element: BridgeShape, theme_colors: dict[str, str] | None = None, default_text_color: str | None = None) -> list[dict[str, Any]]:
    shape_type = _shape_type_from_preset(element.shape_identification.geometry_preset)
    shape = slide.shapes.add_shape(shape_type, *_box(element))
    _apply_transform(shape, element)
    _apply_adjustments(shape, element.shape_identification.geometry_adjustments)
    _apply_fill(shape.fill, element.fill, theme_colors)
    _apply_line(shape.line, element.line, theme_colors)
    if element.text_content.has_text:
        _apply_text_frame(shape.text_frame, element.text_content.paragraphs, theme_colors, default_text_color, slide_part=safe_get(lambda: slide.part))
        _apply_shape_text_settings(shape.text_frame, element.text_frame)
    _apply_outer_shadow(shape, getattr(element, "shadow", None), theme_colors)
    return _unsupported_color_diags(element, [element.fill.color, element.line.color], theme_colors)


def _add_connector(slide: Any, element: BridgeConnector, theme_colors: dict[str, str] | None = None) -> list[dict[str, Any]]:
    connector = slide.shapes.add_connector(
        _connector_type(element.connector_type),
        Inches(element.endpoints.start_x),
        Inches(element.endpoints.start_y),
        Inches(element.endpoints.end_x),
        Inches(element.endpoints.end_y),
    )
    _apply_line(connector.line, element.line, theme_colors)
    _apply_line_ends(connector.line, getattr(element.line, 'head_end', None), getattr(element.line, 'tail_end', None), getattr(element.line, 'head_size', None), getattr(element.line, 'tail_size', None))
    return _unsupported_color_diags(element, [element.line.color], theme_colors)


def _add_freeform(slide: Any, element: BridgeFreeform, theme_colors: dict[str, str] | None = None, default_text_color: str | None = None) -> list[dict[str, Any]]:
    if element.geometry_xml:
        return _add_freeform_from_xml(slide, element, theme_colors)
    if not _can_rebuild_freeform(element):
        _add_generic_shape(slide, element)
        return [_diag(element, "complex_freeform_rebuild", "Complex freeform rebuilt as generic rectangle.")]
    path = element.paths[0]
    commands = path.commands
    start = commands[0].points[0]
    x_scale = _scale(element.position.width, path.width)
    y_scale = _scale(element.position.height, path.height)
    builder = slide.shapes.build_freeform(start[0], start[1], scale=(x_scale, y_scale))
    for command in commands[1:]:
        if command.command == "moveTo" and command.points:
            builder.move_to(*command.points[0])
        elif command.command == "lnTo" and command.points:
            builder.add_line_segments([command.points[0]], close=False)
        elif command.command == "close":
            builder.add_line_segments([], close=True)
    shape = builder.convert_to_shape(Inches(element.position.left), Inches(element.position.top))
    _apply_transform(shape, element)
    _apply_fill(shape.fill, element.fill, theme_colors)
    _apply_line(shape.line, element.line, theme_colors)
    return _unsupported_color_diags(element, [element.fill.fill_color, element.line.line_color], theme_colors)


def _add_freeform_from_xml(slide: Any, element: BridgeFreeform, theme_colors: dict[str, str] | None = None) -> list[dict[str, Any]]:
    """Rebuild a freeform shape by injecting the stored custGeom XML into a placeholder shape."""
    try:
        from lxml import etree
        from pptx.oxml.ns import qn

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *_box(element))
        spPr = shape.element.spPr

        # Remove the prstGeom added by add_shape
        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            spPr.remove(prstGeom)
        # Remove any noFill from default rectangle
        noFill = spPr.find(qn("a:noFill"))
        if noFill is not None:
            spPr.remove(noFill)

        # Insert stored custGeom — goes right after xfrm (index 1 if xfrm is first)
        custGeom_el = etree.fromstring(element.geometry_xml)
        xfrm = spPr.find(qn("a:xfrm"))
        insert_pos = list(spPr).index(xfrm) + 1 if xfrm is not None else 0
        spPr.insert(insert_pos, custGeom_el)

        _apply_transform(shape, element)
        _apply_fill(shape.fill, element.fill, theme_colors)
        _apply_line(shape.line, element.line, theme_colors)
        return _unsupported_color_diags(element, [element.fill.fill_color, element.line.line_color], theme_colors)
    except Exception as exc:
        _add_generic_shape(slide, element)
        return [_diag(element, "freeform_xml_inject_failed", f"custGeom XML inject failed: {exc}")]


def _can_rebuild_freeform(element: BridgeFreeform) -> bool:
    if len(element.paths) != 1 or not element.paths[0].commands:
        return False
    commands = element.paths[0].commands
    if commands[0].command != "moveTo" or not commands[0].points:
        return False
    return all(command.command in {"moveTo", "lnTo", "close"} for command in commands)


def _scale(size_inches: float, local_extent: int) -> float:
    if not local_extent:
        return 1.0
    return Inches(size_inches) / local_extent


def _add_generic_shape(slide: Any, element: Any) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *_box(element))
    text = getattr(getattr(element, "text_content", None), "text_content", None)
    if text:
        shape.text = text


def _apply_transform(shape: Any, element: Any) -> None:
    """Apply rotation and flip from element.transforms to a python-pptx shape."""
    t = getattr(element, "transforms", None)
    if t is None:
        return
    if t.rotation:
        try:
            shape.rotation = t.rotation
        except Exception:
            pass
    if t.flip_h or t.flip_v:
        try:
            from pptx.oxml.ns import qn
            xfrm = shape.element.spPr.xfrm
            if xfrm is not None:
                if t.flip_h:
                    xfrm.set("flipH", "1")
                if t.flip_v:
                    xfrm.set("flipV", "1")
        except Exception:
            pass


def _add_chart(slide: Any, element: BridgeChart, theme_colors: dict[str, str] | None = None) -> list[dict[str, Any]]:
    chart_type = _chart_type(element.chart_type)
    if chart_type is None:
        _add_chart_placeholder(slide, element)
        return [_diag(element, "chart_rebuild_placeholder", f"Unknown chart type: {element.chart_type}")]

    chart_data = _build_chart_data(element)
    if chart_data is None:
        # Empty chart (no series) — add invisible transparent shape to preserve z-order
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *_box(element))
        shape.fill.background()
        shape.line.fill.background()
        return [_diag(element, "chart_rebuild_placeholder", f"No data for chart type: {element.chart_type}")]

    chart_shape = slide.shapes.add_chart(chart_type, *_box(element), chart_data)

    # Replace the fresh embedded workbook with an updated version of the original,
    # so workbook structure (named ranges, formatting, sheet names) is preserved.
    _inject_workbook(chart_shape, element)

    chart = chart_shape.chart
    _inject_chart_overlay_files(chart_shape, element)
    _apply_chart_title(chart, element, theme_colors)
    _apply_chart_legend(chart, element, theme_colors or {})
    _apply_chart_axes(chart, element, theme_colors or {})
    _apply_chart_plot_area(chart, element, theme_colors or {})
    _apply_chart_gap_width(chart, element)
    _apply_chart_series_colors(chart, element, theme_colors or {})
    _apply_chart_plot_properties(chart, element)
    _apply_chart_space_txpr(chart, element, theme_colors or {})
    _apply_chart_disp_blanks_as(chart, element)
    _fix_data_label_num_fmt(chart)
    return _chart_rebuild_debt(element)


_ACCENT_ORDER = ["ACCENT_1", "ACCENT_2", "ACCENT_3", "ACCENT_4", "ACCENT_5", "ACCENT_6"]

_CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
_DML_NS  = "http://schemas.openxmlformats.org/drawingml/2006/main"

_SERIES_DATA_TAGS = frozenset({"cat", "val", "xVal", "yVal", "bubbleSize"})


def _insert_before_data_elements(ser_el: Any, new_el: Any) -> None:
    """Insert new_el before cat/val/xVal/yVal in the series element (OOXML ordering)."""
    for child in ser_el:
        local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if local in _SERIES_DATA_TAGS:
            child.addprevious(new_el)
            return
    ser_el.append(new_el)


_DL_POSITION_MAP = {
    "INSIDE_END": "inEnd",
    "INSIDE_BASE": "inBase",
    "OUTSIDE_END": "outEnd",
    "CENTER": "ctr",
    "BEST_FIT": "bestFit",
    "LEFT": "l",
    "RIGHT": "r",
    "ABOVE": "t",
    "BELOW": "b",
}


_DL_POS_MAP = {
    "INSIDE_END": "inEnd", "INSIDE_BASE": "inBase", "OUTSIDE_END": "outEnd",
    "BEST_FIT": "bestFit", "CENTER": "ctr", "LEFT": "l", "RIGHT": "r",
    "ABOVE": "t", "BELOW": "b",
}

def _apply_series_data_labels(pptx_series: Any, dl: Any, theme_colors: dict[str, str], custom_labels: "dict[int,str]|None" = None) -> None:
    """Inject <c:dLbls> into a series element to show data labels."""
    from lxml import etree
    _A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    if dl is None and not custom_labels:
        return
    ser_el = pptx_series._element
    # Remove existing dLbls if any
    for old in ser_el.findall(f"{{{_CHART_NS}}}dLbls"):
        ser_el.remove(old)
    # dLbls must come before cat/val/xVal/yVal in the OOXML schema
    dLbls = etree.Element(f"{{{_CHART_NS}}}dLbls")
    _insert_before_data_elements(ser_el, dLbls)
    # Per-point custom text labels: dLbl elements come FIRST in CT_DLbls schema
    if custom_labels:
        for pt_idx in sorted(custom_labels.keys()):
            txt = custom_labels[pt_idx]
            dLbl = etree.SubElement(dLbls, f"{{{_CHART_NS}}}dLbl")
            idx_el = etree.SubElement(dLbl, f"{{{_CHART_NS}}}idx")
            idx_el.set("val", str(pt_idx))
            tx_el = etree.SubElement(dLbl, f"{{{_CHART_NS}}}tx")
            rich = etree.SubElement(tx_el, f"{{{_CHART_NS}}}rich")
            etree.SubElement(rich, f"{{{_A}}}bodyPr")
            etree.SubElement(rich, f"{{{_A}}}lstStyle")
            p = etree.SubElement(rich, f"{{{_A}}}p")
            r = etree.SubElement(p, f"{{{_A}}}r")
            rPr = etree.SubElement(r, f"{{{_A}}}rPr")
            rPr.set("lang", "en-US")
            if dl.font_size:
                rPr.set("sz", str(int(dl.font_size * 100)))
            if dl.font_bold is not None:
                rPr.set("b", "1" if dl.font_bold else "0")
            t_el = etree.SubElement(r, f"{{{_A}}}t")
            t_el.text = txt
            # Per-dLbl show flags mirror series-level but showVal=0 since tx has the text
            for tag, bridge_val in [
                ("showLegendKey", getattr(dl, "show_legend_key", False)),
                ("showVal", False),
                ("showCatName", getattr(dl, "show_cat_name", False)),
                ("showSerName", getattr(dl, "show_ser_name", False)),
                ("showPercent", getattr(dl, "show_percent", False)),
                ("showBubbleSize", getattr(dl, "show_bubble_size", False)),
            ]:
                flag = etree.SubElement(dLbl, f"{{{_CHART_NS}}}{tag}")
                flag.set("val", "1" if bridge_val else "0")
    # OOXML CT_DLbls child order: numFmt → spPr → txPr → dLblPos → showXxx → separator → leaderLines
    # Number format
    if dl and getattr(dl, "format", None):
        numFmt = etree.SubElement(dLbls, f"{{{_CHART_NS}}}numFmt")
        numFmt.set("formatCode", dl.format)
        numFmt.set("sourceLinked", "0")
    # Font formatting via txPr
    has_font = dl and (dl.font_name or dl.font_size or dl.font_bold or dl.font_color)
    if has_font:
        txPr = etree.SubElement(dLbls, f"{{{_CHART_NS}}}txPr")
        bodyPr = etree.SubElement(txPr, f"{{{_A}}}bodyPr")
        bodyPr.set("rot", "0")
        bodyPr.set("vert", "horz")
        etree.SubElement(txPr, f"{{{_A}}}lstStyle")
        p = etree.SubElement(txPr, f"{{{_A}}}p")
        pPr = etree.SubElement(p, f"{{{_A}}}pPr")
        defRPr = etree.SubElement(pPr, f"{{{_A}}}defRPr")
        if dl.font_bold is not None:
            defRPr.set("b", "1" if dl.font_bold else "0")
        if dl.font_size:
            defRPr.set("sz", str(int(dl.font_size * 100)))
        if dl.font_color:  # solidFill must precede latin in CT_TextCharacterProperties
            solid_fill = etree.SubElement(defRPr, f"{{{_A}}}solidFill")
            color_spec = dl.font_color
            color_val = color_spec.value if hasattr(color_spec, "value") else str(color_spec)
            if color_val and color_val.startswith("scheme:"):
                scheme_name = color_val[7:]
                _SCHEME_NAME_MAP = {
                    "LIGHT_1": "lt1", "DARK_1": "dk1", "LIGHT_2": "lt2", "DARK_2": "dk2",
                    "ACCENT_1": "accent1", "ACCENT_2": "accent2", "ACCENT_3": "accent3",
                    "ACCENT_4": "accent4", "ACCENT_5": "accent5", "ACCENT_6": "accent6",
                    "BACKGROUND_1": "lt1", "TEXT_1": "dk1", "BACKGROUND_2": "lt2", "TEXT_2": "dk2",
                }
                scheme_clr = etree.SubElement(solid_fill, f"{{{_A}}}schemeClr")
                scheme_clr.set("val", _SCHEME_NAME_MAP.get(scheme_name, scheme_name.lower()))
            elif color_val and color_val.startswith("#"):
                srgb = etree.SubElement(solid_fill, f"{{{_A}}}srgbClr")
                srgb.set("val", color_val.lstrip("#").upper())
            elif theme_colors and color_val in theme_colors:
                srgb = etree.SubElement(solid_fill, f"{{{_A}}}srgbClr")
                srgb.set("val", theme_colors[color_val].lstrip("#").upper())
        if dl.font_name:  # latin comes after solidFill in schema order
            latin = etree.SubElement(defRPr, f"{{{_A}}}latin")
            latin.set("typeface", dl.font_name)
    # Label position — only emit if valid for the series' actual chart element type
    pos = getattr(dl, "position", None) if dl else None
    if pos:
        xml_pos = _DL_POS_MAP.get(pos.upper(), pos.lower())
        # Positions t/b/l/r are only valid for line/scatter/radar/pie charts, not bar/column
        parent_tag = ser_el.getparent().tag.split("}")[-1] if ser_el.getparent() is not None else ""
        _bar_only_positions = {"ctr", "inBase", "inEnd", "outEnd", "bestFit"}
        _line_only_positions = {"t", "b", "l", "r"}
        is_bar = "barChart" in parent_tag or "areaChart" in parent_tag
        if not (is_bar and xml_pos in _line_only_positions):
            dLblPos_el = etree.SubElement(dLbls, f"{{{_CHART_NS}}}dLblPos")
            dLblPos_el.set("val", xml_pos)
    # Show flags
    for tag, bridge_val, default_val in [
        ("showLegendKey", getattr(dl, "show_legend_key", False), "0"),
        ("showVal", getattr(dl, "show_val", True), "1"),
        ("showCatName", getattr(dl, "show_cat_name", False), "0"),
        ("showSerName", getattr(dl, "show_ser_name", False), "0"),
        ("showPercent", getattr(dl, "show_percent", False), "0"),
        ("showBubbleSize", getattr(dl, "show_bubble_size", False), "0"),
    ]:
        el = etree.SubElement(dLbls, f"{{{_CHART_NS}}}{tag}")
        el.set("val", "1" if bridge_val else "0")
    if getattr(dl, "show_leader_lines", False):
        ll = etree.SubElement(dLbls, f"{{{_CHART_NS}}}showLeaderLines")
        ll.set("val", "1")
    sep = getattr(dl, "separator", None)
    if sep:
        sep_el = etree.SubElement(dLbls, f"{{{_CHART_NS}}}separator")
        sep_el.text = sep


def _apply_chart_series_colors(chart: Any, element: BridgeChart, theme_colors: dict[str, str]) -> None:
    ct = (element.chart_type or "").upper()
    use_line = ct in _XY_TYPES or ct in _LINE_TYPES
    for i, bridge_series in enumerate(element.series):
        if i >= len(chart.series):
            break
        pptx_series = chart.series[i]
        # Resolve the series color: explicit → theme accent fallback
        series_color = bridge_series.color
        if not series_color and theme_colors:
            series_color = theme_colors.get(_ACCENT_ORDER[i % len(_ACCENT_ORDER)])
        if series_color:
            try:
                if use_line:
                    _set_color(pptx_series.format.line.color, series_color, theme_colors)
                else:
                    pptx_series.format.fill.solid()
                    _set_color(pptx_series.format.fill.fore_color, series_color, theme_colors)
            except Exception:
                pass
        if bridge_series.line.line_color:
            try:
                _set_color(pptx_series.format.line.color, bridge_series.line.line_color, theme_colors)
            except Exception:
                pass
        # Series line width and style
        if bridge_series.line.line_width:
            try:
                pptx_series.format.line.width = Pt(bridge_series.line.line_width)
            except Exception:
                pass
        if bridge_series.line.line_style:
            try:
                ds = _line_dash(bridge_series.line.line_style)
                if ds is not None:
                    pptx_series.format.line.dash_style = ds
            except Exception:
                pass
        if not bridge_series.line.line_visible:
            try:
                pptx_series.format.line.fill.background()
            except Exception:
                pass
        if bridge_series.marker.marker_color:
            try:
                pptx_series.marker.format.fill.solid()
                _set_color(pptx_series.marker.format.fill.fore_color, bridge_series.marker.marker_color, theme_colors)
            except Exception:
                pass
        for pt_idx, pt_fmt in (bridge_series.point_formatting or {}).items():
            pt_color = pt_fmt.get("fill_color") if pt_fmt else None
            if not pt_color:
                continue
            try:
                point = pptx_series.points[pt_idx]
                point.format.fill.solid()
                _set_color(point.format.fill.fore_color, pt_color, theme_colors)
            except Exception:
                pass
        if getattr(bridge_series, "smooth", False):
            try:
                from pptx.oxml.ns import qn
                from lxml import etree
                ser_el = pptx_series._element
                existing_smooth = ser_el.find(qn("c:smooth"))
                if existing_smooth is None:
                    smooth_el = etree.SubElement(ser_el, qn("c:smooth"))
                    smooth_el.set("val", "1")
            except Exception:
                pass
        dl = bridge_series.data_labels
        custom_labels = getattr(bridge_series, "custom_labels", None) or {}
        if (dl and dl.show) or custom_labels:
            try:
                _apply_series_data_labels(pptx_series, dl, theme_colors, custom_labels=custom_labels)
            except Exception:
                pass


def _fix_data_label_num_fmt(chart: Any) -> None:
    """Ensure data labels that show values use the value axis number format.
    When dLbls has no numFmt, PowerPoint shows raw floats. We inject numFmt
    with sourceLinked='0' and the value axis formatCode so labels display correctly.
    Handles combo charts by mapping each chart-type element to its own valAx.
    """
    from lxml import etree
    try:
        chart_el = chart._element
        plotArea = chart_el.find(f"{{{_CHART_NS}}}chart/{{{_CHART_NS}}}plotArea")
        if plotArea is None:
            return
        # Build axId → formatCode map from all valAx elements
        ax_fmt: dict[str, str] = {}
        for valAx in plotArea.findall(f"{{{_CHART_NS}}}valAx"):
            axId_el = valAx.find(f"{{{_CHART_NS}}}axId")
            numFmt_el = valAx.find(f"{{{_CHART_NS}}}numFmt")
            if axId_el is not None and numFmt_el is not None:
                fmt = numFmt_el.get("formatCode", "")
                if fmt and fmt != "General":
                    ax_fmt[axId_el.get("val", "")] = fmt
        if not ax_fmt:
            return
        # Walk each chart-type element and apply its valAx format to its series
        _chart_type_tags = {
            "barChart", "lineChart", "pieChart", "areaChart", "doughnutChart",
            "radarChart", "bubbleChart", "scatterChart", "stockChart",
        }
        for ct_el in plotArea:
            tag = ct_el.tag.split("}")[-1] if "}" in ct_el.tag else ct_el.tag
            if tag not in _chart_type_tags:
                continue
            # Get this chart type's value axis id
            fmt = None
            for axId_el in ct_el.findall(f"{{{_CHART_NS}}}axId"):
                aid = axId_el.get("val", "")
                if aid in ax_fmt:
                    fmt = ax_fmt[aid]
                    break
            if not fmt:
                # Fallback: use any available format
                fmt = next(iter(ax_fmt.values()), None)
            if not fmt:
                continue
            for ser in ct_el.findall(f"{{{_CHART_NS}}}ser"):
                dLbls = ser.find(f"{{{_CHART_NS}}}dLbls")
                if dLbls is None:
                    continue
                showVal = dLbls.find(f"{{{_CHART_NS}}}showVal")
                if showVal is None or showVal.get("val", "1") == "0":
                    continue
                existing = dLbls.find(f"{{{_CHART_NS}}}numFmt")
                if existing is not None:
                    continue
                numFmt_new = etree.Element(f"{{{_CHART_NS}}}numFmt")
                numFmt_new.set("formatCode", fmt)
                numFmt_new.set("sourceLinked", "0")
                dLbls.insert(0, numFmt_new)
    except Exception:
        pass


def _add_chart_placeholder(slide: Any, element: BridgeChart) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *_box(element))
    shape.fill.background()
    shape.line.fill.background()
    shape.text = f"Chart placeholder: {element.chart_type or 'unknown'}"


_XY_TYPES = frozenset({"XY_SCATTER", "XY_SCATTER_LINES", "XY_SCATTER_LINES_NO_MARKERS",
                        "XY_SCATTER_SMOOTH", "XY_SCATTER_SMOOTH_NO_MARKERS"})
_BUBBLE_TYPES = frozenset({"BUBBLE", "BUBBLE_THREE_D_EFFECT"})
_LINE_TYPES = frozenset({
    "LINE", "LINE_MARKERS", "LINE_STACKED", "LINE_STACKED_100",
    "LINE_MARKERS_STACKED", "LINE_MARKERS_STACKED_100",
    "LINE_INVERSE",
})


def _chart_type(chart_type: str | None) -> Any:
    if not chart_type:
        return None
    # XL_CHART_TYPE enum names are uppercase (e.g. COLUMN_CLUSTERED). The agent
    # builders produce lowercase ('column_clustered') matching the chart-data
    # PATCH API; the onboarded form may be either. Try both.
    return (getattr(XL_CHART_TYPE, chart_type, None)
            or getattr(XL_CHART_TYPE, chart_type.upper(), None))


def _build_chart_data(element: BridgeChart) -> Any:
    """Dispatch to the correct ChartData class based on chart_type."""
    ct = (element.chart_type or "").upper()
    if ct in _XY_TYPES:
        return _xy_chart_data(element)
    if ct in _BUBBLE_TYPES:
        return _bubble_chart_data(element)
    return _category_chart_data(element)


def _category_chart_data(element: BridgeChart) -> CategoryChartData | None:
    if not element.series:
        return None
    chart_data = CategoryChartData()
    max_points = max((len(s.values) for s in element.series), default=0)
    raw_cats = element.categories.categories or [str(i + 1) for i in range(max_points)]
    # Preserve numeric types so Excel doesn't treat numbers as text labels.
    if element.categories.categories_are_numeric:
        categories: list[Any] = [_to_number(c) for c in raw_cats]
    else:
        categories = raw_cats
    chart_data.categories = categories
    for series in element.series:
        chart_data.add_series(series.name or "", series.values)
    return chart_data


def _xy_chart_data(element: BridgeChart) -> XyChartData | None:
    if not element.series:
        return None
    chart_data = XyChartData()
    for series in element.series:
        s = chart_data.add_series(series.name or "")
        x_vals = series.x_values or list(range(len(series.values)))
        for x, y in zip(x_vals, series.values):
            s.add_data_point(
                x if x is not None else 0.0,
                y if y is not None else 0.0,
            )
    return chart_data


def _bubble_chart_data(element: BridgeChart) -> BubbleChartData | None:
    if not element.series:
        return None
    chart_data = BubbleChartData()
    for series in element.series:
        s = chart_data.add_series(series.name or "")
        x_vals = series.x_values or list(range(len(series.values)))
        # Bubble size: use point_formatting or default to 1.0
        for i, (x, y) in enumerate(zip(x_vals, series.values)):
            size = series.point_formatting.get(i, {}).get("bubble_size", 1.0)
            s.add_data_point(
                x if x is not None else 0.0,
                y if y is not None else 0.0,
                size,
            )
    return chart_data


def _to_number(val: str) -> int | float | str:
    try:
        f = float(val)
        return int(f) if f.is_integer() else f
    except (ValueError, TypeError):
        return val


def _inject_chart_overlay_files(chart_shape: Any, element: BridgeChart) -> None:
    """Inject chartStyle.xml and chartColors.xml blobs into the rebuilt chart part."""
    ov = getattr(element, "overlay_files", None)
    if ov is None:
        return
    try:
        from pptx.opc.package import Part
        from pptx.opc.packuri import PackURI
        chart_part = chart_shape.chart.part
        pkg = chart_part.package
        _STYLE_RELTYPE = "http://schemas.microsoft.com/office/drawing/2012/chartStyle"
        _COLORS_RELTYPE = "http://schemas.microsoft.com/office/drawing/2012/chartColorStyle"

        if ov.chart_style:
            for rId, rel in list(chart_part.rels.items()):
                if "chartstyle" in (rel.reltype or "").lower():
                    chart_part.drop_rel(rId)
            style_part = Part(
                partname=PackURI("/ppt/charts/style%d.xml" % id(chart_part)),
                content_type="application/vnd.ms-office.chartstyle+xml",
                package=pkg,
                blob=ov.chart_style,
            )
            chart_part.relate_to(style_part, _STYLE_RELTYPE)

        if ov.chart_colors:
            for rId, rel in list(chart_part.rels.items()):
                if "colorstyle" in (rel.reltype or "").lower() or ("color" in (rel.reltype or "").lower() and "chart" in (rel.reltype or "").lower()):
                    chart_part.drop_rel(rId)
            colors_part = Part(
                partname=PackURI("/ppt/charts/colors%d.xml" % id(chart_part)),
                content_type="application/vnd.ms-office.chartcolorstyle+xml",
                package=pkg,
                blob=ov.chart_colors,
            )
            chart_part.relate_to(colors_part, _COLORS_RELTYPE)
    except Exception:
        pass


def _inject_workbook(chart_shape: Any, element: BridgeChart) -> None:
    """
    Replace the fresh embedded Excel workbook with an updated version of the
    original, preserving workbook structure while writing in current series data.

    Falls back silently if openpyxl is not installed or the workbook snapshot
    is not available.
    """
    original_bytes = element.data_source.embedded_workbook_bytes
    sheets = element.data_source.workbook_sheets
    if not original_bytes or not sheets:
        return
    try:
        updated = _apply_bridge_data_to_xlsx(original_bytes, element, sheets)
        chart_shape.chart.part.chart_workbook.xlsx_part.blob = updated
    except Exception:
        pass  # leave the fresh workbook created by CategoryChartData / XyChartData


def _apply_bridge_data_to_xlsx(
    xlsx_bytes: bytes,
    element: BridgeChart,
    sheets: list[Any],
) -> bytes:
    """
    Open the embedded workbook, locate category and series value cells by
    comparing the snapshot against element data, update them, return new bytes.

    Strategy: scan the first sheet's cells for values that match the current
    category list to find the category column/start-row; series columns follow
    immediately to the right.
    """
    import openpyxl
    from io import BytesIO

    wb = openpyxl.load_workbook(BytesIO(xlsx_bytes))
    snapshot = sheets[0]
    if snapshot.name not in wb.sheetnames:
        target_sheet = wb.active
    else:
        target_sheet = wb[snapshot.name]

    # Build address→value map from snapshot to locate category column
    snap_map: dict[str, Any] = {c.address: c.value for c in snapshot.cells}

    cats = element.categories.categories or element.categories.categories_raw
    if not cats:
        buf = BytesIO(); wb.save(buf); return buf.getvalue()

    # Find the cell whose value matches the FIRST category value
    cat_addr = next(
        (addr for addr, val in snap_map.items() if str(val) == str(cats[0])),
        None,
    )
    if cat_addr is None:
        buf = BytesIO(); wb.save(buf); return buf.getvalue()

    # Derive category column / start row from the found address
    from openpyxl.utils.cell import column_index_from_string, coordinate_from_string
    cat_col_letter, cat_start_row = coordinate_from_string(cat_addr)
    cat_col = column_index_from_string(cat_col_letter)

    # Write categories
    for i, cat_val in enumerate(cats):
        target_sheet.cell(row=cat_start_row + i, column=cat_col, value=cat_val)

    # Write series values — each series occupies the next column to the right
    for j, series in enumerate(element.series):
        ser_col = cat_col + 1 + j
        # Header row (one above data start)
        header_row = cat_start_row - 1
        if header_row >= 1:
            target_sheet.cell(row=header_row, column=ser_col, value=series.name)
        for i, val in enumerate(series.values):
            target_sheet.cell(row=cat_start_row + i, column=ser_col, value=val)

    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _apply_chart_title(chart: Any, element: BridgeChart, theme_colors: dict[str, str] | None = None) -> None:
    title = element.title.title
    auto_title_deleted = getattr(element.title, "auto_title_deleted", None)
    if title is None:
        chart.has_title = False
        return
    # title='' means chart has a <c:title> element but no explicit <c:tx> text
    # (auto-title from series names). Set has_title=True but don't assign text.
    chart.has_title = True
    t = element.title
    if title:
        chart.chart_title.text_frame.text = title
        try:
            para = chart.chart_title.text_frame.paragraphs[0]
            font = para.font
            if t.title_font_name:
                font.name = t.title_font_name
            if t.title_font_size is not None:
                font.size = Pt(t.title_font_size)
            if t.title_font_bold is not None:
                font.bold = t.title_font_bold
            if t.title_font_italic is not None:
                font.italic = t.title_font_italic
            if t.title_font_color:
                _set_color(font.color, t.title_font_color, theme_colors or {})
        except Exception:
            pass
    # Apply manual layout position if present
    if t.title_position_x is not None or t.title_position_y is not None:
        try:
            from lxml import etree
            ct = chart.chart_title._element
            layout_el = ct.find(f"{{{_CHART_NS}}}layout")
            if layout_el is not None:
                ct.remove(layout_el)
            layout_el = etree.Element(f"{{{_CHART_NS}}}layout")
            # Insert layout as first child of title element
            ct.insert(0, layout_el)
            ml = etree.SubElement(layout_el, f"{{{_CHART_NS}}}manualLayout")
            for tag, val in [("xMode", "edge"), ("yMode", "edge"),
                             ("x", t.title_position_x), ("y", t.title_position_y),
                             ("w", t.title_width), ("h", t.title_height)]:
                if val is not None:
                    e = etree.SubElement(ml, f"{{{_CHART_NS}}}{tag}")
                    e.set("val", str(val) if isinstance(val, str) else repr(val))
        except Exception:
            pass


def _apply_chart_legend(chart: Any, element: BridgeChart, theme_colors: dict[str, str] | None = None) -> None:
    chart.has_legend = element.legend.visible
    if not element.legend.visible or chart.legend is None:
        return
    position = _legend_position(element.legend.position)
    if position is not None:
        chart.legend.position = position
    chart.legend.include_in_layout = element.legend.overlay
    lg = element.legend
    # Font properties
    try:
        font = chart.legend.font
        if lg.font_name:
            font.name = lg.font_name
        if lg.font_size is not None:
            font.size = Pt(lg.font_size)
        if lg.font_bold is not None:
            font.bold = lg.font_bold
        if lg.font_color:
            _set_color(font.color, lg.font_color, theme_colors)
    except Exception:
        pass
    # Fill
    if lg.fill_type == "noFill":
        try:
            chart.legend.format.fill.background()
        except Exception:
            pass
    elif lg.fill_color:
        try:
            _set_color_spec(chart.legend.format.fill.fore_color, lg.fill_color, theme_colors)
        except Exception:
            pass
    # Border
    if lg.border_width:
        try:
            chart.legend.format.line.width = Pt(lg.border_width)
        except Exception:
            pass
    # Manual layout
    if getattr(lg, "manual_layout_x", None) is not None:
        try:
            from pptx.oxml.ns import qn
            from lxml import etree
            leg_el = chart.legend._element
            # Remove existing layout
            existing = leg_el.find(qn("c:layout"))
            if existing is not None:
                leg_el.remove(existing)
            # Insert after legendPos (first child)
            layout_el = etree.Element(qn("c:layout"))
            ml = etree.SubElement(layout_el, qn("c:manualLayout"))
            for tag, val in [
                ("c:xMode", lg.manual_layout_x_mode or "edge"),
                ("c:yMode", lg.manual_layout_y_mode or "edge"),
            ]:
                e = etree.SubElement(ml, qn(tag))
                e.set("val", val)
            for tag, val in [
                ("c:x", lg.manual_layout_x),
                ("c:y", lg.manual_layout_y),
                ("c:w", lg.manual_layout_w),
                ("c:h", lg.manual_layout_h),
            ]:
                if val is not None:
                    e = etree.SubElement(ml, qn(tag))
                    e.set("val", str(val))
            # Insert layout after legendPos
            pos_el = leg_el.find(qn("c:legendPos"))
            insert_idx = list(leg_el).index(pos_el) + 1 if pos_el is not None else 0
            leg_el.insert(insert_idx, layout_el)
        except Exception:
            pass


def _apply_chart_axes(chart: Any, element: BridgeChart, theme_colors: dict[str, str] | None = None) -> None:
    _apply_axis(safe_get(lambda: chart.category_axis), element.category_axis, theme_colors)
    _apply_axis(safe_get(lambda: chart.value_axis), element.value_axis, theme_colors)


def _apply_axis(axis: Any, data: Any, theme_colors: dict[str, str] | None = None) -> None:
    if axis is None:
        return
    axis.visible = data.visible
    if data.min_value is not None:
        axis.minimum_scale = data.min_value
    if data.max_value is not None:
        axis.maximum_scale = data.max_value
    if data.tick_labels.number_format:
        axis.tick_labels.number_format = data.tick_labels.number_format
    # Tick label font properties (font color makes labels invisible when set to background color)
    tl = data.tick_labels
    if tl.tick_label_font_color or tl.tick_label_font_name or tl.tick_label_font_size or tl.tick_label_font_bold is not None:
        try:
            font = axis.tick_labels.font
            if tl.tick_label_font_color:
                _set_color(font.color, tl.tick_label_font_color, theme_colors)
            if tl.tick_label_font_name:
                font.name = tl.tick_label_font_name
            if tl.tick_label_font_size is not None:
                from pptx.util import Pt as _Pt
                font.size = _Pt(tl.tick_label_font_size)
            if tl.tick_label_font_bold is not None:
                font.bold = tl.tick_label_font_bold
        except Exception:
            pass
    if data.units.major_unit is not None and hasattr(axis, "major_unit"):
        axis.major_unit = data.units.major_unit
    if data.units.minor_unit is not None and hasattr(axis, "minor_unit"):
        axis.minor_unit = data.units.minor_unit
    # Axis line
    ax_line = getattr(data, "axis_line", None)
    if ax_line is not None:
        try:
            if not ax_line.line_visible:
                axis.format.line.fill.background()
            else:
                if ax_line.line_color:
                    _set_color(axis.format.line.color, ax_line.line_color, theme_colors)
                if ax_line.line_width:
                    axis.format.line.width = Pt(ax_line.line_width)
        except Exception:
            pass
    # Reverse order
    if getattr(data, "reverse_order", False):
        try:
            from pptx.oxml.ns import qn as _qn
            from lxml import etree
            ax_el = axis._element
            existing = ax_el.find(_qn("c:scaling"))
            if existing is None:
                scaling = etree.SubElement(ax_el, _qn("c:scaling"))
            else:
                scaling = existing
            orient_el = scaling.find(_qn("c:orientation"))
            if orient_el is None:
                orient_el = etree.SubElement(scaling, _qn("c:orientation"))
            orient_el.set("val", "maxMin")
        except Exception:
            pass
    # Delete flag (hides axis without removing it)
    if getattr(data, "delete", False):
        try:
            from pptx.oxml.ns import qn as _qn
            from lxml import etree
            ax_el = axis._element
            del_el = ax_el.find(_qn("c:delete"))
            if del_el is None:
                del_el = etree.SubElement(ax_el, _qn("c:delete"))
            del_el.set("val", "1")
        except Exception:
            pass
    # axPos — axis position (b/l/r/t)
    ax_pos = getattr(data, "ax_pos", None)
    if ax_pos is not None:
        try:
            from pptx.oxml.ns import qn as _qn
            from lxml import etree
            ax_el = axis._element
            existing = ax_el.find(_qn("c:axPos"))
            if existing is None:
                existing = etree.SubElement(ax_el, _qn("c:axPos"))
            existing.set("val", ax_pos)
        except Exception:
            pass

    # noMultiLvlLbl — suppress multi-level category labels
    if getattr(data, "no_multi_lvl_lbl", False):
        try:
            from pptx.oxml.ns import qn as _qn
            from lxml import etree
            ax_el = axis._element
            el = ax_el.find(_qn("c:noMultiLvlLbl"))
            if el is None:
                el = etree.SubElement(ax_el, _qn("c:noMultiLvlLbl"))
            el.set("val", "1")
        except Exception:
            pass

    # lblOffset
    lbl_offset = getattr(data, "lbl_offset", None)
    if lbl_offset is not None and lbl_offset != 100:
        try:
            from pptx.oxml.ns import qn as _qn
            from lxml import etree
            ax_el = axis._element
            el = ax_el.find(_qn("c:lblOffset"))
            if el is None:
                el = etree.SubElement(ax_el, _qn("c:lblOffset"))
            el.set("val", str(int(lbl_offset)))
        except Exception:
            pass

    # lblAlgn
    lbl_algn = getattr(data, "lbl_algn", None)
    if lbl_algn is not None:
        try:
            from pptx.oxml.ns import qn as _qn
            from lxml import etree
            ax_el = axis._element
            el = ax_el.find(_qn("c:lblAlgn"))
            if el is None:
                el = etree.SubElement(ax_el, _qn("c:lblAlgn"))
            el.set("val", lbl_algn)
        except Exception:
            pass

    # crossBetween
    cross_between = getattr(data, "cross_between", None)
    if cross_between is not None:
        try:
            from pptx.oxml.ns import qn as _qn
            from lxml import etree
            ax_el = axis._element
            el = ax_el.find(_qn("c:crossBetween"))
            if el is None:
                el = etree.SubElement(ax_el, _qn("c:crossBetween"))
            el.set("val", cross_between)
        except Exception:
            pass

    # Gridlines
    gridlines = getattr(data, "gridlines", None)
    if gridlines is not None:
        try:
            if not gridlines.has_major_gridlines:
                # Explicitly remove any default gridlines from the rebuilt chart
                from pptx.oxml.ns import qn as _qn
                ax_el = axis._element
                for gl_el in ax_el.findall(_qn("c:majorGridlines")):
                    ax_el.remove(gl_el)
            else:
                if gridlines.gridline_no_fill:
                    # Gridlines exist but are explicitly invisible (noFill line)
                    from lxml import etree
                    gl = axis.major_gridlines
                    sp_pr = gl.format._element.get_or_add_spPr()
                    from pptx.oxml.ns import qn as _qn2
                    ln_el = sp_pr.get_or_add_ln()
                    for child in list(ln_el):
                        ln_el.remove(child)
                    ln_el.append(etree.fromstring(f'<a:noFill xmlns:a="{_A_NS}"/>'))
                else:
                    if gridlines.gridline_color:
                        _set_color(axis.major_gridlines.format.line.color, gridlines.gridline_color, theme_colors)
                    if gridlines.gridline_width:
                        axis.major_gridlines.format.line.width = Pt(gridlines.gridline_width)
        except Exception:
            pass


def _apply_chart_gap_width(chart: Any, element: BridgeChart) -> None:
    """Apply gap width (bar spacing) and overlap from plot_properties."""
    pp = getattr(element, "plot_properties", None)
    if pp is None:
        return
    gap = getattr(pp, "bar_width_ratio", None)
    overlap = getattr(pp, "overlap", None)
    try:
        plot = chart.plots[0]
        if gap is not None and hasattr(plot, "gap_width"):
            plot.gap_width = int(gap)
        if overlap is not None and hasattr(plot, "overlap"):
            plot.overlap = int(overlap)
    except Exception:
        pass


def _apply_chart_plot_properties(chart: Any, element: BridgeChart) -> None:
    """Apply pie/donut-specific plot properties: firstSliceAng, holeSize, varyColors."""
    pp = getattr(element, "plot_properties", None)
    if pp is None:
        return
    try:
        from lxml import etree
        from pptx.oxml.ns import qn as _qn
        plot_el = chart._element.plotArea
        # Find the first plot type element (doughnutChart, pieChart, etc.)
        chart_els = [
            el for el in plot_el
            if el.tag.split("}")[-1].endswith("Chart")
        ]
        if not chart_els:
            return
        ct_el = chart_els[0]

        # firstSliceAng
        fsa = getattr(pp, "first_slice_ang", None)
        if fsa is not None:
            existing = ct_el.find(_qn("c:firstSliceAng"))
            if existing is None:
                existing = etree.SubElement(ct_el, _qn("c:firstSliceAng"))
            existing.set("val", str(int(fsa)))

        # holeSize
        hs = getattr(pp, "hole_size", None)
        if hs is not None:
            existing = ct_el.find(_qn("c:holeSize"))
            if existing is None:
                existing = etree.SubElement(ct_el, _qn("c:holeSize"))
            existing.set("val", str(int(hs)))

        # varyColors
        vc = getattr(pp, "vary_colors", None)
        if vc is not None:
            existing = ct_el.find(_qn("c:varyColors"))
            if existing is None:
                existing = etree.SubElement(ct_el, _qn("c:varyColors"))
            existing.set("val", "1" if vc else "0")
    except Exception:
        pass


def _apply_chart_space_txpr(chart: Any, element: BridgeChart, theme_colors: dict[str, str]) -> None:
    """Write chart-space c:txPr default text properties (size/bold/color only — no font name to avoid substitution artifacts on Windows)."""
    # Font name is intentionally skipped: the stored font may not exist on the render OS,
    # and substitution produces worse results than the chart style default.
    font_size = getattr(element, "chart_txpr_font_size", None)
    font_bold = getattr(element, "chart_txpr_font_bold", None)
    font_color = getattr(element, "chart_txpr_font_color", None)
    font_name = None  # never apply — see comment above
    if not any(v is not None for v in [font_size, font_bold, font_color]):
        return
    try:
        from lxml import etree
        from pptx.oxml.ns import qn as _qn
        cs = chart._chartSpace
        # Remove existing txPr if present
        existing = cs.find(_qn("c:txPr"))
        if existing is not None:
            cs.remove(existing)
        txpr_el = etree.Element(_qn("c:txPr"))
        etree.SubElement(txpr_el, _qn("a:bodyPr"))
        etree.SubElement(txpr_el, _qn("a:lstStyle"))
        p_el = etree.SubElement(txpr_el, _qn("a:p"))
        ppr_el = etree.SubElement(p_el, _qn("a:pPr"))
        defrpr = etree.SubElement(ppr_el, _qn("a:defRPr"))
        if font_size is not None:
            defrpr.set("sz", str(int(font_size * 100)))
        if font_bold is not None:
            defrpr.set("b", "1" if font_bold else "0")
        if font_color:
            hex_val = _resolve_color_hex(font_color, theme_colors)
            if hex_val:
                solid = etree.SubElement(defrpr, _qn("a:solidFill"))
                etree.SubElement(solid, _qn("a:srgbClr")).set("val", hex_val.upper().lstrip("#"))
        if font_name:
            latin = etree.SubElement(defrpr, _qn("a:latin"))
            latin.set("typeface", font_name)
        etree.SubElement(p_el, _qn("a:endParaRPr"))
        # Insert txPr before c:externalData (correct OOXML ordering)
        ext_data = cs.find(_qn("c:externalData"))
        if ext_data is not None:
            ext_data.addprevious(txpr_el)
        else:
            cs.append(txpr_el)
    except Exception:
        pass


def _apply_chart_disp_blanks_as(chart: Any, element: BridgeChart) -> None:
    """Apply c:chart/c:dispBlanksAs value."""
    val = getattr(element, "disp_blanks_as", None)
    if val is None:
        return
    try:
        from lxml import etree
        from pptx.oxml.ns import qn as _qn
        _C = _CHART_NS
        chart_el = chart._chartSpace.find(_qn("c:chart"))
        if chart_el is None:
            return
        existing = chart_el.find(_qn("c:dispBlanksAs"))
        if existing is not None:
            existing.set("val", val)
        else:
            el = etree.SubElement(chart_el, _qn("c:dispBlanksAs"))
            el.set("val", val)
    except Exception:
        pass


def _apply_chart_plot_area(chart: Any, element: BridgeChart, theme_colors: dict[str, str]) -> None:
    """Apply plot area border and fill for each series plot."""
    for bridge_series in element.series:
        ab = getattr(bridge_series, "plot_properties", None)
        if ab is None:
            ab = getattr(element, "plot_properties", None)
        if ab is None:
            continue
        area_border = getattr(ab, "area_border", None)
        if area_border is None:
            continue
        try:
            # Apply to the chart plot area (spPr of plotArea element)
            plot_area = chart._element.plotArea
            sp_pr = plot_area.get_or_add_spPr()
            from pptx.oxml.ns import qn
            from lxml import etree
            if area_border.no_line:
                ln = sp_pr.find(qn("a:ln"))
                if ln is None:
                    ln = etree.SubElement(sp_pr, qn("a:ln"))
                for child in list(ln):
                    ln.remove(child)
                etree.SubElement(ln, qn("a:noFill"))
            elif area_border.has_border:
                ln = sp_pr.find(qn("a:ln"))
                if ln is None:
                    ln = etree.SubElement(sp_pr, qn("a:ln"))
                if area_border.border_width:
                    ln.set("w", str(int(Pt(area_border.border_width))))
                if area_border.border_color:
                    hex_val = _resolve_color_hex(area_border.border_color, theme_colors)
                    if hex_val:
                        solid = etree.SubElement(ln, qn("a:solidFill"))
                        etree.SubElement(solid, qn("a:srgbClr")).set("val", hex_val.upper().lstrip("#"))
            if area_border.has_fill and area_border.fill_color:
                hex_val = _resolve_color_hex(area_border.fill_color, theme_colors)
                if hex_val:
                    for old_fill in [sp_pr.find(qn(t)) for t in ("a:solidFill", "a:noFill", "a:gradFill")]:
                        if old_fill is not None:
                            sp_pr.remove(old_fill)
                    solid = etree.SubElement(sp_pr, qn("a:solidFill"))
                    etree.SubElement(solid, qn("a:srgbClr")).set("val", hex_val.upper().lstrip("#"))
        except Exception:
            pass
        break  # Only need one pass — plot_properties is chart-level
    # Apply plot area manual layout if present.
    _ct = (element.chart_type or "").upper()
    # For horizontal bar charts, category axis label widths in the rebuilt chart differ
    # from the original (Excel workbook vs strCache rendering), so the original x/w fractions
    # don't map correctly. Skip layout entirely for these chart types.
    _bar_types = {"BAR", "BAR_CLUSTERED", "BAR_STACKED", "BAR_STACKED_100"}
    _skip_layout = _ct in _bar_types
    if getattr(element, "plot_area_x", None) is not None and not _skip_layout:
        try:
            from lxml import etree
            plot_area = chart._element.plotArea
            # Remove old layout element if any
            old_layout = plot_area.find(f"{{{_CHART_NS}}}layout")
            if old_layout is not None:
                plot_area.remove(old_layout)
            layout_el = etree.Element(f"{{{_CHART_NS}}}layout")
            # Insert layout as first child of plotArea (before chart type elements)
            plot_area.insert(0, layout_el)
            ml = etree.SubElement(layout_el, f"{{{_CHART_NS}}}manualLayout")
            lt = getattr(element, "plot_area_layout_target", None) or "inner"
            lt_el = etree.SubElement(ml, f"{{{_CHART_NS}}}layoutTarget")
            lt_el.set("val", lt)
            xm = getattr(element, "plot_area_x_mode", None) or "edge"
            ym = getattr(element, "plot_area_y_mode", None) or "edge"
            dims = [("xMode", xm), ("yMode", ym),
                    ("x", element.plot_area_x), ("y", element.plot_area_y),
                    ("w", element.plot_area_w), ("h", element.plot_area_h)]
            for tag, val in dims:
                if val is not None:
                    e = etree.SubElement(ml, f"{{{_CHART_NS}}}{tag}")
                    e.set("val", str(val) if isinstance(val, str) else repr(val))
        except Exception:
            pass


def _legend_position(position: str | None) -> Any:
    if position is None:
        return None
    from pptx.enum.chart import XL_LEGEND_POSITION

    return getattr(XL_LEGEND_POSITION, position, None)


def _chart_rebuild_debt(element: BridgeChart) -> list[dict[str, Any]]:
    diagnostics = []
    for debt in element.custom_properties.get("semantic_debt", []):
        if debt == "chart_external_ole_link":
            diagnostics.append(_diag(element, "chart_external_ole_link", "External workbook link not preserved; chart rebuilt from cached series/category data."))
        elif debt == "chart_embedded_workbook":
            diagnostics.append(_diag(element, "chart_embedded_workbook", "Original embedded workbook captured; chart rebuilt from cached series/category data."))
        elif debt == "chart_combo_plots":
            diagnostics.append(_diag(element, "chart_combo_plots", "Combo chart rebuilt as a single chart type."))
        else:
            diagnostics.append(_diag(element, debt, f"Chart semantic debt: {debt}"))
    return diagnostics


def _shape_type_from_preset(preset: str | None) -> Any:
    if not preset:
        return MSO_SHAPE.RECTANGLE
    try:
        return AutoShapeType.id_from_prst(preset)
    except Exception:
        return MSO_SHAPE.RECTANGLE


def _connector_type(connector_type: str | None) -> Any:
    if connector_type == "elbow":
        return MSO_CONNECTOR.ELBOW
    if connector_type == "curved":
        return MSO_CONNECTOR.CURVE
    return MSO_CONNECTOR.STRAIGHT


def _apply_fill(fill: Any, data: Any, theme_colors: dict[str, str] | None = None) -> None:
    fill_type = getattr(data, "fill_type", None)
    color = getattr(data, "color", None) or getattr(data, "fill_color", None)
    gradient_stops = getattr(data, "gradient_stops", None) or []
    gradient_angle = getattr(data, "gradient_angle", 0.0) or 0.0
    pattern_preset = getattr(data, "pattern_preset", None)
    bg_color = getattr(data, "bg_color", None)
    if fill_type == "BACKGROUND":
        fill.background()
        return
    if fill_type in ("gradient", "GRADIENT") and gradient_stops:
        _apply_gradient_fill(fill, gradient_stops, gradient_angle, theme_colors or {})
        return
    if fill_type == "PATTERNED" and pattern_preset:
        _apply_pattern_fill(fill, pattern_preset, color, bg_color, theme_colors)
        return
    if color is None:
        # No explicit fill recorded — make transparent to avoid inheriting a wrong theme style fill
        # GROUP fill (grpFill) inherits from parent group; when groups are flattened, treat as transparent.
        if fill_type is None or fill_type == "GROUP":
            fill.background()
        return
    fill.solid()
    _set_color(fill.fore_color, color, theme_colors)
    # Inject alpha modifier directly into the color XML element (python-pptx's
    # fill.transparency setter doesn't reliably persist to the a:alpha child element)
    from percy.bridge.elements import ColorSpec as _CS
    alpha_val: int | None = None
    transparency = getattr(data, "transparency", 0.0) or 0.0
    if transparency > 0.0:
        alpha_val = int((1.0 - min(1.0, transparency)) * 100000)
    elif isinstance(color, _CS) and color.alpha is not None and color.alpha < 100000:
        alpha_val = color.alpha
    if alpha_val is not None:
        try:
            from pptx.oxml.ns import qn
            from lxml import etree
            spPr = fill._xPr
            solidFill = spPr.find(qn("a:solidFill"))
            if solidFill is not None:
                clr_el = solidFill.find(qn("a:srgbClr"))
                if clr_el is None:
                    clr_el = solidFill.find(qn("a:schemeClr"))
                if clr_el is None:
                    clr_el = solidFill.find(qn("a:sysClr"))
                if clr_el is not None:
                    existing = clr_el.find(qn("a:alpha"))
                    if existing is not None:
                        clr_el.remove(existing)
                    alpha_el = etree.SubElement(clr_el, qn("a:alpha"))
                    alpha_el.set("val", str(alpha_val))
        except Exception:
            pass


def _apply_pattern_fill(
    fill: Any,
    pattern_preset: str,
    fg_color: Any,
    bg_color: Any,
    theme_colors: dict[str, str] | None = None,
) -> None:
    """Inject <a:pattFill> into the shape's spPr element."""
    try:
        from pptx.oxml.ns import qn
        from lxml import etree

        sp_pr = fill._xPr
        for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:blipFill", "a:pattFill"):
            old = sp_pr.find(qn(tag))
            if old is not None:
                sp_pr.remove(old)

        patt_fill = etree.SubElement(sp_pr, qn("a:pattFill"))
        patt_fill.set("prst", pattern_preset)

        if fg_color:
            fg_hex = _resolve_color_hex(fg_color, theme_colors or {})
            if fg_hex:
                fg_el = etree.SubElement(patt_fill, qn("a:fgClr"))
                clr = etree.SubElement(fg_el, qn("a:srgbClr"))
                clr.set("val", fg_hex.upper().lstrip("#"))

        if bg_color:
            bg_hex = _resolve_color_hex(bg_color, theme_colors or {})
            if bg_hex:
                bg_el = etree.SubElement(patt_fill, qn("a:bgClr"))
                clr = etree.SubElement(bg_el, qn("a:srgbClr"))
                clr.set("val", bg_hex.upper().lstrip("#"))
    except Exception:
        pass


def _apply_gradient_fill(fill: Any, stops: list, angle_deg: float, theme_colors: dict[str, str]) -> None:
    """Inject a linear gradient fill directly via OOXML since python-pptx has no gradient API."""
    try:
        from pptx.oxml.ns import qn
        from lxml import etree

        sp_pr = fill._xPr
        # Remove any existing fill child elements
        for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:blipFill", "a:pattFill"):
            old = sp_pr.find(qn(tag))
            if old is not None:
                sp_pr.remove(old)

        # Build <a:gradFill>
        grad_fill = etree.SubElement(sp_pr, qn("a:gradFill"))
        gs_lst = etree.SubElement(grad_fill, qn("a:gsLst"))
        from percy.bridge.elements import ColorSpec as _CS
        for stop in sorted(stops, key=lambda s: s.position):
            pos_val = str(int(stop.position * 100000))
            gs = etree.SubElement(gs_lst, qn("a:gs"), pos=pos_val)
            hex_val = _resolve_color_hex(stop.color, theme_colors)
            clr_el = etree.SubElement(gs, qn("a:srgbClr"))
            clr_el.set("val", (hex_val or "888888").upper().lstrip("#"))
            # Preserve alpha/transparency from the color spec
            if isinstance(stop.color, _CS) and stop.color.alpha is not None:
                alpha_val = max(0, min(100000, stop.color.alpha))
                etree.SubElement(clr_el, qn("a:alpha")).set("val", str(alpha_val))

        # Linear gradient direction — OOXML angle is in 60000ths of a degree, measured from 3-o'clock CW
        # PowerPoint angle 270 = top-to-bottom; 0 = left-to-right
        ooxml_angle = int(((90 - angle_deg) % 360) * 60000)
        lin = etree.SubElement(grad_fill, qn("a:lin"))
        lin.set("ang", str(ooxml_angle))
        lin.set("scaled", "0")
    except Exception:
        # Fall back to solid with first stop color
        try:
            if stops:
                hex_val = _resolve_color_hex(stops[0].color, theme_colors)
                if hex_val:
                    fill.solid()
                    fill.fore_color.rgb = RGBColor.from_string(hex_val.lstrip("#"))
        except Exception:
            pass


def _resolve_color_hex(color: "Any", theme_colors: dict[str, str]) -> str | None:
    """Resolve a color (ColorSpec or str) to a 6-char hex string."""
    from percy.bridge.elements import ColorSpec
    if not color:
        return None
    if isinstance(color, ColorSpec):
        if not color.value:
            return None
        resolved = color.resolve(theme_colors)
        return resolved.lstrip("#") if resolved != "#888888" else None
    # Legacy str
    normalized = color.lstrip("#")
    if normalized.startswith("scheme:"):
        key = normalized[7:]
        resolved = theme_colors.get(key, "")
        return resolved.lstrip("#") if resolved else None
    if len(normalized) == 6:
        return normalized
    if len(normalized) == 8:
        return normalized[2:]
    return None


def _apply_fill_color(fill: Any, color: Any, theme_colors: dict[str, str] | None = None) -> None:
    fill.solid()
    _set_color(fill.fore_color, color, theme_colors)


def _apply_line(line: Any, data: Any, theme_colors: dict[str, str] | None = None) -> None:
    visible = getattr(data, "visible", getattr(data, "has_border", getattr(data, "line_visible", True)))
    if not visible:
        line.fill.background()
        return
    color = getattr(data, "color", None) or getattr(data, "border_color", None) or getattr(data, "line_color", None)
    if color is not None:
        _set_color(line.color, color, theme_colors)
    width = getattr(data, "width", None) or getattr(data, "border_width", None) or getattr(data, "line_width", None)
    if width is not None:
        line.width = Pt(width)
    dash_style = getattr(data, "dash_style", None) or getattr(data, "line_dash", None)
    dash_enum = _line_dash(dash_style)
    if dash_enum is not None:
        line.dash_style = dash_enum
    # Line cap and join (no python-pptx API — set directly on lxml element)
    line_cap = getattr(data, "line_cap", None)
    line_join = getattr(data, "line_join", None)
    if line_cap or line_join:
        try:
            ln_el = line._ln
            if line_cap and ln_el is not None:
                ln_el.set("cap", line_cap)
            if line_join and ln_el is not None:
                from pptx.oxml.ns import qn
                _VALID_JOINS = {"round", "bevel", "miter"}
                join_val = line_join.lower()
                if join_val in _VALID_JOINS:
                    for old in [ln_el.find(qn(f"a:{t}Join")) for t in ("round", "bevel", "miter")]:
                        if old is not None:
                            ln_el.remove(old)
                    from lxml import etree
                    etree.SubElement(ln_el, qn(f"a:{join_val}Join"))
        except Exception:
            pass
    # Arrow heads (OOXML: headEnd / tailEnd via lxml since python-pptx has no API)
    head_end = getattr(data, "head_end", None)
    tail_end = getattr(data, "tail_end", None)
    head_size = getattr(data, "head_size", None)
    tail_size = getattr(data, "tail_size", None)
    if head_end or tail_end:
        _apply_line_ends(line, head_end, tail_end, head_size, tail_size)


def _apply_adjustments(shape: Any, adjustments: dict[str, Any]) -> None:
    if not adjustments:
        return
    adjustment_collection = getattr(shape, "adjustments", None)
    if adjustment_collection is None:
        return
    name_to_index = {
        adjustment.name: index
        for index, adjustment in enumerate(getattr(adjustment_collection, "_adjustments_", []))
    }
    for name, formula in adjustments.items():
        if name not in name_to_index:
            continue
        try:
            value = str(formula).removeprefix("val ")
            adjustment_collection[name_to_index[name]] = int(value) / 100000.0
        except Exception:
            continue


def _apply_autofit(text_frame: Any, autofit_type: str | None, font_scale: int | None = None,
                   ln_spc_reduction: int | None = None) -> None:
    """Set the bodyPr auto-fit element to match the captured autofit_type."""
    if not autofit_type:
        return
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        body_pr = text_frame._txBody.bodyPr
        _A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        for tag in (qn("a:spAutoFit"), f"{{{_A}}}normAutofit", qn("a:noAutofit")):
            el = body_pr.find(tag)
            if el is not None:
                body_pr.remove(el)
        if autofit_type == "TEXT_TO_FIT_SHAPE":
            norm_el = etree.Element(f"{{{_A}}}normAutofit")
            if font_scale is not None:
                norm_el.set("fontScale", str(font_scale))
            if ln_spc_reduction is not None:
                norm_el.set("lnSpcReduction", str(ln_spc_reduction))
            body_pr.append(norm_el)
        elif autofit_type == "SHAPE_TO_FIT_TEXT":
            body_pr.append(etree.Element(qn("a:spAutoFit")))
        elif autofit_type in ("NONE", "NO_AUTOFIT", "shrink"):
            # "shrink" = original had no explicit autofit element; use noAutofit to
            # prevent python-pptx's default spAutoFit from applying (preserves fixed height)
            body_pr.append(etree.Element(qn("a:noAutofit")))
    except Exception:
        pass


def _apply_text_box_settings(text_frame: Any, data: Any) -> None:
    word_wrap = getattr(data, "word_wrap", None)
    # Default to True — new textboxes don't wrap by default in python-pptx
    text_frame.word_wrap = word_wrap if word_wrap is not None else True
    vertical_anchor = getattr(data, "vertical_anchor", None)
    if vertical_anchor:
        va = _enum_member(MSO_VERTICAL_ANCHOR, vertical_anchor)
        if va is not None:
            try:
                text_frame.vertical_anchor = va
            except Exception:
                pass
    body_insets = getattr(data, "body_insets", {}) or {}
    for attr_name, key in (
        ("margin_left", "left"),
        ("margin_right", "right"),
        ("margin_top", "top"),
        ("margin_bottom", "bottom"),
    ):
        if key in body_insets:
            try:
                setattr(text_frame, attr_name, Inches(body_insets[key]))
            except Exception:
                pass
    _apply_autofit(text_frame, getattr(data, "autofit_type", None), getattr(data, "font_scale", None), getattr(data, "ln_spc_reduction", None))


def _apply_shape_text_settings(text_frame: Any, data: Any) -> None:
    if getattr(data, "word_wrap", None) is not None:
        text_frame.word_wrap = data.word_wrap
    vertical_anchor = getattr(data, "vertical_anchor", None)
    if vertical_anchor:
        va = _enum_member(MSO_VERTICAL_ANCHOR, vertical_anchor)
        if va is not None:
            try:
                text_frame.vertical_anchor = va
            except Exception:
                pass
    insets = getattr(data, "text_insets", {}) or {}
    for attr_name, key in (
        ("margin_left", "left"),
        ("margin_right", "right"),
        ("margin_top", "top"),
        ("margin_bottom", "bottom"),
    ):
        if key in insets:
            setattr(text_frame, attr_name, Inches(insets[key]))
    _apply_autofit(text_frame, getattr(data, "autofit_type", None), getattr(data, "font_scale", None), getattr(data, "ln_spc_reduction", None))


def _line_dash(dash_style: str | None) -> Any:
    if not dash_style:
        return None
    normalized = dash_style.upper()
    if normalized == "SOLID":
        return MSO_LINE.SOLID
    return getattr(MSO_LINE, normalized, None)


_ALIGN_XML_TO_ENUM = {
    "ctr": "CENTER", "l": "LEFT", "r": "RIGHT",
    "just": "JUSTIFY", "dist": "DISTRIBUTE",
    "justLow": "JUSTIFY_LOW", "thaiDist": "THAI_DISTRIBUTE",
}


def _enum_member(enum_cls: Any, name: str | None) -> Any:
    if not name:
        return None
    normalized = _ALIGN_XML_TO_ENUM.get(name, str(name).upper())
    return getattr(enum_cls, normalized, None)


def _set_color(color_format: Any, color: "Any", theme_colors: dict[str, str] | None = None) -> bool:
    """Apply a color (ColorSpec, str, or None) to a python-pptx color format object."""
    from percy.bridge.elements import ColorSpec
    if color is None:
        return False
    if isinstance(color, ColorSpec):
        return _set_color_spec(color_format, color, theme_colors)
    # Legacy str path — always resolve to hex, never write schemeClr
    normalized = color.removeprefix("#")
    if normalized.startswith("scheme:"):
        theme_name = normalized.removeprefix("scheme:")
        if theme_name == "NOT_THEME_COLOR":
            return False
        if theme_colors:
            hex_val = (theme_colors.get(theme_name) or "").lstrip("#")
            if len(hex_val) == 6:
                try:
                    color_format.rgb = RGBColor.from_string(hex_val)
                    return True
                except Exception:
                    pass
        return False
    if len(normalized) != 6:
        return False
    color_format.rgb = RGBColor.from_string(normalized)
    return True


def _set_color_spec(color_format: Any, spec: "ColorSpec", theme_colors: dict[str, str] | None = None) -> bool:
    """Apply a ColorSpec to a python-pptx color format object, always resolving to a concrete hex RGB."""
    if not spec.value:
        return False
    if spec.value.startswith("scheme:"):
        # Always resolve scheme colors to concrete hex — never write schemeClr back into rebuilt XML.
        # Resolve without alpha so we can write alpha as an explicit XML modifier.
        from percy.bridge.elements import ColorSpec as _CS
        spec_no_alpha = _CS(
            value=spec.value, lum_mod=spec.lum_mod, lum_off=spec.lum_off,
            shade=spec.shade, tint=spec.tint, hue_mod=spec.hue_mod, sat_mod=spec.sat_mod,
        )
        hex_color = spec_no_alpha.resolve(theme_colors or {}).lstrip("#")
        if len(hex_color) != 6 or hex_color.upper() == "888888":
            key = spec.value[7:]
            raw = (theme_colors or {}).get(key, "")
            hex_color = raw.lstrip("#")
        if len(hex_color) == 6:
            try:
                color_format.rgb = RGBColor.from_string(hex_color)
                if spec.alpha is not None and spec.alpha < 100000:
                    _add_alpha_to_color_format(color_format, spec.alpha)
                return True
            except Exception:
                pass
        return False
    hex_val = spec.value.lstrip("#")
    if len(hex_val) == 8:
        hex_val = hex_val[2:]
    if len(hex_val) != 6:
        return False
    try:
        color_format.rgb = RGBColor.from_string(hex_val)
        if spec.alpha is not None and spec.alpha < 100000:
            _add_alpha_to_color_format(color_format, spec.alpha)
        return True
    except Exception:
        return False


def _add_alpha_to_color_format(color_format: Any, alpha_val: int) -> None:
    """Add <a:alpha val="N"/> to the srgbClr element of a python-pptx color format."""
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        clr_el = color_format._color.element
        if clr_el is not None:
            existing = clr_el.find(qn("a:alpha"))
            if existing is not None:
                clr_el.remove(existing)
            etree.SubElement(clr_el, qn("a:alpha")).set("val", str(int(alpha_val)))
    except Exception:
        pass


def _theme_color(theme_name: str) -> Any:
    return getattr(MSO_THEME_COLOR, theme_name, None)


def _unsupported_color_diags(element: Any, colors: list, theme_colors: dict[str, str] | None = None) -> list[dict[str, Any]]:
    from percy.bridge.elements import ColorSpec
    diags = []
    for color in colors:
        if color is None:
            continue
        if isinstance(color, ColorSpec) and color.value.startswith("scheme:"):
            # Only flag if the color genuinely can't be resolved
            resolved = color.resolve(theme_colors or {})
            if not resolved or resolved == "#888888":
                diags.append(_diag(element, "scheme_color_unsupported", f"Scheme color: {color.value}"))
        elif isinstance(color, str) and color.startswith("scheme:"):
            key = color[7:]
            if not (theme_colors or {}).get(key):
                diags.append(_diag(element, "scheme_color_unsupported", f"Scheme color: {color}"))
    return diags


def _box(element: Any) -> tuple[Any, Any, Any, Any]:
    position = element.position
    return (
        Inches(position.left),
        Inches(position.top),
        Inches(position.width),
        Inches(position.height),
    )


def _diag(element: Any, code: str, message: str) -> dict[str, Any]:
    return {
        "element_type": getattr(element, "element_type", type(element).__name__),
        "slide_number": getattr(getattr(element, "identification", None), "slide_number", None),
        "source_shape_id": element.custom_properties.get("source_shape_id", None),
        "source_shape_name": element.custom_properties.get("source_shape_name", None),
        "code": code,
        "message": message,
    }


def _set_rgb_color(font: Any, color: str) -> None:
    _set_color(font.color, color)


def _set_run_caps(run: Any, caps: str) -> None:
    """Set the OOXML 'cap' attribute on a run's rPr element."""
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        r_elem = run._r
        rpr = r_elem.find(qn("a:rPr"))
        if rpr is None:
            rpr = etree.SubElement(r_elem, qn("a:rPr"))
            r_elem.insert(0, rpr)
        rpr.set("cap", caps)
    except Exception:
        pass


def _set_run_rpr_attr(run: Any, attr: str, value: str) -> None:
    """Set an arbitrary attribute on a run's <a:rPr> element via lxml."""
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        r_elem = run._r
        rpr = r_elem.find(qn("a:rPr"))
        if rpr is None:
            rpr = etree.SubElement(r_elem, qn("a:rPr"))
            r_elem.insert(0, rpr)
        rpr.set(attr, value)
    except Exception:
        pass


def _set_run_hyperlink(run: Any, url: str) -> None:
    """Add a click hyperlink to a run via OOXML relationship."""
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        r_elem = run._r
        rpr = r_elem.find(qn("a:rPr"))
        if rpr is None:
            rpr = etree.SubElement(r_elem, qn("a:rPr"))
            r_elem.insert(0, rpr)
        # Add relationship on the parent part (slide)
        part = run._r.getparent().getparent().getparent().getparent()
        if hasattr(part, "part"):
            from pptx.opc.constants import RELATIONSHIP_TYPE as RT
            rel = part.part.relate_to(url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink", is_external=True)
            ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
            hlinkClick = etree.SubElement(rpr, qn("a:hlinkClick"))
            hlinkClick.set(f"{{{ns_r}}}id", rel)
    except Exception:
        pass


def _apply_line_ends(line: Any, head_end: str | None, tail_end: str | None,
                     head_size: str | None, tail_size: str | None) -> None:
    """Write headEnd / tailEnd arrowhead elements to a line via lxml."""
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        ln = line._ln
        if ln is None:
            return
        if head_end:
            he = ln.find(qn("a:headEnd"))
            if he is None:
                he = etree.SubElement(ln, qn("a:headEnd"))
            he.set("type", head_end)
            if head_size:
                w, h = (head_size + "/med").split("/")[:2]
                he.set("w", w or "med")
                he.set("len", h or "med")
        if tail_end:
            te = ln.find(qn("a:tailEnd"))
            if te is None:
                te = etree.SubElement(ln, qn("a:tailEnd"))
            te.set("type", tail_end)
            if tail_size:
                w, h = (tail_size + "/med").split("/")[:2]
                te.set("w", w or "med")
                te.set("len", h or "med")
    except Exception:
        pass
