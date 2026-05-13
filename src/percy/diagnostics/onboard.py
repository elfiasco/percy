"""Hacky first-pass PPTX to Bridge onboarding."""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

from percy.bridge import (
    BridgeChart,
    BridgeConnector,
    BridgeFreeform,
    BridgeImage,
    BridgeShape,
    BridgeSlide,
    BridgeTable,
    BridgeText,
    Border,
    GradientStop,
    CellAlignment,
    CellBorders,
    CellFont,
    CellFormat,
    CellMerge,
    AreaBorder,
    AxisLine,
    AxisTitle,
    AxisUnits,
    BridgeAxis,
    ChartCategories,
    ChartDataSource,
    ChartSeries,
    ChartTitle,
    ChartWorkbookCell,
    ChartWorkbookSheet,
    ConnectorEndpoints,
    DataLabels,
    EmbeddedFont,
    FillAndBorder,
    FreeformFill,
    FreeformLine,
    FreeformPath,
    Gridlines,
    ImageBorder,
    ImageCropping,
    ImageData,
    ImageDimensions,
    ImageFileInfo,
    Identification,
    Legend,
    LineFormat,
    Margins,
    MarkerFormat,
    PercyDocument,
    PlotProperties,
    Position,
    PresentationMetadata,
    PathCommand,
    ReconstructionBlobs,
    ShapeIdentification,
    ShapeFill,
    ShapeInfo,
    ShapeLine,
    ShapeShadow,
    ShapeTextContent,
    ShapeTextFrame,
    TableDefaults,
    TableDimensions,
    TableProperties,
    TableStyle,
    TableStyleSection,
    TextParagraph,
    TextFrame,
    TextRun,
    Transform,
    TransformEmus,
)
from percy.diagnostics.common import emu_to_inches, enum_name, safe_get
from percy.diagnostics.inheritance import placeholder_info, resolve_body_pr, resolve_text_shape

@dataclass
class _OnboardContext:
    """Carries per-presentation state through the onboarding call chain."""
    major_font: str | None = None
    minor_font: str | None = None
    theme_colors: dict[str, str] = field(default_factory=dict)
    slide_part: Any = None

    def resolve_font_name(self, name: str | None) -> str | None:
        if name == "+mj-lt":
            return self.major_font or name
        if name == "+mn-lt":
            return self.minor_font or name
        return name


# Maps raw XML scheme names → _THEME-style normalized keys
_XML_SCHEME_NAME: dict[str, str] = {
    "dk1": "DARK_1",    "dk2": "DARK_2",
    "lt1": "LIGHT_1",   "lt2": "LIGHT_2",
    # bg1/bg2 are aliases for lt1/lt2; tx1/tx2 are aliases for dk1/dk2 in OOXML
    "bg1": "LIGHT_1",   "bg2": "LIGHT_2",
    "tx1": "DARK_1",    "tx2": "DARK_2",
    "accent1": "ACCENT_1", "accent2": "ACCENT_2",
    "accent3": "ACCENT_3", "accent4": "ACCENT_4",
    "accent5": "ACCENT_5", "accent6": "ACCENT_6",
    "hlink": "HYPERLINK", "folHlink": "FOLLOWED_HYPERLINK",
}


def _extract_theme_fonts(presentation: Any) -> tuple[str | None, str | None]:
    """Return (major_latin, minor_latin) from the first slide master's theme."""
    try:
        from lxml import etree
        sm = presentation.slide_masters[0]
        for rel in sm.part.rels.values():
            if "theme" not in rel.reltype.lower():
                continue
            el = etree.fromstring(rel.target_part.blob)
            fs = _oxml_find_fontScheme(el)
            if fs is None:
                continue
            maj = _oxml_find_latin(_oxml_find_majorFont(fs))
            min_ = _oxml_find_latin(_oxml_find_minorFont(fs))
            return (
                maj.get("typeface") if maj is not None else None,
                min_.get("typeface") if min_ is not None else None,
            )
    except Exception:
        pass
    return (None, None)


def _extract_theme_colors(presentation: Any) -> dict[str, str]:
    """Extract the color palette from the first slide master's theme.

    Keys are stored under MULTIPLE aliases so any code path looking up a theme
    color resolves regardless of which naming convention it uses:

      - canonical: LIGHT_1, DARK_1, ACCENT_1...
      - python-pptx enum names: BACKGROUND_1, TEXT_1, ACCENT_1...
      - OOXML raw: bg1, tx1, accent1...

    Without these aliases, _fill_hex(...).fore_color.theme_color returns
    `BACKGROUND_1` from python-pptx but our dict only had `LIGHT_1`,
    causing every scheme-color slide background and shape fill to silently
    fall through to None (and Studio defaulted to white). Same applies to
    `TEXT_1` vs `DARK_1`.
    """
    result: dict[str, str] = {}
    # python-pptx enum-name aliases — these are what fill.fore_color.theme_color.name returns
    _PPTX_ENUM_ALIAS = {
        "LIGHT_1":  ["BACKGROUND_1", "bg1"],
        "LIGHT_2":  ["BACKGROUND_2", "bg2"],
        "DARK_1":   ["TEXT_1", "tx1"],
        "DARK_2":   ["TEXT_2", "tx2"],
        "ACCENT_1": ["accent1"],
        "ACCENT_2": ["accent2"],
        "ACCENT_3": ["accent3"],
        "ACCENT_4": ["accent4"],
        "ACCENT_5": ["accent5"],
        "ACCENT_6": ["accent6"],
        "HYPERLINK":         ["HYPERLINK", "hlink"],
        "FOLLOWED_HYPERLINK":["FOLLOWED_HYPERLINK", "folHlink"],
    }
    try:
        from lxml import etree
        sm = presentation.slide_masters[0]
        for rel in sm.part.rels.values():
            if "theme" not in rel.reltype.lower():
                continue
            el = etree.fromstring(rel.target_part.blob)
            cs = _oxml_find_clrScheme(el)
            if cs is None:
                continue
            for child in cs:
                xml_name = child.tag.split("}")[-1]
                normalized = _XML_SCHEME_NAME.get(xml_name, xml_name.upper())
                # Try srgbClr first, then sysClr lastClr
                hex_val = None
                rgb_el = _oxml_find_srgbClr(child)
                if rgb_el is not None:
                    hex_val = rgb_el.get("val")
                if hex_val is None:
                    sys_el = _oxml_find_sysClr(child)
                    if sys_el is not None:
                        hex_val = sys_el.get("lastClr")
                if hex_val:
                    color = "#" + hex_val.lstrip("#")
                    result[normalized] = color
                    # Also store under aliases for resilience
                    for alias in _PPTX_ENUM_ALIAS.get(normalized, []):
                        result[alias] = color
            return result
    except Exception:
        pass
    return result


_FONT_MAGIC = (b'\x00\x01\x00\x00', b'OTTO', b'true', b'typ1', b'wOFF', b'wOF2')


def _guid_to_obfuscation_key(guid_str: str) -> bytes | None:
    """Convert a GUID string to the 16-byte obfuscation key per ECMA-376 §22.6.1.1.

    GUID byte encoding: first three components little-endian, last two big-endian.
    """
    import struct
    h = guid_str.strip("{}").replace("-", "")
    if len(h) != 32:
        return None
    try:
        key = (
            struct.pack("<I", int(h[0:8], 16))
            + struct.pack("<H", int(h[8:12], 16))
            + struct.pack("<H", int(h[12:16], 16))
            + bytes.fromhex(h[16:20])
            + bytes.fromhex(h[20:32])
        )
        return key
    except Exception:
        return None


def _deobfuscate_font(font_bytes: bytes, key: bytes) -> bytes:
    """XOR the first 32 bytes of a font file with the 16-byte obfuscation key (applied twice)."""
    result = bytearray(font_bytes)
    for i in range(min(32, len(result))):
        result[i] ^= key[i % 16]
    return bytes(result)


def _extract_embedded_fonts(pptx_path: str | Path) -> list:
    """Extract embedded fonts from a PPTX zip. Returns list of EmbeddedFont.

    Obfuscated fonts (ECMA-376 §22.6.1.1) are deobfuscated using the document GUID
    from docProps/core.xml so they can be registered with matplotlib.
    """
    from zipfile import ZipFile
    from lxml import etree
    fonts: list = []
    try:
        with ZipFile(str(pptx_path)) as z:
            names = set(z.namelist())
            if "ppt/presentation.xml" not in names:
                return fonts

            # Try to obtain the document GUID for font deobfuscation
            obfuscation_key: bytes | None = None
            if "docProps/core.xml" in names:
                try:
                    core_root = etree.fromstring(z.read("docProps/core.xml"))
                    for child in core_root:
                        if child.tag.split("}")[-1] == "identifier" and child.text:
                            obfuscation_key = _guid_to_obfuscation_key(child.text.strip())
                            break
                except Exception:
                    pass

            prs_xml = z.read("ppt/presentation.xml")
            root = etree.fromstring(prs_xml)
            ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
            ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
            font_lst = root.find(f"{{{ns_p}}}embeddedFontLst")
            if font_lst is None:
                return fonts

            rels_path = "ppt/_rels/presentation.xml.rels"
            rid_to_target: dict[str, str] = {}
            if rels_path in names:
                rels_root = etree.fromstring(z.read(rels_path))
                for rel in rels_root:
                    rid = rel.get("Id", "")
                    target = rel.get("Target", "")
                    if "font" in rel.get("Type", "").lower():
                        rid_to_target[rid] = f"ppt/{target}" if not target.startswith("ppt/") else target

            for ef in font_lst.findall(f"{{{ns_p}}}embeddedFont"):
                font_el = ef.find(f"{{{ns_p}}}font")
                typeface = font_el.get("typeface", "Unknown") if font_el is not None else "Unknown"
                for style in ("regular", "bold", "italic", "boldItalic"):
                    style_el = ef.find(f"{{{ns_p}}}{style}")
                    if style_el is None:
                        continue
                    rid = style_el.get(f"{{{ns_r}}}id", "")
                    target = rid_to_target.get(rid, "")
                    if not target or target not in names:
                        continue
                    raw = z.read(target)
                    is_obf = len(raw) >= 4 and raw[:4] not in _FONT_MAGIC
                    font_data = raw
                    if is_obf and obfuscation_key:
                        deobf = _deobfuscate_font(raw, obfuscation_key)
                        if len(deobf) >= 4 and deobf[:4] in _FONT_MAGIC:
                            font_data = deobf
                            is_obf = False
                    fonts.append(EmbeddedFont(
                        typeface=typeface,
                        style=style,
                        font_bytes=font_data,
                        is_obfuscated=is_obf,
                    ))
    except Exception:
        pass
    return fonts


def _resolve_bg_color(slide: Any, theme_colors: dict[str, str]) -> str | None:
    """Walk slide → layout → master fill chain; return resolved hex or None (→ white)."""
    def _gradient_first_stop_hex(fill: Any) -> str | None:
        """Extract dominant (first) stop of a gradient fill as an approximate solid color."""
        try:
            gs_lst = _oxml_find_gsLst(_oxml_find_gradFill(fill._fill))
            if gs_lst is None:
                gs_lst = _oxml_find_gsLst_descendant(fill._fill)
            if gs_lst is None:
                return None
            stops = list(gs_lst)
            if not stops:
                return None
            # Use first stop (position=0) to get the "top" color
            first = stops[0]
            rgb_el = _oxml_find_descendant_a(first, "srgbClr")
            if rgb_el is not None:
                return "#" + rgb_el.get("val", "")
            sys_el = _oxml_find_descendant_a(first, "sysClr")
            if sys_el is not None:
                return "#" + sys_el.get("lastClr", "")
            scheme_el = _oxml_find_descendant_a(first, "schemeClr")
            if scheme_el is not None:
                key = scheme_el.get("val", "")
                normalized = _XML_SCHEME_NAME.get(key, key.upper())
                return theme_colors.get(normalized)
        except Exception:
            pass
        return None

    def _fill_hex(fill: Any) -> str | None:
        try:
            ftype = fill.type
        except Exception:
            return None
        if ftype is None:
            return None
        ftype_int = int(ftype) if ftype is not None else None
        # MSO_FILL_TYPE.BACKGROUND = 5 → inherit/transparent
        if ftype_int == 5:
            return None
        # MSO_FILL_TYPE.SOLID = 1
        if ftype_int == 1:
            try:
                rgb = fill.fore_color.rgb
                return "#" + str(rgb)
            except Exception:
                pass
            # Scheme color (theme color)
            try:
                tc = fill.fore_color.theme_color
                if tc is not None:
                    key = getattr(tc, "name", None) or str(tc)
                    return theme_colors.get(key)
            except Exception:
                pass
        # MSO_FILL_TYPE.GRADIENT = 3 — return first stop as approximate
        if ftype_int == 3:
            return _gradient_first_stop_hex(fill)
        return None

    try:
        layout = slide.slide_layout
        master = layout.slide_master
        for source in (slide, layout, master):
            try:
                c = _fill_hex(source.background.fill)
                if c is not None:
                    return c
            except Exception:
                pass
    except Exception:
        pass
    return None


def _resolve_bg_gradient(slide: Any, theme_colors: dict[str, str]) -> tuple[list, float]:
    """Return (gradient_stops, angle_deg) for a gradient slide background, or ([], 0.0)."""
    try:
        for source in (slide, slide.slide_layout, slide.slide_layout.slide_master):
            try:
                fill = source.background.fill
                if int(fill.type) != 3:  # not gradient
                    continue
                fill_el = fill._fill
                gs_lst = _oxml_find_gsLst_descendant(fill_el)
                if gs_lst is None:
                    continue
                stops = []
                for gs in gs_lst:
                    pos = int(gs.get("pos", "0")) / 100000.0
                    cs = _extract_color_spec(gs)
                    if cs is not None:
                        stops.append(GradientStop(position=pos, color=cs))
                # Extract angle from <a:lin>
                angle_deg = 0.0
                lin_el = _oxml_find_descendant_a(fill_el, "lin")
                if lin_el is not None:
                    ang_raw = int(lin_el.get("ang", "0"))
                    # OOXML: angle in 60000ths of a degree, measured from 3-o'clock CW
                    # Convert back: PowerPoint angle = (90 - ooxml_deg) % 360
                    ooxml_deg = ang_raw / 60000.0
                    angle_deg = (90.0 - ooxml_deg) % 360.0
                if stops:
                    return stops, angle_deg
            except Exception:
                pass
    except Exception:
        pass
    return [], 0.0


def _layout_default_text_color(slide: Any, theme_colors: dict[str, str]) -> str | None:
    """
    Return the default text color for runs on this slide, resolved from the
    layout's body placeholder lstStyle.  Returns a hex string or None.

    PowerPoint color cascade:
      run.font → para defaults → txBody lstStyle → layout lstStyle → master lstStyle → theme
    This reads the *layout's body placeholder lstStyle*, which is the lowest-effort
    level that covers the common dark-bg / light-text pattern (e.g. Quote layouts).
    """
    try:
        layout = slide.slide_layout
        for ph in layout.placeholders:
            try:
                # Look at the body content placeholder (idx 10 in Snowflake)
                if int(ph.placeholder_format.type) != 2:  # BODY = 2
                    continue
                txBody = ph.text_frame._txBody
                lst = _oxml_find_lstStyle(txBody)
                if lst is None:
                    continue
                # Walk all default/level paragraph properties for a solid fill color
                for el in lst.iter():
                    tag = el.tag.split("}")[-1]
                    if tag not in ("defRPr", "r", "rPr"):
                        continue
                    solid = _oxml_find_solidFill(el)
                    if solid is None:
                        continue
                    rgb_el = _oxml_find_srgbClr(solid)
                    if rgb_el is not None:
                        return "#" + rgb_el.get("val", "").lstrip("#")
                    sch_el = _oxml_find_schemeClr(solid)
                    if sch_el is not None:
                        xml_name = sch_el.get("val", "")
                        key = _XML_SCHEME_NAME.get(xml_name, xml_name.upper())
                        return theme_colors.get(key)
                break
            except Exception:
                pass
    except Exception:
        pass
    return None


def _onboard_inherited_shapes(
    slide: Any, slide_number: int, ctx: _OnboardContext,
) -> list:
    """
    Return bridge elements for shapes inherited from the slide layout (and master).
    Includes non-placeholder decorative shapes AND layout placeholder shapes that are
    NOT overridden by slide content (e.g. background fill rectangles on section dividers).

    Elements are given negative z-indices so they render behind all slide content.
    """
    elements: list = []

    try:
        layout = slide.slide_layout
        master = layout.slide_master
    except Exception:
        return elements

    # Collect placeholder idx values that the slide itself provides content for
    slide_placeholder_idxs: set[int] = set()
    for shape in slide.shapes:
        if safe_get(lambda s=shape: s.is_placeholder, False):
            idx = safe_get(lambda s=shape: s.placeholder_format.idx)
            if idx is not None:
                slide_placeholder_idxs.add(idx)

    def _collect(source: Any, z_base: int, include_unoveridden_placeholders: bool = False) -> None:
        idx = 0
        for shape in source.shapes:
            if safe_get(lambda s=shape: s.is_placeholder, False):
                if not include_unoveridden_placeholders:
                    continue
                # Include layout placeholders only if not overridden by a slide shape
                ph_idx = safe_get(lambda s=shape: s.placeholder_format.idx)
                if ph_idx in slide_placeholder_idxs:
                    continue
                # Skip layout placeholders that have text — layout placeholder text
                # is always prompt/hint text ("Add subtitle", "Insert Chart", etc.)
                # and should never appear as real content in the rebuilt slide.
                if safe_get(lambda s=shape: bool(s.text_frame.text.strip()), False):
                    continue
            for el in _onboard_shape_tree(shape, slide_number, ctx):
                el.stacking.z_index = z_base + idx
                idx += 1
                elements.append(el)

    # Master non-placeholder shapes first (furthest back), then layout shapes
    _collect(master, -2000, include_unoveridden_placeholders=False)
    _collect(layout, -1000, include_unoveridden_placeholders=True)

    return elements


def onboard_pptx(pptx_path: str | Path) -> PercyDocument:
    path = Path(pptx_path)
    presentation = Presentation(str(path))
    major_font, minor_font = _extract_theme_fonts(presentation)
    theme_colors = _extract_theme_colors(presentation)
    ctx = _OnboardContext(
        major_font=major_font,
        minor_font=minor_font,
        theme_colors=theme_colors,
    )
    document = PercyDocument(
        source_path=str(path),
        theme_colors=dict(theme_colors),
        metadata=PresentationMetadata(
            slide_width=emu_to_inches(presentation.slide_width),
            slide_height=emu_to_inches(presentation.slide_height),
            slide_count=len(presentation.slides),
            source_path=str(path),
        ),
        custom_properties={
            "slide_width": emu_to_inches(presentation.slide_width),
            "slide_height": emu_to_inches(presentation.slide_height),
        },
    )
    document.embedded_fonts = _extract_embedded_fonts(path)

    for slide_number, slide in enumerate(presentation.slides, start=1):
        ctx.slide_part = safe_get(lambda: slide.part)
        bg_color = _resolve_bg_color(slide, theme_colors)
        bg_grad_stops, bg_grad_angle = _resolve_bg_gradient(slide, theme_colors)
        bridge_slide = BridgeSlide(
            slide_number=slide_number,
            width=emu_to_inches(presentation.slide_width),
            height=emu_to_inches(presentation.slide_height),
            background_color=bg_color,
            background_gradient_stops=bg_grad_stops,
            background_gradient_angle=bg_grad_angle,
            default_text_color=_layout_default_text_color(slide, theme_colors),
        )
        # Inherited (layout + master) non-placeholder shapes go first (low z)
        for element in _onboard_inherited_shapes(slide, slide_number, ctx):
            bridge_slide.elements.append(element)
        # Slide-own shapes on top
        for shape in slide.shapes:
            for element in _onboard_shape_tree(shape, slide_number, ctx):
                bridge_slide.elements.append(element)
                _capture_metadata_element(document, element)
        document.slides.append(bridge_slide)

    return document


def _group_xfrm_transform(shape: Any, parent_transform: tuple) -> tuple:
    """Compute the accumulated affine transform for children of this group.

    Transform is (off_x_emu, off_y_emu, scale_x, scale_y) where:
        slide_x = off_x + child_x * scale_x
        slide_y = off_y + child_y * scale_y
        slide_cx = child_cx * scale_x
        slide_cy = child_cy * scale_y

    OOXML groups define a child coordinate system via off/ext/chOff/chExt:
        parent_x = off_x + (child_x - chOff_x) * (ext_cx / chExt_cx)
    This is composed with parent_transform to get slide-absolute coordinates.
    """
    parent_off_x, parent_off_y, parent_sx, parent_sy = parent_transform
    try:
        _A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        _P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        el = shape._element
        grpSpPr = el.find(f"{{{_P}}}grpSpPr")
        if grpSpPr is None:
            return parent_transform
        xfrm = grpSpPr.find(f"{{{_A}}}xfrm")
        if xfrm is None:
            return parent_transform
        off = xfrm.find(f"{{{_A}}}off")
        ext_el = xfrm.find(f"{{{_A}}}ext")
        chOff = xfrm.find(f"{{{_A}}}chOff")
        chExt_el = xfrm.find(f"{{{_A}}}chExt")
        if off is None or chOff is None:
            return parent_transform
        off_x = int(off.get("x", 0))
        off_y = int(off.get("y", 0))
        chOff_x = int(chOff.get("x", 0))
        chOff_y = int(chOff.get("y", 0))
        # Scale: how child-coord units map to parent-coord units
        sx = 1.0
        sy = 1.0
        if ext_el is not None and chExt_el is not None:
            ext_cx = int(ext_el.get("cx", 0))
            ext_cy = int(ext_el.get("cy", 0))
            chExt_cx = int(chExt_el.get("cx", 1)) or 1
            chExt_cy = int(chExt_el.get("cy", 1)) or 1
            sx = ext_cx / chExt_cx
            sy = ext_cy / chExt_cy
        # Compose: slide_x = parent_off_x + parent_sx * (off_x + (child_x - chOff_x) * sx)
        #                   = (parent_off_x + parent_sx*(off_x - chOff_x*sx)) + (parent_sx*sx)*child_x
        new_off_x = parent_off_x + parent_sx * (off_x - chOff_x * sx)
        new_off_y = parent_off_y + parent_sy * (off_y - chOff_y * sy)
        return (new_off_x, new_off_y, parent_sx * sx, parent_sy * sy)
    except Exception:
        return parent_transform


def _onboard_shape_tree(
    shape: Any, slide_number: int, ctx: _OnboardContext,
    group_id: str | None = None, group_path: list[str] | None = None,
    _group_transform: tuple = (0, 0, 1.0, 1.0),
    _parent_group_fill: "ColorSpec | None" = None,
):
    if _is_group_shape(shape):
        current_group_id = group_id or f"slide-{slide_number}:group-{safe_get(lambda: shape.shape_id)}"
        current_path = [*(group_path or []), current_group_id]
        child_transform = _group_xfrm_transform(shape, _group_transform)
        # Resolve this group's fill. If it has grpFill itself, inherit from parent.
        group_fill = _group_fill_color(shape, ctx.theme_colors)
        if group_fill is None:
            try:
                _P = "http://schemas.openxmlformats.org/presentationml/2006/main"
                _A = "http://schemas.openxmlformats.org/drawingml/2006/main"
                grpSpPr = shape._element.find(f"{{{_P}}}grpSpPr")
                if grpSpPr is not None and grpSpPr.find(f"{{{_A}}}grpFill") is not None:
                    group_fill = _parent_group_fill
            except Exception:
                pass
        effective_group_fill = group_fill or _parent_group_fill
        elements = []
        for child in shape.shapes:
            elements.extend(_onboard_shape_tree(child, slide_number, ctx, current_group_id, current_path, child_transform, effective_group_fill))
        return elements

    element = _onboard_shape(shape, slide_number, ctx, group_id=group_id, group_path=group_path or [], _group_transform=_group_transform, _parent_group_fill=_parent_group_fill)
    return [element]


def _onboard_shape(
    shape: Any, slide_number: int, ctx: _OnboardContext,
    group_id: str | None = None, group_path: list[str] | None = None,
    _group_transform: tuple = (0, 0, 1.0, 1.0),
    # Legacy alias kept for any callers that pass _group_offset
    _group_offset: tuple | None = None,
    _parent_group_fill: "ColorSpec | None" = None,
):
    if _group_offset is not None:
        _group_transform = (_group_offset[0], _group_offset[1], 1.0, 1.0)
    _gt_off_x, _gt_off_y, _gt_sx, _gt_sy = _group_transform
    # Resolve group fill inheritance: if shape has <a:grpFill>, use parent group's fill.
    _has_grp_fill = _shape_has_grp_fill(shape)
    _resolved_grp_fill = _parent_group_fill if _has_grp_fill else None
    placeholder = placeholder_info(shape)
    _rot_raw = safe_get(lambda: shape.rotation)  # degrees, CW in PPTX
    _flip_h  = safe_get(lambda: shape.element.spPr.xfrm.flipH) or False
    _flip_v  = safe_get(lambda: shape.element.spPr.xfrm.flipV) or False
    _raw_left = safe_get(lambda: shape.left) or 0
    _raw_top  = safe_get(lambda: shape.top) or 0
    _raw_width  = safe_get(lambda: shape.width) or 0
    _raw_height = safe_get(lambda: shape.height) or 0
    common = {
        "position": Position(
            left=emu_to_inches(_gt_off_x + _raw_left * _gt_sx),
            top=emu_to_inches(_gt_off_y + _raw_top * _gt_sy),
            width=emu_to_inches(_raw_width * _gt_sx) if _raw_width else 0.0,
            height=emu_to_inches(_raw_height * _gt_sy) if _raw_height else 0.0,
        ),
        "transforms": Transform(
            rotation=float(_rot_raw) if _rot_raw is not None else 0.0,
            flip_h=bool(_flip_h),
            flip_v=bool(_flip_v),
        ),
        "identification": Identification(
            slide_number=slide_number,
            shape_name=safe_get(lambda: shape.name),
            shape_id=safe_get(lambda: shape.shape_id),
            group_id=group_id,
        ),
        "custom_properties": {
            "source_shape_type": enum_name(safe_get(lambda: shape.shape_type)),
            "source_shape_id": safe_get(lambda: shape.shape_id),
            "source_shape_name": safe_get(lambda: shape.name),
            "onboard_status": "semantic-best-effort",
            "placeholder": placeholder,
            "group_id": group_id,
            "group_path": group_path or [],
            "semantic_role": _semantic_role(shape, placeholder),
            "semantic_debt": [],
        },
    }

    shape_type = enum_name(safe_get(lambda: shape.shape_type))

    if safe_get(lambda: shape.has_chart, False):
        chart = shape.chart
        unsupported = _chart_unsupported_features(chart)
        common["custom_properties"]["semantic_debt"].extend(unsupported)
        _tc = ctx.theme_colors
        _pal = _chart_plot_area_layout(chart)
        _txpr = _chart_space_txpr(chart, _tc)
        return BridgeChart(
            **common,
            chart_type=enum_name(safe_get(lambda: chart.chart_type)),
            title=_chart_title(chart, _tc),
            categories=_chart_categories(chart),
            series=_chart_series(chart, _tc),
            category_axis=_chart_axis(chart, "category", _tc),
            value_axis=_chart_axis(chart, "value", _tc),
            plot_properties=_chart_plot_properties(chart, _tc),
            legend=_chart_legend(chart, _tc),
            chart_space_fill=_chart_space_fill(chart, _tc),
            figsize=(common["position"].width, common["position"].height),
            data_source=_chart_data_source(chart),
            overlay_files=_chart_overlay_files(chart),
            reconstruction_blobs=ReconstructionBlobs(),
            plot_area_x=_pal[0], plot_area_y=_pal[1],
            plot_area_w=_pal[2], plot_area_h=_pal[3],
            plot_area_x_mode=_pal[4], plot_area_y_mode=_pal[5],
            plot_area_layout_target=_pal[6],
            chart_txpr_font_name=_txpr[0],
            chart_txpr_font_size=_txpr[1],
            chart_txpr_font_bold=_txpr[2],
            chart_txpr_font_color=_txpr[3],
            disp_blanks_as=_chart_disp_blanks_as(chart),
        )

    if safe_get(lambda: shape.has_table, False):
        table = shape.table
        unsupported = _table_unsupported_features(table)
        common["custom_properties"]["semantic_debt"].extend(unsupported)
        return BridgeTable(
            **common,
            data=[[_cell_typed_value(cell) for cell in row.cells] for row in table.rows],
            cell_formats=_table_cell_formats(table, shape, ctx),
            dimensions=TableDimensions(
                column_widths=[emu_to_inches(safe_get(lambda col=col: col.width)) or 0.0 for col in table.columns],
                row_heights=[emu_to_inches(safe_get(lambda row=row: row.height)) or 0.0 for row in table.rows],
            ),
            table_properties=_table_properties(table, shape),
            defaults=_table_defaults(table),
        )

    image = safe_get(lambda: shape.image, None)
    if image is not None:
        _img_geom, _img_geom_adj = _image_shape_geometry(shape)
        return BridgeImage(
            **common,
            image_data=ImageData(
                image_bytes=safe_get(lambda: image.blob),
                image_format=safe_get(lambda: image.ext.upper()),
            ),
            file_info=ImageFileInfo(original_filename=safe_get(lambda: image.filename)),
            dimensions=ImageDimensions(
                width_px=safe_get(lambda: image.size[0]),
                height_px=safe_get(lambda: image.size[1]),
                dpi=safe_get(lambda: image.dpi[0]),
            ),
            cropping=ImageCropping(
                crop_left=_image_crop(shape, "l"),
                crop_right=_image_crop(shape, "r"),
                crop_top=_image_crop(shape, "t"),
                crop_bottom=_image_crop(shape, "b"),
            ),
            border=ImageBorder(
                has_border=_line_visible(shape),
                border_color=_line_color(shape),
                border_width=_line_width_points(shape),
            ),
            shadow=_extract_outer_shadow(shape, ctx.theme_colors),
            shape_geometry=_img_geom,
            shape_geometry_adj=_img_geom_adj,
        )

    if shape_type == "LINE":
        # Use the group-transformed position for endpoints — begin_x/end_x from pptx-python
        # returns local group coordinates for connectors inside groups, not slide coords.
        _pos = common["position"]
        _flip_h = common["transforms"].flip_h
        _flip_v = common["transforms"].flip_v
        # Determine which corners are start/end based on flip flags
        if _flip_h and _flip_v:
            _sx, _sy = _pos.left + _pos.width, _pos.top + _pos.height
            _ex, _ey = _pos.left, _pos.top
        elif _flip_h:
            _sx, _sy = _pos.left + _pos.width, _pos.top
            _ex, _ey = _pos.left, _pos.top + _pos.height
        elif _flip_v:
            _sx, _sy = _pos.left, _pos.top + _pos.height
            _ex, _ey = _pos.left + _pos.width, _pos.top
        else:
            _sx, _sy = _pos.left, _pos.top
            _ex, _ey = _pos.left + _pos.width, _pos.top + _pos.height
        return BridgeConnector(
            **common,
            connector_type=_connector_type(shape),
            endpoints=ConnectorEndpoints(
                start_x=_sx, start_y=_sy, end_x=_ex, end_y=_ey,
            ),
            line=_shape_line(shape),
        )

    if shape_type == "FREEFORM":
        paths = _freeform_paths(shape)
        if not paths:
            common["custom_properties"]["semantic_debt"].append("freeform_path_commands")
        return BridgeFreeform(
            **common,
            paths=paths,
            geometry_xml=_freeform_geometry_xml(shape),
            fill=FreeformFill(
                fill_type="solidFill" if _resolved_grp_fill else _fill_type(shape),
                fill_color=_resolved_grp_fill or _fill_color(shape, ctx.theme_colors),
                fill_scheme=None if _resolved_grp_fill else _fill_scheme(shape),
                transparency=_fill_transparency(shape) or 0.0,
                gradient_angle=_fill_gradient_angle(shape),
                gradient_stops=_fill_gradient_stops(shape, ctx.theme_colors),
                pattern_preset=_fill_pattern_preset(shape),
                bg_color=_fill_bg_color(shape, ctx.theme_colors),
            ),
            line=FreeformLine(
                line_visible=_line_visible(shape),
                line_color=_line_color(shape, ctx.theme_colors),
                line_scheme=_line_scheme(shape),
                line_width=_line_width_points(shape),
                line_dash=_line_dash(shape),
                line_cap=_line_cap(shape),
                line_join=_line_join(shape),
            ),
            transform_emus=TransformEmus(
                offset_x=safe_get(lambda: shape.element.spPr.xfrm.off.x),
                offset_y=safe_get(lambda: shape.element.spPr.xfrm.off.y),
                extent_cx=safe_get(lambda: shape.element.spPr.xfrm.ext.cx),
                extent_cy=safe_get(lambda: shape.element.spPr.xfrm.ext.cy),
            ),
            description=safe_get(lambda: shape.element.nvSpPr.cNvPr.descr),
        )

    if _is_text_element(shape, shape_type):
        return BridgeText(
            **common,
            paragraphs=_text_paragraphs(shape, ctx),
            text_frame=_text_frame(shape),
            fill_and_border=FillAndBorder(
                fill_color=_resolved_grp_fill or _fill_color(shape, ctx.theme_colors),
                has_fill=bool(_resolved_grp_fill) or _shape_has_explicit_fill(shape),
                border_color=_line_color(shape, ctx.theme_colors),
                border_width=_line_width_points(shape),
                has_border=_line_visible(shape),
            ),
            shape_info=ShapeInfo(
                shape_type="textbox",
                is_placeholder=placeholder["is_placeholder"],
                placeholder_type=placeholder.get("type"),
                placeholder_idx=placeholder.get("idx"),
            ),
            lst_style_xml=None,
            shadow=_extract_outer_shadow(shape, ctx.theme_colors),
        )

    return BridgeShape(
        **common,
        shape_identification=ShapeIdentification(
            shape_type=shape_type or "auto_shape",
            geometry_preset=_geometry_preset(shape) or "rect",
            geometry_adjustments=_geometry_adjustments(shape),
        ),
        fill=ShapeFill(
            fill_type="solidFill" if _resolved_grp_fill else _fill_type(shape),
            color=_resolved_grp_fill or _fill_color(shape, ctx.theme_colors),
            transparency=_fill_transparency(shape) or 0.0,
            gradient_angle=_fill_gradient_angle(shape),
            gradient_stops=_fill_gradient_stops(shape, ctx.theme_colors),
            pattern_preset=_fill_pattern_preset(shape),
            bg_color=_fill_bg_color(shape, ctx.theme_colors),
        ),
        line=_shape_line(shape, ctx.theme_colors),
        text_content=ShapeTextContent(
            has_text=bool(safe_get(lambda: shape.text, "")),
            text_content=safe_get(lambda: shape.text),
            paragraphs=_text_paragraphs(shape, ctx) if safe_get(lambda: shape.has_text_frame, False) else [],
        ),
        text_frame=_shape_text_frame(shape),
        shadow=_extract_outer_shadow(shape, ctx.theme_colors),
    )


def _capture_metadata_element(document: PercyDocument, element: Any) -> None:
    role = element.custom_properties.get("semantic_role")
    item = {
        "slide_number": element.identification.slide_number,
        "shape_id": element.identification.shape_id,
        "shape_name": element.identification.shape_name,
        "element_type": element.element_type,
        "text": _element_text(element),
    }
    if role == "page_number":
        document.metadata.page_number_elements.append(item)
    elif role == "footer":
        document.metadata.footer_elements.append(item)


def _element_text(element: Any) -> str | None:
    if isinstance(element, BridgeText):
        return "\n".join("".join(run.text for run in paragraph.runs) for paragraph in element.paragraphs)
    return None


def _semantic_role(shape: Any, placeholder: dict[str, Any]) -> str | None:
    name = (safe_get(lambda: shape.name, "") or "").lower()
    placeholder_type = placeholder.get("type")
    if placeholder_type == "SLIDE_NUMBER" or "slide number" in name or "pagenumber" in name or "page number" in name:
        return "page_number"
    if placeholder_type in {"FOOTER", "DATE"} or "footer" in name or "date" in name:
        return "footer"
    return None


def _chart_title(chart: Any, theme_colors: "dict | None" = None) -> ChartTitle:
    title = ChartTitle(auto_title_deleted=safe_get(lambda: chart._chartSpace.chart.autoTitleDeleted.val))
    if not safe_get(lambda: chart.has_title, False):
        return title
    chart_title = safe_get(lambda: chart.chart_title)
    title.title = _text_frame_text(safe_get(lambda: chart_title.text_frame))
    title.title_font_name = safe_get(lambda: chart_title.text_frame.paragraphs[0].font.name)
    title.title_font_size = _points(safe_get(lambda: chart_title.text_frame.paragraphs[0].font.size))
    title.title_font_bold = safe_get(lambda: chart_title.text_frame.paragraphs[0].font.bold)
    title.title_font_italic = safe_get(lambda: chart_title.text_frame.paragraphs[0].font.italic)
    title.title_font_color = _font_color(safe_get(lambda: chart_title.text_frame.paragraphs[0].font), theme_colors)
    try:
        layout_el = _oxml_find_c(chart_title._element, "layout")
        manual_el = _oxml_find_c(layout_el, "manualLayout")
        if manual_el is not None:
            def _ml_val(tag: str) -> "float | None":
                e = _oxml_find_c(manual_el, tag)
                return float(e.get("val")) if e is not None and e.get("val") is not None else None
            title.title_position_x = _ml_val("x")
            title.title_position_y = _ml_val("y")
            title.title_width = _ml_val("w")
            title.title_height = _ml_val("h")
    except Exception:
        pass
    return title


def _chart_categories(chart: Any) -> ChartCategories:
    plot = _first_chart_plot(chart)
    categories = safe_get(lambda: list(plot.categories), []) or []
    raw = [str(category) for category in categories]
    return ChartCategories(
        categories=raw,
        categories_raw=raw,
        categories_are_numeric=all(_is_number(category) for category in categories) if categories else False,
        category_levels=_category_levels(safe_get(lambda: plot.categories)),
    )


def _category_levels(categories: Any) -> list[dict[str, Any]]:
    levels = safe_get(lambda: categories.levels, []) or []
    return [
        {"level": level_index, "items": [{"index": index, "label": str(label)} for index, label in level]}
        for level_index, level in enumerate(levels)
    ]


def _chart_series(chart: Any, theme_colors: "dict | None" = None) -> list[ChartSeries]:
    bridge_series = []
    plots = safe_get(lambda: list(chart.plots), []) or []
    for plot_index, plot in enumerate(plots):
        for series in safe_get(lambda: plot.series, []) or []:
            # Smoothing (line charts)
            smooth = False
            try:
                smooth_el = _oxml_find_c(series._element, "smooth")
                if smooth_el is not None:
                    smooth = smooth_el.get("val", "0") not in ("0", "false")
            except Exception:
                pass

            bridge_series.append(
                ChartSeries(
                    name=safe_get(lambda series=series: series.name),
                    values=[_float_or_none(value) for value in safe_get(lambda series=series: list(series.values), [])],
                    color=_chart_format_fill_color(safe_get(lambda series=series: series.format), theme_colors),
                    negative_color=_series_negative_color(series, theme_colors),
                    point_colors=_point_colors(series, theme_colors),
                    plot_type=_local_name(safe_get(lambda series=series: series._element.getparent().tag, "")),
                    plot_index=plot_index,
                    invert_if_negative=safe_get(lambda series=series: series.invert_if_negative, False),
                    line=_chart_line_format(safe_get(lambda series=series: series.format.line), theme_colors),
                    marker=_marker_format(safe_get(lambda series=series: series.marker), theme_colors),
                    data_labels=_data_labels(safe_get(lambda series=series: series._element.dLbls), theme_colors),
                    x_values=[_float_or_none(value) for value in _series_x_values(series)],
                    point_formatting=_point_formatting(series, theme_colors),
                    custom_labels=_custom_data_labels(series),
                    smooth=smooth,
                    fill_type=_chart_fill_type(safe_get(lambda series=series: series._element.spPr)),
                    gradient_stops=_series_gradient_stops(series),
                )
            )
    return bridge_series


def _axis_minor_gridlines(axis: Any, theme_colors: "dict | None" = None) -> "Gridlines":
    """Build Gridlines bridge object for minor gridlines on an axis."""
    from percy.bridge import Gridlines as _GL
    try:
        minor_el = safe_get(lambda: _oxml_find_c(axis._element, "minorGridlines"))
        if minor_el is None:
            return _GL(has_major_gridlines=False)
        line = safe_get(lambda: axis.minor_gridlines.format.line)
        return _GL(
            has_major_gridlines=True,
            gridline_style=None,
            gridline_color=_chart_line_color(line, theme_colors),
            gridline_width=_chart_line_width(line),
        )
    except Exception:
        return _GL(has_major_gridlines=False)


def _chart_axis(chart: Any, axis_kind: str, theme_colors: "dict | None" = None) -> BridgeAxis:
    axis = safe_get(lambda: chart.category_axis if axis_kind == "category" else chart.value_axis)
    if axis is None:
        return BridgeAxis(visible=False, axis_type=axis_kind)

    # Crosses / crossesAt
    crosses = None
    crosses_at = None
    try:
        crosses_el = _oxml_find_c(axis._element, "crosses")
        if crosses_el is not None:
            crosses = crosses_el.get("val")
        crosses_at_el = _oxml_find_c(axis._element, "crossesAt")
        if crosses_at_el is not None:
            try:
                crosses_at = float(crosses_at_el.get("val", "0"))
            except (ValueError, TypeError):
                pass
    except Exception:
        pass

    # delete attribute
    delete = False
    try:
        delete_el = _oxml_find_c(axis._element, "delete")
        if delete_el is not None:
            delete = delete_el.get("val", "0") not in ("0", "false")
    except Exception:
        pass

    # numFmt on the axis itself (for value axis; separate from tick_labels.number_format)
    axis_num_fmt = None
    try:
        nf = _oxml_find_c(axis._element, "numFmt")
        if nf is not None:
            axis_num_fmt = nf.get("formatCode")
    except Exception:
        pass

    # axPos — where the axis is placed (t/b/l/r)
    ax_pos = None
    try:
        axpos_el = _oxml_find_c(axis._element, "axPos")
        if axpos_el is not None:
            ax_pos = axpos_el.get("val")
    except Exception:
        pass

    # noMultiLvlLbl — suppress multi-level category labels
    no_multi_lvl_lbl = False
    try:
        nml_el = _oxml_find_c(axis._element, "noMultiLvlLbl")
        if nml_el is not None:
            no_multi_lvl_lbl = nml_el.get("val", "0") not in ("0", "false")
    except Exception:
        pass

    # lblOffset — offset of axis tick labels (100 = 100%, default)
    lbl_offset = None
    try:
        lo_el = _oxml_find_c(axis._element, "lblOffset")
        if lo_el is not None:
            lbl_offset = int(lo_el.get("val", 100))
    except Exception:
        pass

    # lblAlgn — label alignment (ctr/l/r)
    lbl_algn = None
    try:
        la_el = _oxml_find_c(axis._element, "lblAlgn")
        if la_el is not None:
            lbl_algn = la_el.get("val")
    except Exception:
        pass

    # crossBetween — where value axis crosses category axis (between/midCat)
    cross_between = None
    try:
        cb_el = _oxml_find_c(axis._element, "crossBetween")
        if cb_el is not None:
            cross_between = cb_el.get("val")
    except Exception:
        pass

    return BridgeAxis(
        visible=safe_get(lambda: axis.visible, True),
        axis_type=_local_name(safe_get(lambda: axis._element.tag, axis_kind)),
        min_value=safe_get(lambda: axis.minimum_scale),
        max_value=safe_get(lambda: axis.maximum_scale),
        gridlines=_axis_gridlines(axis, theme_colors),
        minor_gridlines=_axis_minor_gridlines(axis, theme_colors),
        title=_axis_title(axis),
        tick_labels=_axis_tick_labels(axis),
        tick_marks=_axis_tick_marks(axis),
        units=AxisUnits(
            major_unit=safe_get(lambda: axis.major_unit),
            minor_unit=safe_get(lambda: axis.minor_unit),
            major_time_unit=enum_name(safe_get(lambda: axis._element.majorTimeUnit.val)),
            minor_time_unit=enum_name(safe_get(lambda: axis._element.minorTimeUnit.val)),
            base_time_unit=enum_name(safe_get(lambda: axis._element.baseTimeUnit.val)),
        ),
        axis_line=AxisLine(
            line_visible=safe_get(lambda: axis.format.line.fill.type) is not None,
            line_color=_chart_line_color(safe_get(lambda: axis.format.line), theme_colors),
            line_width=_chart_line_width(safe_get(lambda: axis.format.line)),
        ),
        reverse_order=safe_get(lambda: axis.reverse_order, False) or False,
        number_format=axis_num_fmt,
        crosses=crosses,
        crosses_at=crosses_at,
        delete=delete,
        ax_pos=ax_pos,
        no_multi_lvl_lbl=no_multi_lvl_lbl,
        lbl_offset=lbl_offset,
        lbl_algn=lbl_algn,
        cross_between=cross_between,
    )


def _axis_title(axis: Any) -> AxisTitle:
    if not safe_get(lambda: axis.has_title, False):
        return AxisTitle()
    title = safe_get(lambda: axis.axis_title)
    return AxisTitle(
        title_text=_text_frame_text(safe_get(lambda: title.text_frame)),
        title_font_name=safe_get(lambda: title.text_frame.paragraphs[0].font.name),
        title_font_size=_points(safe_get(lambda: title.text_frame.paragraphs[0].font.size)),
        title_font_bold=safe_get(lambda: title.text_frame.paragraphs[0].font.bold),
    )


def _axis_tick_labels(axis: Any) -> Any:
    tick_labels = safe_get(lambda: axis.tick_labels)
    from percy.bridge import TickLabels

    return TickLabels(
        number_format=safe_get(lambda: tick_labels.number_format),
        tick_label_font_name=safe_get(lambda: tick_labels.font.name),
        tick_label_font_size=_points(safe_get(lambda: tick_labels.font.size)),
        tick_label_font_bold=safe_get(lambda: tick_labels.font.bold),
        tick_label_font_color=_font_color(safe_get(lambda: tick_labels.font)),
        tick_label_position=enum_name(safe_get(lambda: axis.tick_label_position)),
        tick_label_rotation=_bodypr_rotation(safe_get(lambda: axis._element.txPr.bodyPr)),
        tick_label_bodypr_attrs=dict(safe_get(lambda: axis._element.txPr.bodyPr.attrib, {}) or {}),
    )


def _axis_tick_marks(axis: Any) -> Any:
    from percy.bridge import TickMarks

    return TickMarks(
        major_tick_mark=enum_name(safe_get(lambda: axis.major_tick_mark)),
        minor_tick_mark=enum_name(safe_get(lambda: axis.minor_tick_mark)),
        tick_label_skip=safe_get(lambda: axis._element.tickLblSkip.val),
        tick_mark_skip=safe_get(lambda: axis._element.tickMarkSkip.val),
    )


def _chart_plot_properties(chart: Any, theme_colors: "dict | None" = None) -> PlotProperties:
    plot = _first_chart_plot(chart)
    if plot is None:
        return PlotProperties()
    # firstSliceAng and holeSize — present on doughnutChart and pieChart elements
    first_slice_ang = None
    hole_size = None
    vary_colors = None
    try:
        plot_el = plot._element
        fsa_el = _oxml_find_c(plot_el, "firstSliceAng")
        if fsa_el is not None:
            first_slice_ang = int(fsa_el.get("val", 0))
        hs_el = _oxml_find_c(plot_el, "holeSize")
        if hs_el is not None:
            hole_size = int(hs_el.get("val", 50))
        vc_el = _oxml_find_c(plot_el, "varyColors")
        if vc_el is not None:
            vary_colors = vc_el.get("val", "0") not in ("0", "false")
    except Exception:
        pass
    return PlotProperties(
        grouping=enum_name(safe_get(lambda: plot._element.grouping.val)),
        bar_width_ratio=safe_get(lambda: plot.gap_width),
        overlap=safe_get(lambda: plot.overlap),
        is_horizontal=enum_name(safe_get(lambda: plot._element.barDir.val)) == "BAR",
        area_border=AreaBorder(
            has_border=safe_get(lambda: plot._element.spPr.ln) is not None,
            border_width=_chart_line_width(safe_get(lambda: plot.format.line)),
            border_color=_chart_line_color(safe_get(lambda: plot.format.line), theme_colors),
            has_fill=safe_get(lambda: plot.format.fill.type) is not None,
            fill_color=_chart_format_fill_color(safe_get(lambda: plot.format), theme_colors),
            no_line=safe_get(lambda: plot._element.spPr.ln.noFill) is not None,
        ),
        first_slice_ang=first_slice_ang,
        hole_size=hole_size,
        vary_colors=vary_colors,
    )


def _chart_plot_area_layout(chart: Any) -> "tuple[float|None,float|None,float|None,float|None,str|None,str|None,str|None]":
    """Returns (x, y, w, h, xMode, yMode, layoutTarget) from plotArea/layout/manualLayout."""
    try:
        chart_el = chart._element
        plot_layout = chart_el.find(f".//{_qc('plotArea')}/{_qc('layout')}/{_qc('manualLayout')}")
        if plot_layout is None:
            return (None, None, None, None, None, None, None)
        def _v(tag: str):
            e = _oxml_find_c(plot_layout, tag)
            return float(e.get("val")) if e is not None and e.get("val") else None
        def _s(tag: str):
            e = _oxml_find_c(plot_layout, tag)
            return e.get("val") if e is not None else None
        return (_v("x"), _v("y"), _v("w"), _v("h"), _s("xMode"), _s("yMode"), _s("layoutTarget"))
    except Exception:
        return (None, None, None, None, None, None, None)


def _chart_legend(chart: Any, theme_colors: "dict | None" = None) -> Legend:
    legend = safe_get(lambda: chart.legend)
    if legend is None:
        return Legend(visible=False)
    # Extract manual layout if present
    ml_x = ml_y = ml_w = ml_h = ml_xm = ml_ym = None
    try:
        leg_el = legend._element
        layout_el = _oxml_find_c(leg_el, "layout")
        if layout_el is not None:
            ml_el = _oxml_find_c(layout_el, "manualLayout")
            if ml_el is not None:
                def _ml_val(tag: str):
                    e = _oxml_find_c(ml_el, tag)
                    return float(e.get("val")) if e is not None and e.get("val") else None
                def _ml_mode(tag: str):
                    e = _oxml_find_c(ml_el, tag)
                    return e.get("val") if e is not None else None
                ml_x = _ml_val("x")
                ml_y = _ml_val("y")
                ml_w = _ml_val("w")
                ml_h = _ml_val("h")
                ml_xm = _ml_mode("xMode")
                ml_ym = _ml_mode("yMode")
    except Exception:
        pass
    return Legend(
        visible=True,
        position=enum_name(safe_get(lambda: legend.position)),
        overlay=safe_get(lambda: legend.include_in_layout, True),
        font_name=safe_get(lambda: legend.font.name),
        font_size=_points(safe_get(lambda: legend.font.size)),
        font_bold=safe_get(lambda: legend.font.bold),
        font_color=_font_color(safe_get(lambda: legend.font), theme_colors),
        fill_type=_chart_fill_type(safe_get(lambda: legend._element.spPr)),
        fill_color=_chart_fill_color_from_parent(safe_get(lambda: legend._element.spPr), theme_colors),
        border_width=_chart_line_width(safe_get(lambda: legend.format.line)),
        manual_layout_x=ml_x,
        manual_layout_y=ml_y,
        manual_layout_w=ml_w,
        manual_layout_h=ml_h,
        manual_layout_x_mode=ml_xm,
        manual_layout_y_mode=ml_ym,
    )


def _chart_space_fill(chart: Any, theme_colors: "dict | None" = None) -> dict[str, Any]:
    sp_pr = safe_get(lambda: chart._chartSpace.spPr)
    return {
        "fill_type": _chart_fill_type(sp_pr),
        "fill_color": _chart_fill_color_from_parent(sp_pr, theme_colors),
        "style": safe_get(lambda: chart.chart_style),
    }


def _chart_space_txpr(chart: Any, theme_colors: "dict | None" = None) -> "tuple[str|None, float|None, bool|None, ColorSpec|None]":
    """Extract chart-space c:txPr default text properties (font name, size, bold, color)."""
    font_name = None
    font_size = None
    font_bold = None
    font_color = None
    try:
        txpr_el = _oxml_find_c(chart._chartSpace, "txPr")
        if txpr_el is None:
            return (None, None, None, None)
        defrpr = _oxml_find_descendant_a(txpr_el, "defRPr")
        if defrpr is not None:
            sz = defrpr.get("sz")
            if sz is not None:
                font_size = round(int(sz) / 100.0, 2)
            b = defrpr.get("b")
            if b is not None:
                font_bold = b not in ("0", "false")
            latin = _oxml_find_latin(defrpr)
            if latin is not None:
                font_name = latin.get("typeface")
            font_color = _extract_color_spec(defrpr, theme_colors) if theme_colors else _extract_color_spec(defrpr)
    except Exception:
        pass
    return (font_name, font_size, font_bold, font_color)


def _chart_disp_blanks_as(chart: Any) -> "str | None":
    """Extract c:chart/c:dispBlanksAs value."""
    try:
        el = chart._chartSpace.find(f"{_qc('chart')}/{_qc('dispBlanksAs')}")
        if el is not None:
            return el.get("val")
    except Exception:
        pass
    return None


def _chart_overlay_files(chart: Any) -> "OverlayFiles":
    """Extract chartStyle.xml, chartColors.xml, userShapes, and themeOverride blobs from chart rels."""
    chart_style_blob: bytes | None = None
    chart_colors_blob: bytes | None = None
    user_shapes_blob: bytes | None = None
    theme_override_blob: bytes | None = None
    try:
        for rel in chart.part.rels.values():
            reltype = safe_get(lambda r=rel: r.reltype, "") or ""
            part = safe_get(lambda r=rel: r.target_part)
            if part is None:
                continue
            blob = safe_get(lambda p=part: p.blob)
            if not blob:
                continue
            tail = reltype.split("/")[-1].lower()
            if tail == "chartstyle" or "chartstyle" in tail:
                chart_style_blob = blob
            elif tail in ("chartcolorstyle", "colors") or "color" in tail:
                chart_colors_blob = blob
            elif "usershapes" in tail or "drawing" in tail:
                user_shapes_blob = blob
            elif "theme" in tail and "override" in tail:
                theme_override_blob = blob
    except Exception:
        pass
    from percy.bridge import OverlayFiles
    return OverlayFiles(
        chart_style=chart_style_blob,
        chart_colors=chart_colors_blob,
        chart_user_shapes=user_shapes_blob,
        theme_override=theme_override_blob,
    )


def _chart_data_source(chart: Any) -> ChartDataSource:
    external_data = safe_get(lambda: chart._chartSpace.externalData)
    relationship_id = safe_get(lambda: chart._chartSpace.xlsx_part_rId)
    relationship = safe_get(lambda: chart.part.rels.get(relationship_id)) if relationship_id else None
    embedded_bytes = None
    embedded_filename = None
    sheet_names: list[str] = []
    sheet_dimensions: dict[str, str] = {}
    workbook_sheets: list[ChartWorkbookSheet] = []
    if relationship is not None and not safe_get(lambda: relationship.is_external, False):
        embedded_part = safe_get(lambda: relationship.target_part)
        embedded_bytes = safe_get(lambda: embedded_part.blob)
        embedded_filename = Path(str(safe_get(lambda: relationship.target_ref, ""))).name or None
        sheet_names, sheet_dimensions, workbook_sheets = _xlsx_snapshot(embedded_bytes)

    source_kind = "cache_only"
    if relationship is not None:
        if safe_get(lambda: relationship.is_external, False):
            source_kind = "external_ole_link"
        elif safe_get(lambda: relationship.reltype, "").endswith("/package"):
            source_kind = "embedded_workbook"
        else:
            source_kind = "related_part"

    formulas = _chart_formulas(chart)
    series = safe_get(lambda: list(chart.series), []) or []
    plot = _first_chart_plot(chart)
    categories = safe_get(lambda: list(plot.categories), []) if plot is not None else []
    return ChartDataSource(
        has_external_data=external_data is not None,
        relationship_id=relationship_id,
        relationship_type=safe_get(lambda: relationship.reltype),
        target=safe_get(lambda: relationship.target_ref),
        target_mode="External" if safe_get(lambda: relationship.is_external, False) else "Internal" if relationship else None,
        source_kind=source_kind,
        auto_update=safe_get(lambda: external_data.autoUpdate.val) if external_data is not None else None,
        has_embedded_workbook=embedded_bytes is not None,
        embedded_workbook_filename=embedded_filename,
        embedded_workbook_bytes=embedded_bytes,
        workbook_sheet_names=sheet_names,
        workbook_dimensions=sheet_dimensions,
        workbook_sheets=workbook_sheets,
        cache_series_count=len(series),
        cache_category_count=len(categories),
        cache_point_count=sum(len(safe_get(lambda series=series: series.values, []) or []) for series in series),
        formulas=formulas,
    )


def _table_properties(table: Any, shape: Any) -> TableProperties:
    return TableProperties(
        first_row_header=safe_get(lambda: table.first_row, False),
        first_col_header=safe_get(lambda: table.first_col, False),
        last_row_total=safe_get(lambda: table.last_row, False),
        last_col_total=safe_get(lambda: table.last_col, False),
        banded_rows=safe_get(lambda: table.horz_banding, False),
        banded_cols=safe_get(lambda: table.vert_banding, False),
        style=_parse_table_style_xml(shape),
        conditional_formatting=_table_style_flags(table),
    )


def _table_defaults(table: Any) -> TableDefaults:
    first_cell = safe_get(lambda: table.cell(0, 0))
    first_run = safe_get(lambda: first_cell.text_frame.paragraphs[0].runs[0])
    return TableDefaults(
        text_autofit=enum_name(safe_get(lambda: first_cell.text_frame.auto_size)),
        default_font_name=safe_get(lambda: first_run.font.name),
        default_font_size=_points(safe_get(lambda: first_run.font.size)),
    )


def _table_cell_formats(table: Any, shape: Any, ctx: _OnboardContext) -> list[list[CellFormat]]:
    n_rows = len(table.rows)
    first_row_flag = safe_get(lambda: table.first_row, False)
    last_row_flag = safe_get(lambda: table.last_row, False)
    first_col_flag = safe_get(lambda: table.first_col, False)
    last_col_flag = safe_get(lambda: table.last_col, False)
    banded_rows = safe_get(lambda: table.horz_banding, False)
    banded_cols = safe_get(lambda: table.vert_banding, False)
    table_style = _parse_table_style_xml(shape)

    result = []
    for row_index, row in enumerate(table.rows):
        row_cells = list(row.cells)
        n_cols = len(row_cells)
        row_formats = []
        for col_index, cell in enumerate(row_cells):
            style_section = (
                table_style.resolve_cell(
                    row_index, col_index, n_rows, n_cols,
                    first_row_flag, last_row_flag, first_col_flag, last_col_flag,
                    banded_rows, banded_cols,
                )
                if table_style is not None else None
            )
            row_formats.append(_table_cell_format(cell, row_index, col_index, ctx, style_section, n_rows, n_cols,
                                                   first_row_flag, last_row_flag, first_col_flag, last_col_flag))
        result.append(row_formats)
    return result


def _table_cell_format(
    cell: Any, row_index: int, col_index: int, ctx: _OnboardContext,
    style_section: "TableStyleSection | None" = None,
    n_rows: int = 1, n_cols: int = 1,
    first_row_flag: bool = False, last_row_flag: bool = False,
    first_col_flag: bool = False, last_col_flag: bool = False,
) -> CellFormat:
    paragraphs = _text_paragraphs_from_frame(safe_get(lambda: cell.text_frame), ctx)
    # Apply style section defaults to runs where properties aren't explicitly set
    if style_section is not None and paragraphs:
        for para in paragraphs:
            for run in para.runs:
                if run.font_bold is None and style_section.bold is not None:
                    run.font_bold = style_section.bold
                if run.font_color is None and style_section.font_color is not None:
                    run.font_color = style_section.font_color
                if run.font_name is None and style_section.font_name is not None:
                    run.font_name = style_section.font_name
    explicit_fill_color = _cell_fill_color(cell, ctx.theme_colors)
    explicit_fill_type = _cell_fill_type(cell)
    # Fall back to style section fill when cell has no explicit fill
    if explicit_fill_type is None and explicit_fill_color is None and style_section is not None and style_section.fill_color is not None:
        effective_fill_color = style_section.fill_color
        effective_fill_type = "solidFill"
    else:
        effective_fill_color = explicit_fill_color
        effective_fill_type = explicit_fill_type
    # Resolve borders: explicit cell borders take priority; fall back to style section borders
    explicit_borders = _cell_borders(cell, ctx.theme_colors)
    if style_section is not None:
        # Determine which style border applies to each side (outer vs inside)
        is_left_outer = (col_index == 0)
        is_right_outer = (col_index == n_cols - 1)
        is_top_outer = (row_index == 0)
        is_bottom_outer = (row_index == n_rows - 1)
        # Separator edges: firstRow/lastRow/firstCol/lastCol sections define their own
        # border_bottom/border_top/border_right/border_left as separators between the
        # special section and the body. Use border_bottom (not insideH) for the bottom
        # edge of a header row, border_top for the top edge of a footer row, etc.
        is_header_sep_bottom = (first_row_flag and row_index == 0 and n_rows > 1)
        is_footer_sep_top = (last_row_flag and row_index == n_rows - 1 and n_rows > 1)
        is_first_col_sep_right = (first_col_flag and col_index == 0 and n_cols > 1)
        is_last_col_sep_left = (last_col_flag and col_index == n_cols - 1 and n_cols > 1)
        style_left = style_section.border_left if (is_left_outer or is_last_col_sep_left) else style_section.border_inside_v
        style_right = style_section.border_right if (is_right_outer or is_first_col_sep_right) else style_section.border_inside_v
        style_top = style_section.border_top if (is_top_outer or is_footer_sep_top) else style_section.border_inside_h
        style_bottom = style_section.border_bottom if (is_bottom_outer or is_header_sep_bottom) else style_section.border_inside_h
        # Apply style borders only where no explicit border is set
        from percy.bridge.elements import CellBorders as _CB
        resolved_borders = _CB(
            border_left=explicit_borders.border_left if explicit_borders.border_left is not None else style_left,
            border_right=explicit_borders.border_right if explicit_borders.border_right is not None else style_right,
            border_top=explicit_borders.border_top if explicit_borders.border_top is not None else style_top,
            border_bottom=explicit_borders.border_bottom if explicit_borders.border_bottom is not None else style_bottom,
            diagonal_down=explicit_borders.diagonal_down,
            diagonal_up=explicit_borders.diagonal_up,
        )
    else:
        resolved_borders = explicit_borders
    first_run = safe_get(lambda: cell.text_frame.paragraphs[0].runs[0])
    return CellFormat(
        text=safe_get(lambda: cell.text),
        paragraphs=paragraphs,
        font=CellFont(
            font_name=safe_get(lambda: first_run.font.name),
            font_size=_points(safe_get(lambda: first_run.font.size)),
            font_bold=safe_get(lambda: first_run.font.bold),
            font_italic=safe_get(lambda: first_run.font.italic),
            text_color=_font_color(safe_get(lambda: first_run.font), ctx.theme_colors),
        ),
        alignment=CellAlignment(
            text_alignment=enum_name(safe_get(lambda: cell.text_frame.paragraphs[0].alignment)) or "left",
            vertical_alignment=enum_name(safe_get(lambda: cell.vertical_anchor)) or "top",
        ),
        fill_color=effective_fill_color,
        fill_type=effective_fill_type,
        fill_transparency=safe_get(lambda: cell.fill.transparency),
        borders=resolved_borders,
        margins=Margins(
            margin_left=emu_to_inches(safe_get(lambda: cell.margin_left)),
            margin_right=emu_to_inches(safe_get(lambda: cell.margin_right)),
            margin_top=emu_to_inches(safe_get(lambda: cell.margin_top)),
            margin_bottom=emu_to_inches(safe_get(lambda: cell.margin_bottom)),
        ),
        merge=CellMerge(
            is_merged=safe_get(lambda: cell.is_merge_origin, False) or safe_get(lambda: cell.is_spanned, False),
            is_merge_origin=safe_get(lambda: cell.is_merge_origin, False),
            is_spanned=safe_get(lambda: cell.is_spanned, False),
            merge_span_rows=safe_get(lambda: cell.span_height, 1) if safe_get(lambda: cell.is_merge_origin, False) else 1,
            merge_span_cols=safe_get(lambda: cell.span_width, 1) if safe_get(lambda: cell.is_merge_origin, False) else 1,
        ),
        grid_row=row_index,
        grid_col=col_index,
        text_direction=safe_get(lambda: cell.text_frame._txBody.bodyPr.vert),
        word_wrap=safe_get(lambda: cell.text_frame.word_wrap),
        anchor=safe_get(lambda: cell._tc.tcPr.anchor) if safe_get(lambda: cell._tc.tcPr) is not None else None,
        raw_properties=_cell_raw_properties(cell),
    )


def _run_font_caps(run: Any) -> str | None:
    """Extract the OOXML 'cap' attribute from a run's rPr, walking up to defRPr if absent."""
    try:
        rpr = _oxml_find_rPr(run._r)
        if rpr is not None:
            cap = rpr.get("cap")
            if cap:
                return cap
    except Exception:
        pass
    return None


def _resolve_run_bold(run: Any, paragraph: Any) -> bool | None:
    """Resolve bold for a run by walking run → paragraph defRPr → raw XML attrib."""
    # 1. Explicit run-level bold
    val = safe_get(lambda: run.font.bold)
    if val is not None:
        return val
    # 2. Paragraph-level default run properties
    val = safe_get(lambda: paragraph.font.bold)
    if val is not None:
        return val
    # 3. Raw XML rPr 'b' attribute (handles cases python-pptx doesn't walk)
    try:
        rpr = _oxml_find_rPr(run._r)
        if rpr is not None:
            b_attr = rpr.get("b")
            if b_attr is not None:
                return b_attr not in ("0", "false")
    except Exception:
        pass
    # 4. Font name implies bold weight (e.g. "GothamBold", "DIN Next LT Pro Bold")
    font_name = safe_get(lambda: run.font.name) or ""
    fn_lower = font_name.lower()
    if any(w in fn_lower for w in ("bold", "black", "heavy", "extrabold", "demibold")):
        return True
    return None


from percy.oxml import (
    A_NS as _A_NS,
    C_NS as _C_NS,
    R_NS as _R_NS_CONST,
    qa as _qa,
    qc as _qc,
    qp as _qp,
    find_a as _oxml_find_a,
    find_c as _oxml_find_c,
    find_p as _oxml_find_p,
    find_descendant_a as _oxml_find_descendant_a,
    findall_a as _oxml_findall_a,
    find_pPr as _oxml_find_pPr,
    find_rPr as _oxml_find_rPr,
    find_lstStyle as _oxml_find_lstStyle,
    find_lvl_pPr as _oxml_find_lvl_pPr,
    find_bodyPr as _oxml_find_bodyPr,
    find_buChar as _oxml_find_buChar,
    find_buFont as _oxml_find_buFont,
    find_buBlip as _oxml_find_buBlip,
    bullet_type_from_pPr as _oxml_bullet_type,
    bullet_char_from_pPr as _oxml_bullet_char,
    find_solidFill as _oxml_find_solidFill,
    find_srgbClr as _oxml_find_srgbClr,
    find_sysClr as _oxml_find_sysClr,
    find_schemeClr as _oxml_find_schemeClr,
    find_noFill as _oxml_find_noFill,
    find_grpFill as _oxml_find_grpFill,
    find_gradFill as _oxml_find_gradFill,
    find_pattFill as _oxml_find_pattFill,
    find_blipFill as _oxml_find_blipFill,
    find_gsLst as _oxml_find_gsLst,
    find_gsLst_descendant as _oxml_find_gsLst_descendant,
    findall_gs as _oxml_findall_gs,
    find_lin as _oxml_find_lin,
    find_bgClr as _oxml_find_bgClr,
    find_ln as _oxml_find_ln,
    find_prstGeom as _oxml_find_prstGeom,
    find_avLst as _oxml_find_avLst,
    findall_gd as _oxml_findall_gd,
    find_effectLst as _oxml_find_effectLst,
    find_outerShdw as _oxml_find_outerShdw,
    find_fontScheme as _oxml_find_fontScheme,
    find_clrScheme as _oxml_find_clrScheme,
    find_majorFont as _oxml_find_majorFont,
    find_minorFont as _oxml_find_minorFont,
    find_latin as _oxml_find_latin,
    find_p_spPr as _oxml_find_p_spPr,
)


def _run_rPr(run: Any) -> Any:
    try:
        return _oxml_find_rPr(run._r)
    except Exception:
        return None


def _rPr_attr_int(rPr: Any, attr: str) -> int | None:
    if rPr is None:
        return None
    val = rPr.get(attr)
    if val is None:
        return None
    try:
        return int(val)
    except (ValueError, TypeError):
        return None


def _ooxml_baseline_to_fraction(raw: int | None) -> float | None:
    """Convert OOXML baseline attr (thousandths of %) to bridge fraction.

    OOXML: positive=up (superscript), e.g. 30000 = 30% up.
    Bridge convention: negative=up, positive=down, unit = fraction of font-size.
    So: bridge = -(raw / 100000).
    """
    if raw is None:
        return None
    return -(raw / 100000.0)


def _rPr_attr_str(rPr: Any, attr: str) -> str | None:
    if rPr is None:
        return None
    return rPr.get(attr) or None


def _run_hyperlink(run: Any) -> str | None:
    try:
        rPr = run._r.find(f"{{{_A_NS}}}rPr")
        if rPr is None:
            return None
        hlinkClick = rPr.find(f"{{{_A_NS}}}hlinkClick")
        if hlinkClick is None:
            # also check r namespace
            r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
            return hlinkClick.get(f"{{{r_ns}}}id") if hlinkClick is not None else None
        r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        rel_id = hlinkClick.get(f"{{{r_ns}}}id")
        if rel_id and hasattr(run, "_r"):
            part = safe_get(lambda: run._r.getparent().getparent().getparent().getparent())
            if part is not None and hasattr(part, "part"):
                url = safe_get(lambda: part.part.target_ref(rel_id))
                return url
        return rel_id
    except Exception:
        return None


def _paragraph_pPr(paragraph: Any) -> Any:
    try:
        return _oxml_find_pPr(paragraph._p)
    except Exception:
        return None


def _para_line_spacing(paragraph: Any) -> float | None:
    try:
        ls = paragraph.line_spacing
        if ls is None:
            return None
        # python-pptx returns Pt value for fixed, or a float multiplier for proportional
        from pptx.util import Pt as _Pt
        if isinstance(ls, (int, float)):
            return float(ls)  # proportional multiplier, store as-is
        return float(ls.pt)  # fixed Pt value
    except Exception:
        return None


def _para_space_pt(paragraph: Any, attr: str) -> float | None:
    try:
        val = getattr(paragraph, attr, None)
        if val is None:
            return None
        return float(val.pt)
    except Exception:
        return None


def _pPr_emu_to_inches(pPr: Any, attr: str) -> float | None:
    if pPr is None:
        return None
    val = pPr.get(attr)
    if val is None:
        return None
    try:
        return int(val) / 914400.0  # EMUs to inches
    except (ValueError, TypeError):
        return None


def _lstStyle_pPr_for_level(pPr: Any, level: int) -> Any:
    """Return the txBody's <a:lstStyle>/<a:lvlNpPr> matching `level` (0-indexed).
    PowerPoint stores per-list bullet defaults there when individual paragraphs
    omit a <a:buChar> — e.g. reside slide 4 puts "•" on the lstStyle only.
    """
    if pPr is None:
        return None
    p_el = pPr.getparent()
    if p_el is None: return None
    txBody = p_el.getparent()
    if txBody is None: return None
    return _oxml_find_lvl_pPr(_oxml_find_lstStyle(txBody), level)


def _para_bullet_type(pPr: Any, level: int = 0) -> str:
    bt = _oxml_bullet_type(pPr)
    if bt is not None:
        return bt
    bt = _oxml_bullet_type(_lstStyle_pPr_for_level(pPr, level))
    return bt if bt is not None else "none"


def _para_bullet_char(pPr: Any, level: int = 0) -> str | None:
    val = _oxml_bullet_char(pPr)
    if val is not None:
        return val
    return _oxml_bullet_char(_lstStyle_pPr_for_level(pPr, level))


def _para_bullet_font(pPr: Any) -> str | None:
    if pPr is None:
        return None
    el = _oxml_find_buFont(pPr)
    if el is not None:
        return el.get("typeface")
    return None


def _para_bullet_blip(pPr: Any, part: Any) -> tuple[bytes | None, str | None]:
    """Extract image bytes and extension from a buBlip picture bullet."""
    if pPr is None:
        return None, None
    blip_el = _oxml_find_buBlip(pPr)
    if blip_el is None:
        return None, None
    inner = blip_el.find(_qa("blip"))
    if inner is None:
        return None, None
    rId = inner.get(f"{{{_R_NS_CONST}}}embed")
    if not rId:
        return None, None
    try:
        image_part = part.related_part(rId)
        blob = image_part.blob
        ct = image_part.content_type or ""
        ext = "png" if "png" in ct else "jpeg" if ("jpeg" in ct or "jpg" in ct) else "png"
        return blob, ext
    except Exception:
        return None, None


def _runs_with_breaks(paragraph: Any, make_run_fn: Any) -> list[TextRun]:
    """Iterate paragraph XML children to preserve <a:br> line breaks."""
    result: list[TextRun] = []
    p_el = safe_get(lambda: paragraph._p)
    if p_el is None:
        return [make_run_fn(r) for r in paragraph.runs]
    run_iter = iter(paragraph.runs)
    for child in p_el:
        local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if local == "r":
            try:
                result.append(make_run_fn(next(run_iter)))
            except StopIteration:
                pass
        elif local == "br":
            result.append(TextRun(is_line_break=True))
    return result


def _text_paragraphs_from_frame(text_frame: Any, ctx: _OnboardContext) -> list[TextParagraph]:
    if text_frame is None:
        return []
    paragraphs = []
    for paragraph in safe_get(lambda: text_frame.paragraphs, []) or []:
        pPr = _paragraph_pPr(paragraph)

        def _make_run(run: Any) -> TextRun:
            rPr = _run_rPr(run)
            return TextRun(
                text=safe_get(lambda r=run: r.text, "") or "",
                font_name=ctx.resolve_font_name(safe_get(lambda r=run: r.font.name)),
                font_size=_points(safe_get(lambda r=run: r.font.size)),
                font_bold=_resolve_run_bold(run, paragraph),
                font_italic=safe_get(lambda r=run: r.font.italic),
                font_underline=safe_get(lambda r=run: r.font.underline),
                font_color=_font_color(safe_get(lambda r=run: r.font), ctx.theme_colors),
                font_caps=_run_font_caps(run),
                char_spacing=_rPr_attr_int(rPr, "spc"),
                baseline_shift=_ooxml_baseline_to_fraction(_rPr_attr_int(rPr, "baseline")),
                strikethrough=_rPr_attr_str(rPr, "strike"),
                hyperlink=_run_hyperlink(run),
            )

        runs = _runs_with_breaks(paragraph, _make_run)
        blip_bytes, blip_ext = _para_bullet_blip(pPr, ctx.slide_part) if _para_bullet_type(pPr) == "image" else (None, None)
        paragraphs.append(TextParagraph(
            runs=runs,
            alignment=enum_name(safe_get(lambda p=paragraph: p.alignment)),
            indent_level=safe_get(lambda p=paragraph: p.level, 0) or 0,
            line_spacing=_para_line_spacing(paragraph),
            space_before=_para_space_pt(paragraph, "space_before"),
            space_after=_para_space_pt(paragraph, "space_after"),
            left_indent=_pPr_emu_to_inches(pPr, "marL"),
            first_line_indent=_pPr_emu_to_inches(pPr, "indent"),
            bullet_type=_para_bullet_type(pPr, safe_get(lambda p=paragraph: p.level, 0) or 0),
            bullet_char=_para_bullet_char(pPr, safe_get(lambda p=paragraph: p.level, 0) or 0),
            bullet_font=_para_bullet_font(pPr),
            bullet_blip_bytes=blip_bytes,
            bullet_blip_ext=blip_ext,
            end_para_font_size=_end_para_font_size(paragraph),
        ))
    return paragraphs



def _table_style_flags(table: Any) -> list[dict[str, Any]]:
    flags = {
        "first_row": safe_get(lambda: table.first_row, False),
        "first_col": safe_get(lambda: table.first_col, False),
        "last_row": safe_get(lambda: table.last_row, False),
        "last_col": safe_get(lambda: table.last_col, False),
        "banded_rows": safe_get(lambda: table.horz_banding, False),
        "banded_cols": safe_get(lambda: table.vert_banding, False),
    }
    return [{"kind": "table_style_option", "name": name, "enabled": value} for name, value in flags.items() if value]


def _cell_fill_type(cell: Any) -> str | None:
    tc_pr = safe_get(lambda: cell._tc.tcPr)
    if tc_pr is None:
        return None
    for child in tc_pr:
        local_name = _local_name(safe_get(lambda child=child: child.tag, ""))
        if local_name.endswith("Fill"):
            return local_name
    return None


def _cell_fill_color(cell: Any, theme_colors: "dict | None" = None) -> "ColorSpec | None":
    tc_pr = safe_get(lambda: cell._tc.tcPr)
    return _extract_color_spec(tc_pr, theme_colors)


def _cell_borders(cell: Any, theme_colors: "dict | None" = None) -> CellBorders:
    tc_pr = safe_get(lambda: cell._tc.tcPr)
    return CellBorders(
        border_left=_table_border(_first_child(tc_pr, "lnL"), theme_colors),
        border_right=_table_border(_first_child(tc_pr, "lnR"), theme_colors),
        border_top=_table_border(_first_child(tc_pr, "lnT"), theme_colors),
        border_bottom=_table_border(_first_child(tc_pr, "lnB"), theme_colors),
        diagonal_down=_table_border(_first_child(tc_pr, "lnTlToBr"), theme_colors),
        diagonal_up=_table_border(_first_child(tc_pr, "lnBlToTr"), theme_colors),
    )


def _table_border(line: Any, theme_colors: "dict | None" = None) -> Border | None:
    if line is None:
        return None
    dash = _first_child(line, "prstDash")
    return Border(
        style=safe_get(lambda: dash.get("val")),
        width=_emu_to_points(safe_get(lambda: line.get("w"))),
        color=_line_xml_color(line, theme_colors),
        visible=_first_child(line, "noFill") is None,
        dash_style=safe_get(lambda: dash.get("val")),
    )


def _line_xml_color(line: Any, theme_colors: "dict | None" = None) -> "ColorSpec | None":
    if line is None:
        return None
    return _extract_color_spec(line, theme_colors)


def _extract_color_spec(element: Any, theme_colors: "dict | None" = None) -> "ColorSpec | None":
    """Extract a ColorSpec from any OOXML element containing color info."""
    from percy.bridge.elements import ColorSpec
    from percy.oxml import find_solidFill, find_srgbClr, find_sysClr, find_schemeClr
    if element is None:
        return None
    tag = element.tag.split("}")[-1] if hasattr(element, "tag") else ""
    if tag not in ("solidFill", "schemeClr", "srgbClr", "sysClr"):
        solid = find_solidFill(element)
        if solid is None:
            return None
        element = solid
        tag = "solidFill"
    def _read_alpha(clr_el: Any) -> "int | None":
        alpha_el = clr_el.find(_qa("alpha"))
        if alpha_el is not None:
            try:
                return int(alpha_el.get("val", ""))
            except (ValueError, TypeError):
                pass
        return None

    if tag == "solidFill":
        srgb = find_srgbClr(element)
        if srgb is not None:
            val = srgb.get("val", "")
            return ColorSpec(value=f"#{val.upper()}", alpha=_read_alpha(srgb)) if val else None
        sys_el = find_sysClr(element)
        if sys_el is not None:
            last_clr = sys_el.get("lastClr", "")
            return ColorSpec(value=f"#{last_clr.upper()}", alpha=_read_alpha(sys_el)) if last_clr else None
        scheme = find_schemeClr(element)
        if scheme is not None:
            return _extract_scheme_color_spec(scheme, theme_colors)
        return None
    if tag == "srgbClr":
        val = element.get("val", "")
        return ColorSpec(value=f"#{val.upper()}", alpha=_read_alpha(element)) if val else None
    if tag == "sysClr":
        last_clr = element.get("lastClr", "")
        return ColorSpec(value=f"#{last_clr.upper()}", alpha=_read_alpha(element)) if last_clr else None
    if tag == "schemeClr":
        return _extract_scheme_color_spec(element, theme_colors)
    return None


def _extract_scheme_color_spec(scheme_el: Any, theme_colors: "dict | None" = None) -> "ColorSpec | None":
    """Extract ColorSpec from a <a:schemeClr> element.

    When theme_colors is provided, the scheme base color is resolved immediately to a
    concrete hex value. Modifiers (lum_mod, tint, shade, etc.) are preserved so
    ColorSpec.resolve() can apply them. This ensures bridge objects never contain
    ``scheme:X`` references — all colors are self-contained hex values.
    """
    from percy.bridge.elements import ColorSpec
    xml_name = scheme_el.get("val", "")
    if not xml_name:
        return None
    normalized = _XML_SCHEME_NAME.get(xml_name, xml_name.upper())

    def _int_val(tag: str) -> int | None:
        el = _oxml_find_a(scheme_el, tag)
        if el is not None:
            try:
                return int(el.get("val", ""))
            except (ValueError, TypeError):
                pass
        return None

    lum_mod = _int_val("lumMod")
    lum_off = _int_val("lumOff")
    shade = _int_val("shade")
    tint = _int_val("tint")
    alpha = _int_val("alpha")
    hue_mod = _int_val("hueMod")
    sat_mod = _int_val("satMod")

    # Fully resolve to concrete hex at extraction time when theme_colors available.
    # Apply ALL modifiers (lum_mod, tint, shade, etc.) so bridge objects are self-contained
    # and never contain scheme: references.
    if theme_colors:
        raw = theme_colors.get(normalized, "")
        if raw:
            base_hex = raw if raw.startswith("#") else f"#{raw}"
            if len(base_hex.lstrip("#")) == 6:
                # Build a temporary spec with the resolved base color + modifiers to get final hex
                temp = ColorSpec(
                    value=base_hex,
                    lum_mod=lum_mod, lum_off=lum_off,
                    shade=shade, tint=tint,
                    hue_mod=hue_mod, sat_mod=sat_mod,
                )
                resolved_hex = temp.resolve().lstrip("#")
                if len(resolved_hex) == 6 and resolved_hex.upper() != "888888":
                    return ColorSpec(value=f"#{resolved_hex.upper()}", alpha=alpha)

    return ColorSpec(
        value=f"scheme:{normalized}",
        lum_mod=lum_mod, lum_off=lum_off, shade=shade, tint=tint,
        alpha=alpha, hue_mod=hue_mod, sat_mod=sat_mod,
    )


def _emu_to_points(value: Any) -> float | None:
    if value is None:
        return None
    return round(float(value) / 12700.0, 2)


def _cell_raw_properties(cell: Any) -> dict[str, Any]:
    tc = safe_get(lambda: cell._tc)
    tc_pr = safe_get(lambda: tc.tcPr)
    return {
        "grid_span": safe_get(lambda: tc.gridSpan),
        "row_span": safe_get(lambda: tc.rowSpan),
        "h_merge": safe_get(lambda: tc.hMerge),
        "v_merge": safe_get(lambda: tc.vMerge),
        "tcPr_attrs": dict(safe_get(lambda: tc_pr.attrib, {}) or {}),
    }


def _cell_typed_value(cell: Any) -> Any:
    """Return cell text coerced to int/float where unambiguous, otherwise str."""
    text = safe_get(lambda: cell.text, "") or ""
    stripped = text.strip().replace(",", "")
    if stripped:
        try:
            as_float = float(stripped)
            return int(as_float) if as_float.is_integer() else as_float
        except ValueError:
            pass
    return text


def _table_unsupported_features(table: Any) -> list[str]:
    debt = []
    if any(safe_get(lambda cell=cell: cell.is_merge_origin, False) for cell in table.iter_cells()):
        debt.append("table_merged_cells")
    if any(_has_table_diagonal_border(cell) for cell in table.iter_cells()):
        debt.append("table_diagonal_borders")
    return debt


# ---------------------------------------------------------------------------
# Table style resolution — parse tableStyles.xml into TableStyle objects
# ---------------------------------------------------------------------------

_TBL_SECTION_MAP = {
    "wholeTbl": "whole_tbl",
    "band1H": "band1_h", "band2H": "band2_h",
    "band1V": "band1_v", "band2V": "band2_v",
    "firstRow": "first_row", "lastRow": "last_row",
    "firstCol": "first_col", "lastCol": "last_col",
    "nwCell": "nw_cell", "neCell": "ne_cell",
    "swCell": "sw_cell", "seCell": "se_cell",
}


def _parse_table_style_xml(shape: Any) -> "TableStyle | None":
    """Read tableStyles.xml from the PPTX package and parse the style matching the table.

    Returns a fully-structured TableStyle object with per-section formatting rules,
    or None if no style is found or parsing fails.
    """
    try:
        from lxml import etree as _etree2
        _A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        table = shape.table
        tbl_pr = table._tbl.tblPr
        if tbl_pr is None:
            return None
        style_id_el = tbl_pr.find(f"{{{_A}}}tableStyleId")
        if style_id_el is None or not style_id_el.text:
            return None
        style_id = style_id_el.text.strip()

        # Find tableStyles.xml in the package
        pkg = shape.part.package
        tbl_styles_part = None
        for part in pkg.iter_parts():
            if hasattr(part, 'partname') and 'tableStyles' in str(part.partname):
                tbl_styles_part = part
                break
        if tbl_styles_part is None:
            return None

        root = _etree2.fromstring(tbl_styles_part.blob)
        tbl_style_el = None
        for style_el in root.findall(f"{{{_A}}}tblStyle"):
            if style_el.get("styleId") == style_id:
                tbl_style_el = style_el
                break
        if tbl_style_el is None:
            return None

        from percy.diagnostics.inheritance import _theme_info
        theme = _theme_info(shape)
        theme_colors = theme.get("colors", {})
        minor_font = theme.get("minor_latin")
        major_font = theme.get("major_latin")

        table_style = TableStyle(
            style_id=style_id,
            style_name=tbl_style_el.get("styleName"),
        )
        for xml_name, attr_name in _TBL_SECTION_MAP.items():
            section_el = tbl_style_el.find(f"{{{_A}}}{xml_name}")
            if section_el is None:
                continue
            section = _parse_tbl_style_section(section_el, theme_colors, minor_font, major_font)
            setattr(table_style, attr_name, section)
        return table_style
    except Exception:
        return None


def _parse_tbl_style_section(section_el: Any, theme_colors: dict, minor_font: str | None, major_font: str | None) -> TableStyleSection:
    _A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    bold: bool | None = None
    font_name: str | None = None
    font_color = None
    fill_color = None

    tx_style = section_el.find(f"{{{_A}}}tcTxStyle")
    if tx_style is not None:
        b_attr = tx_style.get("b")
        if b_attr == "on":
            bold = True
        elif b_attr == "off":
            bold = False
        font_ref = tx_style.find(f"{{{_A}}}fontRef")
        if font_ref is not None:
            idx = font_ref.get("idx")
            if idx == "minor":
                font_name = minor_font
            elif idx == "major":
                font_name = major_font or minor_font
        font_color = _extract_tx_style_color(tx_style, theme_colors)

    border_left = border_right = border_top = border_bottom = border_inside_h = border_inside_v = None
    tc_style = section_el.find(f"{{{_A}}}tcStyle")
    if tc_style is not None:
        fill_color = _extract_tc_style_fill(tc_style, theme_colors)
        tc_bdr = tc_style.find(f"{{{_A}}}tcBdr")
        if tc_bdr is not None:
            def _parse_bdr(tag: str) -> "Border | None":
                from percy.bridge.elements import Border, ColorSpec
                edge = tc_bdr.find(f"{{{_A}}}{tag}")
                if edge is None:
                    return None
                ln = edge.find(f"{{{_A}}}ln")
                if ln is None:
                    return None
                w_emu = int(ln.get("w", "0") or "0")
                width_pt = w_emu / 12700.0 if w_emu else None
                no_fill = ln.find(f"{{{_A}}}noFill") is not None
                if no_fill:
                    return Border(visible=False)
                # Color is inside solidFill → schemeClr or srgbClr
                solid = ln.find(f"{{{_A}}}solidFill")
                clr = None
                if solid is not None:
                    scheme = solid.find(f"{{{_A}}}schemeClr")
                    if scheme is not None:
                        val = scheme.get("val", "")
                        raw = theme_colors.get(val, "")
                        if raw:
                            clr = ColorSpec(value=f"#{raw}" if not raw.startswith("#") else raw)
                    srgb = solid.find(f"{{{_A}}}srgbClr")
                    if srgb is not None and clr is None:
                        v = srgb.get("val", "")
                        if v:
                            clr = ColorSpec(value=f"#{v}")
                return Border(visible=True, width=width_pt, color=clr)
            border_left = _parse_bdr("left")
            border_right = _parse_bdr("right")
            border_top = _parse_bdr("top")
            border_bottom = _parse_bdr("bottom")
            border_inside_h = _parse_bdr("insideH")
            border_inside_v = _parse_bdr("insideV")

    return TableStyleSection(bold=bold, font_name=font_name, font_color=font_color, fill_color=fill_color,
                             border_left=border_left, border_right=border_right,
                             border_top=border_top, border_bottom=border_bottom,
                             border_inside_h=border_inside_h, border_inside_v=border_inside_v)


def _extract_tx_style_color(tx_style: Any, theme_colors: dict) -> "ColorSpec | None":
    """Extract text color from a tcTxStyle element (schemeClr or srgbClr)."""
    try:
        _A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        scheme_clr = tx_style.find(f"{{{_A}}}schemeClr")
        if scheme_clr is not None:
            val = scheme_clr.get("val")
            if val and val in theme_colors:
                from percy.bridge.elements import ColorSpec
                raw = theme_colors[val]
                return ColorSpec(value=f"#{raw}" if not raw.startswith("#") else raw)
        srgb = tx_style.find(f"{{{_A}}}srgbClr")
        if srgb is not None:
            val = srgb.get("val")
            if val:
                from percy.bridge.elements import ColorSpec
                return ColorSpec(value=f"#{val}")
    except Exception:
        pass
    return None


def _extract_tc_style_fill(tc_style: Any, theme_colors: dict) -> "ColorSpec | None":
    """Extract fill color from a tcStyle element."""
    try:
        _A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        fill = tc_style.find(f"{{{_A}}}fill")
        if fill is None:
            return None
        solid = fill.find(f"{{{_A}}}solidFill")
        if solid is None:
            return None
        # schemeClr with optional tint/shade
        scheme = solid.find(f"{{{_A}}}schemeClr")
        if scheme is not None:
            val = scheme.get("val")
            raw = theme_colors.get(val, "") if val else ""
            if not raw:
                return None
            hex_color = raw.lstrip("#")
            # Apply tint/shade modifiers
            tint_el = scheme.find(f"{{{_A}}}tint")
            shade_el = scheme.find(f"{{{_A}}}shade")
            if tint_el is not None:
                tint = int(tint_el.get("val", "0")) / 100000.0
                hex_color = _apply_tint(hex_color, tint)
            elif shade_el is not None:
                shade = int(shade_el.get("val", "0")) / 100000.0
                hex_color = _apply_shade(hex_color, shade)
            from percy.bridge.elements import ColorSpec
            return ColorSpec(value=f"#{hex_color}")
        # srgbClr
        srgb = solid.find(f"{{{_A}}}srgbClr")
        if srgb is not None:
            val = srgb.get("val")
            if val:
                from percy.bridge.elements import ColorSpec
                return ColorSpec(value=f"#{val}")
    except Exception:
        pass
    return None


def _apply_tint(hex_color: str, tint: float) -> str:
    """Apply OOXML tint to a hex color (blend toward white, tint=1.0 = full color)."""
    try:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        r = int(r + (255 - r) * (1.0 - tint))
        g = int(g + (255 - g) * (1.0 - tint))
        b = int(b + (255 - b) * (1.0 - tint))
        return f"{r:02X}{g:02X}{b:02X}"
    except Exception:
        return hex_color


def _apply_shade(hex_color: str, shade: float) -> str:
    """Apply OOXML shade to a hex color (blend toward black, shade=1.0 = full color)."""
    try:
        r = int(int(hex_color[0:2], 16) * shade)
        g = int(int(hex_color[2:4], 16) * shade)
        b = int(int(hex_color[4:6], 16) * shade)
        return f"{r:02X}{g:02X}{b:02X}"
    except Exception:
        return hex_color



def _first_chart_plot(chart: Any) -> Any:
    plots = safe_get(lambda: list(chart.plots), None)
    if not plots:
        return None
    return plots[0]


def _has_table_diagonal_border(cell: Any) -> bool:
    tc_pr = safe_get(lambda: cell._tc.tcPr)
    return _first_child(tc_pr, "lnTlToBr") is not None or _first_child(tc_pr, "lnBlToTr") is not None


def _first_child(parent: Any, local_name: str) -> Any:
    if parent is None:
        return None
    for child in parent:
        if _local_name(safe_get(lambda child=child: child.tag, "")) == local_name:
            return child
    return None


def _xlsx_snapshot(blob: bytes | None) -> tuple[list[str], dict[str, str], list[ChartWorkbookSheet]]:
    if not blob:
        return [], {}, []
    try:
        from io import BytesIO
        import xml.etree.ElementTree as ET

        ns = {
            "main": "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
            "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
        }
        with ZipFile(BytesIO(blob)) as workbook_zip:
            workbook = ET.fromstring(workbook_zip.read("xl/workbook.xml"))
            rels = ET.fromstring(workbook_zip.read("xl/_rels/workbook.xml.rels"))
            rid_to_target = {rel.get("Id"): rel.get("Target") for rel in rels}
            shared_strings = _xlsx_shared_strings(workbook_zip, ET, ns)
            names = []
            dimensions = {}
            sheets = []
            for sheet in workbook.findall("main:sheets/main:sheet", ns):
                name = sheet.get("name") or ""
                names.append(name)
                target = rid_to_target.get(sheet.get(f"{{{ns['r']}}}id"))
                if not target:
                    continue
                sheet_path = _xlsx_part_path(target, "xl")
                sheet_root = ET.fromstring(workbook_zip.read(sheet_path))
                dimension = sheet_root.find("main:dimension", ns)
                if dimension is not None and dimension.get("ref"):
                    dimensions[name] = dimension.get("ref") or ""
                sheets.append(
                    ChartWorkbookSheet(
                        name=name,
                        dimension=dimensions.get(name),
                        cells=_xlsx_cells(sheet_root, ns, shared_strings),
                    )
                )
            return names, dimensions, sheets
    except Exception:
        return [], {}, []


def _xlsx_part_path(target: str, base_dir: str) -> str:
    target = target.replace("\\", "/")
    if target.startswith("/"):
        return target.lstrip("/")
    while target.startswith("../"):
        target = target[3:]
    return f"{base_dir}/{target.lstrip('/')}"


def _xlsx_shared_strings(workbook_zip: ZipFile, et_module: Any, ns: dict[str, str]) -> list[str]:
    try:
        shared_strings_root = et_module.fromstring(workbook_zip.read("xl/sharedStrings.xml"))
    except Exception:
        return []
    strings = []
    for item in shared_strings_root.findall("main:si", ns):
        text_nodes = item.findall(".//main:t", ns)
        strings.append("".join(node.text or "" for node in text_nodes))
    return strings


def _xlsx_cells(sheet_root: Any, ns: dict[str, str], shared_strings: list[str]) -> list[ChartWorkbookCell]:
    cells = []
    for cell in sheet_root.findall(".//main:c", ns):
        address = cell.get("r") or ""
        row, column = _cell_address_parts(address)
        data_type = cell.get("t")
        style_id = _safe_int(cell.get("s"))
        formula_node = cell.find("main:f", ns)
        value_node = cell.find("main:v", ns)
        inline_nodes = cell.findall(".//main:is/main:t", ns)
        raw_value = value_node.text if value_node is not None else None
        if inline_nodes:
            value: Any = "".join(node.text or "" for node in inline_nodes)
        else:
            value = _xlsx_cell_value(raw_value, data_type, shared_strings)
        cells.append(
            ChartWorkbookCell(
                address=address,
                row=row,
                column=column,
                value=value,
                formula=formula_node.text if formula_node is not None else None,
                data_type=data_type,
                style_id=style_id,
            )
        )
    return cells


def _xlsx_cell_value(raw_value: str | None, data_type: str | None, shared_strings: list[str]) -> Any:
    if raw_value is None:
        return None
    if data_type == "s":
        index = _safe_int(raw_value)
        return shared_strings[index] if index is not None and 0 <= index < len(shared_strings) else raw_value
    if data_type == "b":
        return raw_value == "1"
    if data_type in {"str", "inlineStr"}:
        return raw_value
    return _number_or_text(raw_value)


def _cell_address_parts(address: str) -> tuple[int, int]:
    letters = "".join(char for char in address if char.isalpha()).upper()
    digits = "".join(char for char in address if char.isdigit())
    column = 0
    for char in letters:
        column = column * 26 + (ord(char) - ord("A") + 1)
    return _safe_int(digits) or 0, column


def _number_or_text(value: str) -> int | float | str:
    try:
        number = float(value)
    except ValueError:
        return value
    return int(number) if number.is_integer() else number


def _safe_int(value: Any) -> int | None:
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def _chart_formulas(chart: Any) -> list[str]:
    formulas = safe_get(lambda: chart._chartSpace.xpath(".//c:f/text()"), []) or []
    return sorted({str(formula) for formula in formulas if formula})


def _chart_unsupported_features(chart: Any) -> list[str]:
    debt = []
    plots = safe_get(lambda: list(chart.plots), None)
    if plots is None:
        debt.append("chart_type_rebuild_not_supported")
        plots = []
    if len(plots or []) > 1:
        debt.append("chart_combo_plots")
    if any(_local_name(safe_get(lambda: plot._element.tag, "")) in {"bubbleChart", "surfaceChart", "bar3DChart", "area3DChart", "line3DChart", "pie3DChart"} for plot in plots or []):
        debt.append("chart_type_rebuild_not_supported")
    data_source = _chart_data_source(chart)
    if data_source.source_kind == "external_ole_link":
        debt.append("chart_external_ole_link")
    elif data_source.source_kind == "embedded_workbook":
        debt.append("chart_embedded_workbook")
    if safe_get(lambda: chart._chartSpace.style) is not None:
        debt.append("chart_style_theme_resolution")
    return debt


def _is_group_shape(shape: Any) -> bool:
    return enum_name(safe_get(lambda: shape.shape_type)) == "GROUP" and safe_get(lambda: shape.shapes, None) is not None


def _group_fill_color(group_shape: Any, theme_colors: "dict | None" = None) -> "ColorSpec | None":
    """Extract resolved fill color from a group shape's grpSpPr element."""
    try:
        _P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        el = group_shape._element
        grpSpPr = el.find(f"{{{_P}}}grpSpPr")
        if grpSpPr is None:
            return None
        return _extract_color_spec(grpSpPr, theme_colors)
    except Exception:
        return None


def _shape_has_grp_fill(shape: Any) -> bool:
    """Return True if the shape's spPr contains <a:grpFill> (inherit fill from parent group)."""
    try:
        spPr = shape.element.spPr
        if spPr is None:
            return False
        return _oxml_find_grpFill(spPr) is not None
    except Exception:
        return False


def _text_frame_text(text_frame: Any) -> str | None:
    if text_frame is None:
        return None
    paragraphs = safe_get(lambda: text_frame.paragraphs, []) or []
    return "\n".join("".join(run.text for run in paragraph.runs) for paragraph in paragraphs)


def _font_color(font: Any, theme_colors: "dict | None" = None) -> "ColorSpec | None":
    from percy.bridge.elements import ColorSpec
    try:
        rpr_el = font._element
        spec = _extract_color_spec(rpr_el, theme_colors)
        if spec is not None:
            return spec
    except Exception:
        pass
    rgb = safe_get(lambda: font.color.rgb)
    if rgb is not None:
        return ColorSpec(value=f"#{str(rgb)}")
    theme = safe_get(lambda: font.color.theme_color)
    if theme is not None:
        name = enum_name(theme)
        normalized = _XML_SCHEME_NAME.get(name, name) if name else ""
        if not normalized:
            return None
        if theme_colors:
            raw = theme_colors.get(normalized, "")
            if raw:
                return ColorSpec(value=raw if raw.startswith("#") else f"#{raw}")
        return ColorSpec(value=f"scheme:{normalized}")
    return None


def _chart_format_fill_color(chart_format: Any, theme_colors: "dict | None" = None) -> "ColorSpec | None":
    return _chart_fill_color(safe_get(lambda: chart_format.fill), theme_colors)


def _chart_fill_color(fill: Any, theme_colors: "dict | None" = None) -> "ColorSpec | None":
    from percy.bridge.elements import ColorSpec
    try:
        sp_pr = getattr(fill, "_xPr", None)
        if sp_pr is None:
            sp_pr = getattr(fill, "_element", None)
        spec = _extract_color_spec(sp_pr, theme_colors)
        if spec is not None:
            return spec
    except Exception:
        pass
    rgb = safe_get(lambda: fill.fore_color.rgb)
    if rgb is not None:
        return ColorSpec(value=f"#{str(rgb)}")
    theme = safe_get(lambda: fill.fore_color.theme_color)
    if theme is not None:
        name = enum_name(theme)
        normalized = _XML_SCHEME_NAME.get(name, name) if name else ""
        if normalized and theme_colors:
            raw = theme_colors.get(normalized, "")
            if raw:
                return ColorSpec(value=raw if raw.startswith("#") else f"#{raw}")
        return ColorSpec(value=f"scheme:{normalized}") if normalized else None
    return None


def _chart_fill_color_from_parent(parent: Any, theme_colors: "dict | None" = None) -> "ColorSpec | None":
    if parent is None:
        return None
    return _extract_color_spec(parent, theme_colors)


def _chart_fill_type(parent: Any) -> str | None:
    if parent is None:
        return None
    for name in ("solidFill", "gradFill", "pattFill", "blipFill", "noFill"):
        if safe_get(lambda name=name: getattr(parent, name)) is not None:
            return name
    return None


def _chart_line_format(line: Any, theme_colors: "dict | None" = None) -> LineFormat:
    return LineFormat(
        line_visible=safe_get(lambda: line.fill.type) is not None,
        line_width=_chart_line_width(line),
        line_style=enum_name(safe_get(lambda: line.dash_style)),
        line_color=_chart_line_color(line, theme_colors),
    )


def _chart_line_color(line: Any, theme_colors: "dict | None" = None) -> "ColorSpec | None":
    from percy.bridge.elements import ColorSpec
    try:
        ln_el = safe_get(lambda: line._element)
        if ln_el is not None:
            spec = _extract_color_spec(ln_el, theme_colors)
            if spec is not None:
                return spec
    except Exception:
        pass
    rgb = safe_get(lambda: line.color.rgb)
    if rgb is not None:
        return ColorSpec(value=f"#{str(rgb)}")
    theme = safe_get(lambda: line.color.theme_color)
    if theme is not None:
        name = enum_name(theme)
        normalized = _XML_SCHEME_NAME.get(name, name) if name else ""
        if normalized and theme_colors:
            raw = theme_colors.get(normalized, "")
            if raw:
                return ColorSpec(value=raw if raw.startswith("#") else f"#{raw}")
        return ColorSpec(value=f"scheme:{normalized}") if normalized else None
    return None


def _chart_line_width(line: Any) -> float | None:
    width = safe_get(lambda: line.width)
    return _points(width)


def _axis_gridlines(axis: Any, theme_colors: "dict | None" = None) -> Gridlines:
    """Build Gridlines bridge object for an axis, reading noFill BEFORE accessing color."""
    has_gl = safe_get(lambda: axis.has_major_gridlines, False)
    # noFill must be read before _chart_line_color, because line.color.rgb destroys noFill XML
    no_fill = _axis_gridline_no_fill(axis)
    line = safe_get(lambda: axis.major_gridlines.format.line)
    return Gridlines(
        has_major_gridlines=has_gl,
        gridline_color=_chart_line_color(line, theme_colors),
        gridline_width=_chart_line_width(line),
        gridline_no_fill=no_fill,
    )


def _axis_gridline_no_fill(axis: Any) -> bool:
    """Return True if the axis major gridlines line has explicit <a:noFill/> in XML."""
    try:
        from lxml import etree
        ax_el = safe_get(lambda: axis._element)
        if ax_el is None:
            return False
        _C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
        _A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        gl_el = ax_el.find(f"{{{_C}}}majorGridlines")
        if gl_el is None:
            return False
        gl_xml = etree.tostring(gl_el).decode()
        return "noFill" in gl_xml
    except Exception:
        return False


def _chart_line_is_no_fill(line: Any) -> bool:
    """Return True if the <a:ln> element explicitly has <a:noFill/>."""
    if line is None:
        return False
    try:
        # LineFormat._ln gives the <a:ln> element (or None if absent)
        ln_el = safe_get(lambda: line._ln)
        if ln_el is not None and _oxml_find_noFill(ln_el) is not None:
            return True
        # Fallback: parent element may contain <a:ln>
        parent_el = safe_get(lambda: line._parent)
        if parent_el is not None:
            ln_el2 = _oxml_find_ln(parent_el)
            if ln_el2 is not None and _oxml_find_noFill(ln_el2) is not None:
                return True
    except Exception:
        pass
    return False


def _series_gradient_stops(series: Any) -> list:
    """Extract gradient stops from a chart series' spPr/gradFill."""
    try:
        sp_pr = _oxml_find_p_spPr(series._element) or _oxml_find_c(series._element, "spPr")
        if sp_pr is None:
            return []
        grad = _oxml_find_gradFill(sp_pr)
        if grad is None:
            return []
        gs_lst = _oxml_find_gsLst(grad)
        if gs_lst is None:
            return []
        stops = []
        for gs in _oxml_findall_gs(gs_lst):
            pos = int(gs.get("pos", "0")) / 100000.0
            cs = _extract_color_spec(gs)
            if cs is not None:
                stops.append(GradientStop(position=pos, color=cs))
        return stops
    except Exception:
        return []


def _marker_format(marker: Any, theme_colors: "dict | None" = None) -> MarkerFormat:
    return MarkerFormat(
        marker_style=enum_name(safe_get(lambda: marker.style)),
        marker_size=safe_get(lambda: marker.size),
        marker_color=_chart_format_fill_color(safe_get(lambda: marker.format), theme_colors),
        marker_line_visible=safe_get(lambda: marker.format.line.fill.type) is not None,
    )


def _dl_flag(d_lbls: Any, tag: str) -> bool:
    """Read a show-flag (showVal, showCatName, etc.) from dLbls XML, default False."""
    try:
        el = _oxml_find_c(d_lbls, tag)
        if el is not None:
            return el.get("val", "0") not in ("0", "false")
    except Exception:
        pass
    return False


def _data_labels(d_lbls: Any, theme_colors: "dict | None" = None) -> DataLabels:
    if d_lbls is None:
        return DataLabels(show=False)

    show_val = _dl_flag(d_lbls, "showVal")
    show_cat = _dl_flag(d_lbls, "showCatName")
    show_ser = _dl_flag(d_lbls, "showSerName")
    show_pct = _dl_flag(d_lbls, "showPercent")
    show_key = _dl_flag(d_lbls, "showLegendKey")
    show_bub = _dl_flag(d_lbls, "showBubbleSize")
    show_lead = _dl_flag(d_lbls, "showLeaderLines")
    # If no show flag is set (e.g. old XML without explicit flags), default to showing values
    if not any([show_val, show_cat, show_ser, show_pct, show_bub]):
        show_val = True

    sep = None
    try:
        sep_el = _oxml_find_c(d_lbls, "separator")
        if sep_el is not None and sep_el.text:
            sep = sep_el.text
    except Exception:
        pass

    return DataLabels(
        show=True,
        format=safe_get(lambda: d_lbls.numFmt.formatCode),
        position=enum_name(safe_get(lambda: d_lbls.dLblPos.val)),
        point_positions={
            safe_get(lambda d_lbl=d_lbl: d_lbl.idx.val): enum_name(safe_get(lambda d_lbl=d_lbl: d_lbl.dLblPos.val))
            for d_lbl in safe_get(lambda: d_lbls.dLbl_lst, []) or []
        },
        font_name=safe_get(lambda: d_lbls.defRPr.latin.typeface),
        font_size=_font_size_from_defrpr(safe_get(lambda: d_lbls.defRPr)),
        font_bold=safe_get(lambda: d_lbls.defRPr.b),
        font_color=_color_from_defrpr(safe_get(lambda: d_lbls.defRPr), theme_colors),
        show_val=show_val,
        show_cat_name=show_cat,
        show_ser_name=show_ser,
        show_percent=show_pct,
        show_legend_key=show_key,
        show_bubble_size=show_bub,
        show_leader_lines=show_lead,
        separator=sep,
    )


def _series_x_values(series: Any) -> list[Any]:
    x_val = safe_get(lambda: series._element.xVal)
    if x_val is None:
        cat = safe_get(lambda: series._element.cat)
        return _cached_values(cat)
    return _cached_values(x_val)


def _series_negative_color(series: Any, theme_colors: "dict | None" = None) -> "ColorSpec | None":
    return _chart_fill_color_from_parent(safe_get(lambda: series._element.invertSolidFillFmt.spPr), theme_colors)


def _point_colors(series: Any, theme_colors: "dict | None" = None) -> list:
    points = safe_get(lambda: series._element.dPt_lst, []) or []
    return [
        color
        for color in (_chart_fill_color_from_parent(safe_get(lambda point=point: point.spPr), theme_colors) for point in points)
        if color is not None
    ]


def _point_formatting(series: Any, theme_colors: "dict | None" = None) -> dict[int, dict[str, Any]]:
    formatting = {}
    for point in safe_get(lambda: series._element.dPt_lst, []) or []:
        index = safe_get(lambda point=point: point.idx.val)
        if index is None:
            continue
        formatting[int(index)] = {
            "fill_color": _chart_fill_color_from_parent(safe_get(lambda point=point: point.spPr), theme_colors),
            "fill_type": _chart_fill_type(safe_get(lambda point=point: point.spPr)),
        }
    return formatting


def _custom_data_labels(series: Any) -> dict[int, str]:
    from lxml import etree as _etree
    _A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    labels = {}
    for d_lbl in safe_get(lambda: series._element.dLbls.dLbl_lst, []) or []:
        index = safe_get(lambda d_lbl=d_lbl: d_lbl.idx.val)
        if index is None:
            continue
        try:
            raw = _etree.fromstring(_etree.tostring(d_lbl))
            texts = [t.text or "" for t in raw.findall(f".//{{{_A}}}t")]
            text = "".join(texts)
        except Exception:
            text = ""
        if text:
            labels[int(index)] = text
    return labels


def _cached_values(parent: Any) -> list[Any]:
    if parent is None:
        return []
    values = safe_get(lambda: parent.xpath(".//c:pt/c:v/text()"), []) or []
    return values


def _font_size_from_defrpr(defrpr: Any) -> float | None:
    size = safe_get(lambda: defrpr.sz)
    return round(size / 100.0, 2) if size is not None else None


def _color_from_defrpr(defrpr: Any, theme_colors: "dict | None" = None) -> "ColorSpec | None":
    return _extract_color_spec(defrpr, theme_colors)


def _bodypr_rotation(body_pr: Any) -> float | None:
    rotation = safe_get(lambda: body_pr.rot)
    return round(rotation / 60000.0, 2) if rotation is not None else None


def _points(value: Any) -> float | None:
    points = getattr(value, "pt", None)
    return round(float(points), 2) if points is not None else None


def _float_or_none(value: Any) -> float | None:
    try:
        return None if value is None else float(value)
    except Exception:
        return None


def _is_number(value: Any) -> bool:
    try:
        float(value)
        return True
    except Exception:
        return False


def _is_text_element(shape: Any, shape_type: str | None) -> bool:
    if not safe_get(lambda: shape.has_text_frame, False):
        return False
    return shape_type in {"TEXT_BOX", "PLACEHOLDER"}


def _end_para_font_size(paragraph: Any) -> float | None:
    """Return the endParaRPr sz (in points) if present, else None."""
    try:
        _A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        end_pr = paragraph._p.find(f"{{{_A}}}endParaRPr")
        if end_pr is not None:
            sz = end_pr.get("sz")
            if sz is not None:
                return int(sz) / 100.0  # hundredths of a point → points
    except Exception:
        pass
    return None


def _text_paragraphs(shape: Any, ctx: _OnboardContext) -> list[TextParagraph]:
    if not safe_get(lambda: shape.has_text_frame, False):
        return []
    result = []
    for paragraph in shape.text_frame.paragraphs:
        runs = _runs_with_breaks(paragraph, lambda run: _onboard_text_run(shape, paragraph, run, ctx))
        resolved = _resolved_paragraph_values(shape, paragraph)
        pPr = _paragraph_pPr(paragraph)
        btype = resolved.get("bullet_type") or "none"
        blip_bytes, blip_ext = _para_bullet_blip(pPr, ctx.slide_part) if btype == "image" else (None, None)
        result.append(TextParagraph(
            runs=runs,
            alignment=resolved.get("alignment"),
            indent_level=safe_get(lambda p=paragraph: p.level, 0) or 0,
            line_spacing=resolved.get("line_spacing"),
            space_before=resolved.get("space_before") if resolved.get("space_before") is not None else _para_space_pt(paragraph, "space_before"),
            space_after=resolved.get("space_after") if resolved.get("space_after") is not None else _para_space_pt(paragraph, "space_after"),
            left_indent=resolved.get("left_indent"),
            first_line_indent=resolved.get("first_line_indent"),
            bullet_type=btype,
            bullet_char=resolved.get("bullet_char"),
            bullet_font=resolved.get("bullet_font"),
            bullet_blip_bytes=blip_bytes,
            bullet_blip_ext=blip_ext,
            end_para_font_size=_end_para_font_size(paragraph) or (resolved.get("font_size") if not runs else None),
        ))
    return result


def _resolved_paragraph_values(shape: Any, paragraph: Any) -> dict:
    """Resolve all paragraph-level properties via the inheritance chain."""
    runs = list(paragraph.runs)
    run = runs[0] if runs else None
    return resolve_text_shape(shape, paragraph, run)


def _text_body_insets(shape: Any) -> dict:
    """Read explicit lIns/rIns/tIns/bIns from bodyPr (only present attributes)."""
    try:
        body_pr = shape.text_frame._txBody.bodyPr
        result = {}
        for attr, key in (("lIns", "left"), ("rIns", "right"), ("tIns", "top"), ("bIns", "bottom")):
            val = body_pr.get(attr)
            if val is not None:
                result[key] = emu_to_inches(int(val))
        return result
    except Exception:
        return {}


def _text_frame(shape: Any) -> TextFrame:
    font_scale, ln_spc_reduction = _read_norm_autofit(shape)
    return TextFrame(
        word_wrap=safe_get(lambda: shape.text_frame.word_wrap) is not False,
        autofit_type=enum_name(safe_get(lambda: shape.text_frame.auto_size)) or "shrink",
        vertical_anchor=_resolve_vertical_anchor(shape),
        font_scale=font_scale,
        ln_spc_reduction=ln_spc_reduction,
        body_insets=_text_body_insets(shape),
    )


def _read_font_scale(shape: Any) -> int | None:
    return _read_norm_autofit(shape)[0]


def _read_norm_autofit(shape: Any) -> tuple[int | None, int | None]:
    """Return (fontScale, lnSpcReduction) from normAutofit element, or (None, None)."""
    try:
        body_pr = shape.text_frame._txBody.bodyPr
        _A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        norm = body_pr.find(f"{{{_A}}}normAutofit")
        if norm is not None:
            fs_val = norm.get("fontScale")
            ln_val = norm.get("lnSpcReduction")
            return (int(fs_val) if fs_val is not None else None,
                    int(ln_val) if ln_val is not None else None)
    except Exception:
        pass
    return (None, None)


def _resolve_vertical_anchor(shape: Any) -> str | None:
    """Resolve vertical anchor from shape bodyPr XML (not pptx-python default), falling back through layout/master."""
    _ANCHOR_MAP = {"t": "TOP", "ctr": "MIDDLE", "b": "BOTTOM", "just": "MIXED"}
    # Check the actual XML attribute — pptx-python's vertical_anchor returns a default even when absent
    try:
        body_pr = _oxml_find_bodyPr(shape.text_frame._txBody)
        if body_pr is not None and body_pr.get("anchor"):
            return _ANCHOR_MAP.get(body_pr.get("anchor"), body_pr.get("anchor").upper())
    except Exception:
        pass
    # Fall back through layout → master placeholder for anchor
    if safe_get(lambda: shape.is_placeholder, False):
        inherited = resolve_body_pr(shape)
        anchor = inherited.get("anchor")
        if anchor:
            return _ANCHOR_MAP.get(anchor, anchor.upper())
    return None


def _shape_text_frame(shape: Any) -> ShapeTextFrame:
    if not safe_get(lambda: shape.has_text_frame, False):
        return ShapeTextFrame()
    _font_scale, _ln_spc = _read_norm_autofit(shape)
    return ShapeTextFrame(
        word_wrap=safe_get(lambda: shape.text_frame.word_wrap) is not False,
        autofit_type=enum_name(safe_get(lambda: shape.text_frame.auto_size)),
        vertical_anchor=enum_name(safe_get(lambda: shape.text_frame.vertical_anchor)),
        text_insets={
            "left": emu_to_inches(safe_get(lambda: shape.text_frame.margin_left)) or 0.0,
            "right": emu_to_inches(safe_get(lambda: shape.text_frame.margin_right)) or 0.0,
            "top": emu_to_inches(safe_get(lambda: shape.text_frame.margin_top)) or 0.0,
            "bottom": emu_to_inches(safe_get(lambda: shape.text_frame.margin_bottom)) or 0.0,
        },
        font_scale=_font_scale,
        ln_spc_reduction=_ln_spc,
    )


def _geometry_preset(shape: Any) -> str | None:
    auto_shape_type = safe_get(lambda: shape.auto_shape_type)
    if auto_shape_type is not None:
        return safe_get(lambda: MSO_AUTO_SHAPE_TYPE.to_xml(auto_shape_type))
    preset = safe_get(lambda: shape.element.spPr.prstGeom.prst)
    if preset is None:
        return None
    if not isinstance(preset, str):
        return safe_get(lambda: MSO_AUTO_SHAPE_TYPE.to_xml(preset), str(preset))
    return preset


def _geometry_adjustments(shape: Any) -> dict[str, Any]:
    guides = safe_get(lambda: shape.element.spPr.prstGeom.gd_lst, []) or []
    return {
        safe_get(lambda gd=guide: gd.name): safe_get(lambda gd=guide: gd.fmla)
        for guide in guides
        if safe_get(lambda gd=guide: gd.name) is not None
    }


def _fill_gradient_stops(shape: Any, theme_colors: "dict | None" = None) -> list[GradientStop]:
    """Extract gradient stops from a shape's gradFill XML. Returns [] for non-gradient fills."""
    stops: list[GradientStop] = []
    try:
        spPr = shape.element.spPr
        grad_fill = _oxml_find_gradFill(spPr)
        if grad_fill is None:
            return stops
        gs_lst = _oxml_find_gsLst(grad_fill)
        if gs_lst is None:
            return stops
        for gs in _oxml_findall_gs(gs_lst):
            pos = int(gs.get("pos", 0)) / 100000.0
            # Gradient stop colors are stored directly as srgbClr/schemeClr/sysClr,
            # not wrapped in solidFill like shape fills.
            color_child = None
            for color_tag in ("srgbClr", "schemeClr", "sysClr"):
                color_child = _oxml_find_a(gs, color_tag)
                if color_child is not None:
                    break
            cs = _extract_color_spec(color_child, theme_colors) if color_child is not None else _extract_color_spec(gs, theme_colors)
            if cs is not None:
                stops.append(GradientStop(position=pos, color=cs))
    except Exception:
        pass
    return stops


def _fill_gradient_angle(shape: Any) -> float:
    """Extract linear gradient angle in degrees (OOXML 60ths-of-a-degree → degrees)."""
    try:
        spPr = shape.element.spPr
        grad_fill = _oxml_find_gradFill(spPr)
        if grad_fill is None:
            return 0.0
        lin = _oxml_find_lin(grad_fill)
        if lin is not None:
            return int(lin.get("ang", 0)) / 60000.0
    except Exception:
        pass
    return 0.0


def _fill_type(shape: Any) -> str | None:
    ft = enum_name(safe_get(lambda: shape.fill.type))
    # python-pptx returns None or BACKGROUND when the spPr has a scheme-color solidFill;
    # always check raw XML so explicit spPr fills override style references.
    try:
        spPr = shape.element.spPr
        if spPr is not None:
            if _oxml_find_solidFill(spPr) is not None:
                return "solidFill"
            if _oxml_find_gradFill(spPr) is not None:
                return "gradFill"
            if _oxml_find_noFill(spPr) is not None:
                return "noFill"
            if _oxml_find_grpFill(spPr) is not None:
                return "grpFill"
    except Exception:
        pass
    return ft


def _shape_has_explicit_fill(shape: Any) -> bool:
    """Return True only if the shape's own spPr XML has an explicit fill (not noFill, not inherited).

    Uses raw XML rather than python-pptx's effective fill (which walks theme inheritance and
    can return a solid fill for shapes that have no fill in their own XML).
    """
    try:
        spPr = shape.element.spPr
        if spPr is None:
            return False
        # Explicit noFill → definitely no fill
        if _oxml_find_noFill(spPr) is not None:
            return False
        # Any explicit fill element → has fill
        for tag in ("solidFill", "gradFill", "blipFill", "pattFill"):
            if _oxml_find_a(spPr, tag) is not None:
                return True
        return False
    except Exception:
        # Fall back to python-pptx effective fill, filtering out background/inherit
        return _fill_type(shape) not in (None, "BACKGROUND")


def _fill_color(shape: Any, theme_colors: "dict | None" = None) -> "ColorSpec | None":
    from percy.bridge.elements import ColorSpec
    try:
        spPr = shape.element.spPr
        spec = _extract_color_spec(spPr, theme_colors)
        if spec is not None:
            return spec
    except Exception:
        pass
    color = safe_get(lambda: shape.fill.fore_color)
    rgb = safe_get(lambda: color.rgb)
    if rgb is not None:
        return ColorSpec(value=f"#{str(rgb)}")
    scheme = _fill_scheme(shape)
    if not scheme:
        return None
    if theme_colors:
        raw = theme_colors.get(scheme, "")
        if raw:
            return ColorSpec(value=raw if raw.startswith("#") else f"#{raw}")
    return ColorSpec(value=f"scheme:{scheme}")


def _fill_scheme(shape: Any) -> str | None:
    color = safe_get(lambda: shape.fill.fore_color)
    theme = safe_get(lambda: color.theme_color)
    return enum_name(theme) if theme is not None else None


def _fill_pattern_preset(shape: Any) -> str | None:
    try:
        spPr = shape.element.spPr
        if spPr is None:
            return None
        patt = _oxml_find_pattFill(spPr)
        return patt.get("prst") if patt is not None else None
    except Exception:
        return None


def _fill_bg_color(shape: Any, theme_colors: "dict | None" = None) -> "ColorSpec | None":
    try:
        spPr = shape.element.spPr
        if spPr is None:
            return None
        patt = _oxml_find_pattFill(spPr)
        if patt is None:
            return None
        bg_el = _oxml_find_bgClr(patt)
        if bg_el is None:
            return None
        # bgClr/fgClr contain a color element directly (schemeClr/srgbClr), not solidFill
        for child in bg_el:
            spec = _extract_color_spec(child, theme_colors)
            if spec is not None:
                return spec
        return None
    except Exception:
        return None


def _fill_transparency(shape: Any) -> float | None:
    return safe_get(lambda: shape.fill.transparency)


def _shape_line(shape: Any, theme_colors: "dict | None" = None) -> ShapeLine:
    return ShapeLine(
        visible=_line_visible(shape),
        color=_line_color(shape, theme_colors),
        width=_line_width_points(shape),
        dash_style=_line_dash(shape) or "solid",
        head_end=_line_end(shape, "tailEnd"),
        tail_end=_line_end(shape, "headEnd"),
        head_size=_line_end_size(shape, "tailEnd"),
        tail_size=_line_end_size(shape, "headEnd"),
    )


def _freeform_geometry_xml(shape: Any) -> str | None:
    """Extract the custGeom XML blob from a freeform shape for lossless roundtrip."""
    try:
        from lxml import etree
        custGeom = shape.element.spPr.custGeom
        if custGeom is None:
            return None
        return etree.tostring(custGeom, encoding="unicode")
    except Exception:
        return None


def _freeform_paths(shape: Any) -> list[FreeformPath]:
    paths = []
    path_nodes = safe_get(lambda: shape.element.spPr.custGeom.pathLst)
    if path_nodes is None:
        return paths
    for path_node in path_nodes:
        commands = []
        for child in path_node:
            command_name = _local_name(child.tag)
            if command_name in {"moveTo", "lnTo"}:
                point = _path_point(child)
                commands.append(PathCommand(command=command_name, points=[point] if point else []))
            elif command_name in {"quadBezTo", "cubicBezTo"}:
                commands.append(PathCommand(command=command_name, points=_path_points(child)))
            elif command_name == "arcTo":
                commands.append(
                    PathCommand(
                        command=command_name,
                        arc_params={
                            key: _int_or_zero(child.get(key))
                            for key in ("wR", "hR", "stAng", "swAng")
                            if child.get(key) is not None
                        },
                    )
                )
            elif command_name == "close":
                commands.append(PathCommand(command=command_name))
        paths.append(
            FreeformPath(
                width=_int_or_zero(path_node.get("w")),
                height=_int_or_zero(path_node.get("h")),
                fill_mode=path_node.get("fill"),
                stroke=path_node.get("stroke") != "0",
                commands=commands,
            )
        )
    return paths


def _path_points(node: Any) -> list[tuple[int, int]]:
    """Extract all <a:pt> children from a bezier node (cubicBezTo / quadBezTo).

    Bezier nodes contain <a:pt> elements directly as children, so we read x/y
    attributes from each child whose local tag name is 'pt'.
    """
    pts = []
    for child in node:
        if _local_name(child.tag) == "pt":
            try:
                pts.append((_int_or_zero(child.get("x")), _int_or_zero(child.get("y"))))
            except Exception:
                pass
    return pts


def _path_point(node: Any) -> tuple[int, int] | None:
    """Extract the single <a:pt> coordinate from a moveTo or lnTo node."""
    for child in node:
        if _local_name(child.tag) == "pt":
            try:
                return (_int_or_zero(child.get("x")), _int_or_zero(child.get("y")))
            except Exception:
                return None
    return None


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def _int_or_zero(value: Any) -> int:
    try:
        return int(value)
    except Exception:
        return 0


def _line_visible(shape: Any) -> bool:
    """True only if the shape's OWN XML has an <a:ln> with a solid fill (not noFill)."""
    ln_el = safe_get(lambda: shape.element.spPr.ln)
    if ln_el is None:
        return False
    # noFill child means explicitly no border
    if _oxml_find_noFill(ln_el) is not None:
        return False
    # solidFill or gradFill → visible
    return (_oxml_find_solidFill(ln_el) is not None
            or _oxml_find_gradFill(ln_el) is not None)


def _line_color(shape: Any, theme_colors: "dict | None" = None) -> "ColorSpec | None":
    """Read border/line color only from the shape's OWN <a:ln> XML.

    We deliberately avoid shape.line.color (which resolves through layout/master
    inheritance) — inherited colors must not be treated as the shape's own border.
    """
    try:
        ln_el = safe_get(lambda: shape.element.spPr.ln)
        if ln_el is not None:
            spec = _extract_color_spec(ln_el, theme_colors)
            if spec is not None:
                return spec
    except Exception:
        pass
    return None


def _line_scheme(shape: Any) -> str | None:
    color = safe_get(lambda: shape.line.color)
    theme = safe_get(lambda: color.theme_color)
    return enum_name(theme) if theme is not None else None


def _line_width_points(shape: Any) -> float | None:
    width = safe_get(lambda: shape.line.width)
    return round(width.pt, 2) if width is not None and getattr(width, "pt", None) is not None else None


def _line_dash(shape: Any) -> str | None:
    return enum_name(safe_get(lambda: shape.line.dash_style))


def _line_cap(shape: Any) -> str | None:
    return safe_get(lambda: shape.element.spPr.ln.cap)


def _line_join(shape: Any) -> str | None:
    ln = safe_get(lambda: shape.element.spPr.ln)
    if ln is None:
        return None
    for child in ln:
        name = _local_name(child.tag)
        if name in {"round", "bevel", "miter"}:
            return name
    return None


def _line_end(shape: Any, tag_local_name: str) -> str | None:
    ln = safe_get(lambda: shape.element.spPr.ln)
    if ln is None:
        return None
    for child in ln:
        if child.tag.endswith(tag_local_name):
            return child.get("type")
    return None


def _line_end_size(shape: Any, tag_local_name: str) -> str | None:
    """Return 'w/len' size string for a headEnd/tailEnd element, e.g. 'med/med'."""
    ln = safe_get(lambda: shape.element.spPr.ln)
    if ln is None:
        return None
    for child in ln:
        if child.tag.endswith(tag_local_name):
            w = child.get("w")
            length = child.get("len")
            if w or length:
                return f"{w or 'med'}/{length or 'med'}"
    return None


def _connector_type(shape: Any) -> str:
    preset = _geometry_preset(shape)
    if preset in {"bentConnector2", "bentConnector3", "bentConnector4", "bentConnector5"}:
        return "elbow"
    if preset in {"curvedConnector2", "curvedConnector3", "curvedConnector4", "curvedConnector5"}:
        return "curved"
    return "straight"


def _image_crop(shape: Any, attr_name: str) -> float:
    value = safe_get(lambda: shape.element.blipFill.srcRect.get(attr_name))
    return round(int(value) / 100000.0, 5) if value is not None else 0.0


def _image_shape_geometry(shape: Any) -> tuple[str | None, dict[str, str]]:
    """Return (prstGeom_prst, adj_dict) for a picture shape.

    Checks the shape's own spPr first, then the layout placeholder's spPr when
    the shape is a content placeholder (e.g. roundRect picture placeholders).
    Returns (None, {}) for plain rectangles.
    """
    def _geom_from_sp_pr(sp_pr: Any) -> tuple[str | None, dict[str, str]]:
        if sp_pr is None:
            return None, {}
        geom = _oxml_find_prstGeom(sp_pr)
        if geom is None:
            return None, {}
        prst = geom.get("prst")
        if not prst or prst == "rect":
            return None, {}
        adj: dict[str, str] = {}
        av_lst = _oxml_find_avLst(geom)
        if av_lst is not None:
            for gd in _oxml_findall_gd(av_lst):
                name = gd.get("name")
                fmla = gd.get("fmla")
                if name and fmla:
                    adj[name] = fmla
        return prst, adj

    try:
        prst, adj = _geom_from_sp_pr(safe_get(lambda: shape.element.spPr))
        if prst:
            return prst, adj
    except Exception:
        pass

    # Check layout placeholder inheritance
    try:
        ph_idx = shape.placeholder_format.idx
        for lshape in shape.part.slide_layout.placeholders:
            if lshape.placeholder_format.idx == ph_idx:
                prst, adj = _geom_from_sp_pr(safe_get(lambda: lshape._element.spPr))
                if prst:
                    return prst, adj
                break
    except Exception:
        pass

    return None, {}


def _extract_outer_shadow(shape: Any, theme_colors: "dict | None" = None) -> "ShapeShadow":
    """Extract outerShdw from effectLst or effectDag on a shape element."""
    try:
        sp_pr = safe_get(lambda: shape.element.spPr)
        if sp_pr is None:
            # For pictures the element is a <p:pic>, check directly
            sp_pr = safe_get(lambda: _oxml_find_p_spPr(shape.element))
        if sp_pr is None:
            return ShapeShadow()

        effect_lst = _oxml_find_effectLst(sp_pr)
        if effect_lst is None:
            # Also check direct on element (charts/pics sometimes put it there)
            effect_lst = safe_get(lambda: _oxml_find_descendant_a(shape.element, "effectLst"))
        if effect_lst is None:
            return ShapeShadow()

        shdw = _oxml_find_outerShdw(effect_lst)
        if shdw is None:
            return ShapeShadow()

        blur_rad = shdw.get("blurRad")
        dist = shdw.get("dist")
        dir_ang = shdw.get("dir")
        algn = shdw.get("algn")
        rot_w = shdw.get("rotWithShape")

        color_spec = _extract_color_spec(shdw, theme_colors)

        return ShapeShadow(
            has_shadow=True,
            blur=round(int(blur_rad) / 12700.0, 2) if blur_rad else None,
            distance=round(int(dist) / 12700.0, 2) if dist else None,
            direction=round(int(dir_ang) / 60000.0, 2) if dir_ang else None,
            color=color_spec,
            alpha=color_spec.alpha if color_spec else None,
            align=algn,
            rot_with_shape=(rot_w == "1") if rot_w is not None else False,
        )
    except Exception:
        return ShapeShadow()


def _onboard_text_run(shape: Any, paragraph: Any, run: Any, ctx: _OnboardContext) -> TextRun:
    resolved = resolve_text_shape(shape, paragraph, run)
    return TextRun(
        text=run.text,
        font_name=ctx.resolve_font_name(resolved.get("font_name")),
        font_size=resolved.get("font_size"),
        font_bold=resolved.get("font_bold"),
        font_italic=resolved.get("font_italic"),
        font_underline=resolved.get("font_underline"),
        font_color=resolved.get("font_color"),
        font_caps=resolved.get("font_caps"),
        char_spacing=resolved.get("char_spacing"),
        baseline_shift=_ooxml_baseline_to_fraction(resolved.get("baseline_shift")),
        strikethrough=resolved.get("strikethrough"),
        hyperlink=_run_hyperlink(run),
    )


def _resolved_paragraph_value(shape: Any, paragraph: Any, key: str) -> Any:
    runs = list(paragraph.runs)
    run = runs[0] if runs else None
    return resolve_text_shape(shape, paragraph, run).get(key)
