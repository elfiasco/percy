"""Chart-focused corpus diagnostics."""

from __future__ import annotations

from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation

from percy.bridge import BridgeChart
from percy.diagnostics.common import ensure_dir, enum_name, safe_get, write_json
from percy.diagnostics.onboard import onboard_pptx


def analyze_charts(input_dir: str | Path, out_dir: str | Path | None = None) -> dict[str, Any]:
    pptx_paths = sorted(Path(input_dir).glob("*.pptx"))
    report: dict[str, Any] = {
        "input_dir": str(input_dir),
        "pptx_count": len(pptx_paths),
        "chart_count": 0,
        "chart_types": {},
        "semantic_debt_counts": {},
        "charts": [],
    }
    chart_type_counts: Counter[str] = Counter()
    debt_counts: Counter[str] = Counter()

    for pptx_path in pptx_paths:
        document = onboard_pptx(pptx_path)
        source_facts = _source_chart_facts(pptx_path)
        source_index = 0
        for slide in document.slides:
            for element in slide.elements:
                if not isinstance(element, BridgeChart):
                    continue
                facts = source_facts[source_index] if source_index < len(source_facts) else {}
                source_index += 1
                debt = element.custom_properties.get("semantic_debt", [])
                chart_type = element.chart_type or "unknown"
                chart_type_counts[chart_type] += 1
                debt_counts.update(debt)
                report["charts"].append(
                    {
                        "pptx": str(pptx_path),
                        "slide_number": slide.slide_number,
                        "shape_id": element.identification.shape_id,
                        "shape_name": element.identification.shape_name,
                        "chart_type": chart_type,
                        "plot_count": facts.get("plot_count"),
                        "plot_tags": facts.get("plot_tags", []),
                        "series_count": len(element.series),
                        "category_count": len(element.categories.categories),
                        "has_title": element.title.title is not None,
                        "has_legend": element.legend.visible,
                        "has_category_axis": element.category_axis.visible,
                        "has_value_axis": element.value_axis.visible,
                        "external_data": facts.get("external_data", False),
                        "embedded_workbook": facts.get("embedded_workbook", False),
                        "data_source_kind": element.data_source.source_kind,
                        "relationship_id": element.data_source.relationship_id,
                        "relationship_type": element.data_source.relationship_type,
                        "target_mode": element.data_source.target_mode,
                        "target": element.data_source.target,
                        "auto_update": element.data_source.auto_update,
                        "embedded_workbook_filename": element.data_source.embedded_workbook_filename,
                        "workbook_sheet_names": element.data_source.workbook_sheet_names,
                        "workbook_dimensions": element.data_source.workbook_dimensions,
                        "workbook_cell_count": sum(len(sheet.cells) for sheet in element.data_source.workbook_sheets),
                        "workbook_sheets": [
                            {
                                "name": sheet.name,
                                "dimension": sheet.dimension,
                                "cell_count": len(sheet.cells),
                                "sample_cells": [
                                    {
                                        "address": cell.address,
                                        "value": cell.value,
                                        "formula": cell.formula,
                                        "data_type": cell.data_type,
                                    }
                                    for cell in sheet.cells[:12]
                                ],
                            }
                            for sheet in element.data_source.workbook_sheets
                        ],
                        "cache_series_count": element.data_source.cache_series_count,
                        "cache_category_count": element.data_source.cache_category_count,
                        "cache_point_count": element.data_source.cache_point_count,
                        "formula_count": len(element.data_source.formulas),
                        "formulas": element.data_source.formulas,
                        "combo_chart": facts.get("plot_count", 0) > 1,
                        "series": [
                            {
                                "name": series.name,
                                "value_count": len(series.values),
                                "x_value_count": len(series.x_values),
                                "plot_type": series.plot_type,
                                "has_data_labels": series.data_labels.show,
                                "point_format_count": len(series.point_formatting),
                            }
                            for series in element.series
                        ],
                        "semantic_debt": debt,
                    }
                )

    report["chart_count"] = len(report["charts"])
    report["chart_types"] = dict(chart_type_counts.most_common())
    report["semantic_debt_counts"] = dict(debt_counts.most_common())
    if out_dir is not None:
        output_dir = ensure_dir(out_dir)
        write_json(report, output_dir / "chart-audit.json")
    return report


def _source_chart_facts(pptx_path: Path) -> list[dict[str, Any]]:
    presentation = Presentation(str(pptx_path))
    facts = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if not safe_get(lambda shape=shape: shape.has_chart, False):
                continue
            chart = shape.chart
            plot_tags = [
                _local_name(safe_get(lambda plot=plot: plot._element.tag, ""))
                for plot in safe_get(lambda: chart.plots, []) or []
            ]
            facts.append(
                {
                    "slide_number": slide_number,
                    "shape_id": safe_get(lambda shape=shape: shape.shape_id),
                    "chart_type": enum_name(safe_get(lambda: chart.chart_type)),
                    "plot_count": len(plot_tags),
                    "plot_tags": plot_tags,
                    "external_data": safe_get(lambda: chart._chartSpace.externalData) is not None,
                    "embedded_workbook": safe_get(lambda: chart.part.chart_workbook.xlsx_part) is not None,
                    "style": safe_get(lambda: chart.chart_style),
                }
            )
    return facts


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]
