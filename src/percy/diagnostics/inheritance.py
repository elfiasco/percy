"""Resolve inherited PowerPoint XML formatting into explicit values."""

from __future__ import annotations

from dataclasses import dataclass
from typing import Any

from lxml import etree
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from percy.diagnostics.common import enum_name, length_to_points, safe_get
from percy.oxml import (
    A_NS,
    P_NS,
    NS_MAP as NS,
    qa,
    find_pPr,
    find_rPr,
    find_buChar,
    find_buFont,
    find_solidFill,
    find_srgbClr,
    find_sysClr,
    find_schemeClr,
)

TITLE_PLACEHOLDERS = {
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
    PP_PLACEHOLDER.VERTICAL_TITLE,
}
BODY_PLACEHOLDERS = {
    PP_PLACEHOLDER.BODY,
    PP_PLACEHOLDER.OBJECT,
    PP_PLACEHOLDER.VERTICAL_BODY,
}


@dataclass(frozen=True)
class ResolvedValue:
    value: Any
    source: str


def resolve_text_shape(shape: Any, paragraph: Any | None = None, run: Any | None = None) -> dict[str, Any]:
    """Resolve effective text properties for a run/paragraph on *shape*.

    The returned values are intended to be copied into Bridge elements. The
    ``sources`` map is retained in diagnostics so we can reverse-engineer which
    Open XML layer contributed each value.
    """

    paragraph_level = safe_get(lambda: paragraph.level, 0) or 0
    theme = _theme_info(shape)
    sources = _style_sources(shape, paragraph_level)

    resolved = {
        "font_name": _resolve_font_name(shape, run, sources, theme),
        "font_size": _resolve_attr("font_size", run, sources, _font_size_from_rpr),
        "font_bold": _resolve_bool_attr("b", run, sources, safe_get(lambda: run.font.bold)),
        "font_italic": _resolve_bool_attr("i", run, sources, safe_get(lambda: run.font.italic)),
        "font_underline": _resolve_underline(run, sources),
        "font_color": _resolve_color(shape, run, sources, theme),
        "alignment": _resolve_alignment(paragraph, sources),
        "font_caps": _resolve_run_str_attr("cap", run, sources),
        "char_spacing": _resolve_run_int_attr("spc", run, sources),
        "baseline_shift": _resolve_run_int_attr("baseline", run, sources),
        "strikethrough": _resolve_run_str_attr("strike", run, sources),
        "line_spacing": _resolve_para_line_spacing(paragraph, sources),
        "space_before": _resolve_para_spacing("spcBef", paragraph, sources),
        "space_after": _resolve_para_spacing("spcAft", paragraph, sources),
        "left_indent": _resolve_para_emu_attr("marL", paragraph, sources),
        "first_line_indent": _resolve_para_emu_attr("indent", paragraph, sources),
        "bullet_type": _resolve_para_bullet_type(paragraph, sources),
        "bullet_char": _resolve_para_bullet_char(paragraph, sources),
        "bullet_font": _resolve_para_bullet_font(paragraph, sources),
    }
    unresolved = [key for key, value in resolved.items() if value.value is None]
    return {
        key: value.value for key, value in resolved.items()
    } | {
        "sources": {key: value.source for key, value in resolved.items()},
        "unresolved": unresolved,
        "placeholder": placeholder_info(shape),
        "style_source_chain": [source["name"] for source in sources],
    }


def resolve_body_pr(shape: Any) -> dict[str, Any]:
    """Resolve bodyPr attributes with inheritance: slide → layout placeholder → master placeholder."""
    result: dict[str, Any] = {}
    try:
        txBody = shape.text_frame._txBody
        bodyPr = txBody.find(qa("bodyPr"))
        if bodyPr is not None:
            result["wrap"] = bodyPr.get("wrap")
            if bodyPr.get("anchor"):
                result["anchor"] = bodyPr.get("anchor")
                return result
    except Exception:
        pass
    # Fall through to layout then master placeholder for anchor (shapes use p:txBody not a:txBody)
    from percy.oxml import qp
    for match_fn in (_matching_layout_placeholder, _matching_master_placeholder):
        match = match_fn(shape)
        if match is not None:
            try:
                txBody = match.element.find(qp("txBody"))
                if txBody is not None:
                    bodyPr = txBody.find(qa("bodyPr"))
                    if bodyPr is not None and bodyPr.get("anchor"):
                        result["anchor"] = bodyPr.get("anchor")
                        return result
            except Exception:
                pass
    return result


def placeholder_info(shape: Any) -> dict[str, Any]:
    if not safe_get(lambda: shape.is_placeholder, False):
        return {"is_placeholder": False}
    placeholder = shape.placeholder_format
    return {
        "is_placeholder": True,
        "type": enum_name(safe_get(lambda: placeholder.type)),
        "idx": safe_get(lambda: placeholder.idx),
    }


def _style_sources(shape: Any, paragraph_level: int) -> list[dict[str, Any]]:
    level_tag = f"lvl{min(max(paragraph_level, 0), 8) + 1}pPr"
    sources: list[dict[str, Any]] = []

    shape_element = shape.element
    sources.extend(_shape_text_sources("slide", shape_element, level_tag))

    is_placeholder = safe_get(lambda: shape.is_placeholder, False)

    layout_match = _matching_layout_placeholder(shape)
    if layout_match is not None:
        sources.extend(_shape_text_sources("layout-placeholder", layout_match.element, level_tag))

    master_match = _matching_master_placeholder(shape)
    if master_match is not None:
        sources.extend(_shape_text_sources("master-placeholder", master_match.element, level_tag))

    master = safe_get(lambda: shape.part.slide_layout.slide_master)
    if master is not None:
        if is_placeholder:
            # Placeholder shapes: use master txStyles (titleStyle/bodyStyle/otherStyle)
            style_name = _master_style_name(shape)
            master_style = _first(_xpath(master.element, f".//p:txStyles/p:{style_name}/a:{level_tag}"))
            if master_style is not None:
                sources.append({"name": f"master-txStyles:{style_name}:{level_tag}", "pPr": master_style, "rPr": _def_rpr(master_style)})
            if level_tag != "lvl1pPr":
                lvl1_style = _first(_xpath(master.element, f".//p:txStyles/p:{style_name}/a:lvl1pPr"))
                if lvl1_style is not None:
                    sources.append({"name": f"master-txStyles:{style_name}:lvl1pPr:fallback", "pPr": lvl1_style, "rPr": _def_rpr(lvl1_style)})
        else:
            # Non-placeholder shapes (text boxes, auto shapes): use presentation defaultTextStyle
            prs_el = safe_get(lambda: shape.part.package.presentation_part._element)
            if prs_el is not None:
                from percy.oxml import qp
                dts = prs_el.find(qp("defaultTextStyle"))
                if dts is not None:
                    lvl_pr = dts.find(qa(level_tag))
                    if lvl_pr is not None:
                        sources.append({"name": f"prs:defaultTextStyle:{level_tag}", "pPr": lvl_pr, "rPr": _def_rpr(lvl_pr)})
                    if level_tag != "lvl1pPr":
                        lvl1_pr = dts.find(qa("lvl1pPr"))
                        if lvl1_pr is not None:
                            sources.append({"name": "prs:defaultTextStyle:lvl1pPr:fallback", "pPr": lvl1_pr, "rPr": _def_rpr(lvl1_pr)})

    return sources


def _shape_text_sources(name: str, element: Any, level_tag: str) -> list[dict[str, Any]]:
    sources = []
    # Slide/layout/master shapes use <p:txBody>, not <a:txBody>
    p_pr = _first(_xpath(element, ".//p:txBody/a:p[1]/a:pPr"))
    if p_pr is not None:
        sources.append({"name": f"{name}:paragraph", "pPr": p_pr, "rPr": _def_rpr(p_pr)})
    lvl_pr = _first(_xpath(element, f".//p:txBody/a:lstStyle/a:{level_tag}"))
    if lvl_pr is not None:
        sources.append({"name": f"{name}:lstStyle:{level_tag}", "pPr": lvl_pr, "rPr": _def_rpr(lvl_pr)})
    return sources


def _matching_layout_placeholder(shape: Any) -> Any | None:
    if not safe_get(lambda: shape.is_placeholder, False):
        return None
    return _matching_placeholder(shape, safe_get(lambda: shape.part.slide_layout.placeholders, []))


def _matching_master_placeholder(shape: Any) -> Any | None:
    if not safe_get(lambda: shape.is_placeholder, False):
        return None
    return _matching_placeholder(shape, safe_get(lambda: shape.part.slide_layout.slide_master.placeholders, []))


def _matching_placeholder(shape: Any, placeholders: Any) -> Any | None:
    info = placeholder_info(shape)
    # First pass: exact idx match
    for placeholder in placeholders:
        ph_idx = safe_get(lambda p=placeholder: p.placeholder_format.idx)
        if ph_idx is not None and ph_idx == info.get("idx"):
            return placeholder
    # Second pass: type match only when there's a unique match (avoid wrong fallback)
    type_matches = [
        p for p in placeholders
        if enum_name(safe_get(lambda ph=p: ph.placeholder_format.type)) == info.get("type")
    ]
    return type_matches[0] if len(type_matches) == 1 else None


def _master_style_name(shape: Any) -> str:
    placeholder_type = safe_get(lambda: shape.placeholder_format.type)
    if placeholder_type in TITLE_PLACEHOLDERS:
        return "titleStyle"
    if placeholder_type in BODY_PLACEHOLDERS:
        return "bodyStyle"
    return "otherStyle"


def _theme_info(shape: Any) -> dict[str, Any]:
    master_part = safe_get(lambda: shape.part.slide_layout.slide_master.part)
    if master_part is None:
        return {}
    theme_part = safe_get(lambda: master_part.part_related_by(RT.THEME))
    if theme_part is None:
        return {}
    root = etree.fromstring(theme_part.blob)
    return {
        "major_latin": _attr(root, ".//a:fontScheme/a:majorFont/a:latin", "typeface"),
        "minor_latin": _attr(root, ".//a:fontScheme/a:minorFont/a:latin", "typeface"),
        "colors": _theme_colors(root),
    }


def _theme_colors(root: Any) -> dict[str, str]:
    colors = {}
    for color_name in ("dk1", "lt1", "dk2", "lt2", "accent1", "accent2", "accent3", "accent4", "accent5", "accent6", "hlink", "folHlink"):
        node = _first(root.xpath(f".//a:clrScheme/a:{color_name}/*[1]", namespaces=NS))
        if node is None:
            continue
        value = node.get("lastClr") if node.tag.endswith("sysClr") else node.get("val")
        if value:
            colors[color_name] = value
    colors["tx1"] = colors.get("dk1")
    colors["tx2"] = colors.get("dk2")
    colors["bg1"] = colors.get("lt1")
    colors["bg2"] = colors.get("lt2")
    return {key: value for key, value in colors.items() if value}


def _resolve_font_name(shape: Any, run: Any | None, sources: list[dict[str, Any]], theme: dict[str, Any]) -> ResolvedValue:
    explicit = safe_get(lambda: run.font.name)
    if explicit:
        return ResolvedValue(explicit, "slide:run:rPr")
    for source in sources:
        typeface = _attr(source["rPr"], "./a:latin", "typeface")
        if typeface:
            return ResolvedValue(_resolve_theme_font(typeface, theme), source["name"])
    placeholder_type = safe_get(lambda: shape.placeholder_format.type)
    if placeholder_type in TITLE_PLACEHOLDERS and theme.get("major_latin"):
        return ResolvedValue(theme["major_latin"], "theme:majorFont")
    if theme.get("minor_latin"):
        return ResolvedValue(theme["minor_latin"], "theme:minorFont")
    return ResolvedValue(None, "unresolved")


def _resolve_theme_font(typeface: str, theme: dict[str, Any]) -> str:
    if typeface == "+mj-lt":
        return theme.get("major_latin") or typeface
    if typeface == "+mn-lt":
        return theme.get("minor_latin") or typeface
    return typeface


def _resolve_attr(name: str, run: Any | None, sources: list[dict[str, Any]], extractor) -> ResolvedValue:
    if name == "font_size":
        explicit = length_to_points(safe_get(lambda: run.font.size))
        if explicit is not None:
            return ResolvedValue(explicit, "slide:run:rPr")
    for source in sources:
        value = extractor(source["rPr"])
        if value is not None:
            return ResolvedValue(value, source["name"])
    return ResolvedValue(None, "unresolved")


def _resolve_bool_attr(xml_attr: str, run: Any | None, sources: list[dict[str, Any]], explicit: Any) -> ResolvedValue:
    if explicit is not None:
        return ResolvedValue(bool(explicit), "slide:run:rPr")
    for source in sources:
        value = _bool_attr(source["rPr"], xml_attr)
        if value is not None:
            return ResolvedValue(value, source["name"])
    return ResolvedValue(False, "office-default")


def _resolve_underline(run: Any | None, sources: list[dict[str, Any]]) -> ResolvedValue:
    explicit = safe_get(lambda: run.font.underline)
    if explicit is not None:
        return ResolvedValue(bool(explicit), "slide:run:rPr")
    for source in sources:
        value = source["rPr"].get("u") if source["rPr"] is not None else None
        if value is not None:
            return ResolvedValue(value != "none", source["name"])
    return ResolvedValue(False, "office-default")


def _resolve_color(shape: Any, run: Any | None, sources: list[dict[str, Any]], theme: dict[str, Any]) -> ResolvedValue:
    explicit = _python_color(run, theme)
    if explicit:
        return ResolvedValue(explicit, "slide:run:rPr")
    for source in sources:
        value = _xml_color(source["rPr"], theme)
        if value:
            return ResolvedValue(value, source["name"])
    from percy.bridge.elements import ColorSpec
    placeholder_type = safe_get(lambda: shape.placeholder_format.type)
    scheme = "tx1" if placeholder_type in TITLE_PLACEHOLDERS | BODY_PLACEHOLDERS else "tx1"
    if theme.get("colors", {}).get(scheme):
        raw = theme["colors"][scheme]
        return ResolvedValue(ColorSpec(value=raw if raw.startswith("#") else f"#{raw}"), f"theme:{scheme}")
    return ResolvedValue(None, "unresolved")


def _resolve_alignment(paragraph: Any | None, sources: list[dict[str, Any]]) -> ResolvedValue:
    explicit = enum_name(safe_get(lambda: paragraph.alignment))
    if explicit:
        return ResolvedValue(explicit, "slide:paragraph:pPr")
    for source in sources:
        pPr = source["pPr"]
        if pPr is None:
            continue
        value = pPr.get("algn")
        if not value:
            continue
        return ResolvedValue(value, source["name"])
    return ResolvedValue("left", "office-default")


def _font_size_from_rpr(r_pr: Any) -> float | None:
    if r_pr is None:
        return None
    size = r_pr.get("sz")
    if size is None:
        return None
    return round(int(size) / 100.0, 2)


def _python_color(run: Any | None, theme: dict[str, Any] | None = None) -> "ColorSpec | None":
    from percy.bridge.elements import ColorSpec
    if run is None:
        return None
    try:
        rpr_el = find_rPr(safe_get(lambda: run._r))
        if rpr_el is not None:
            spec = _xml_color(rpr_el, theme or {})
            if spec is not None:
                return spec
    except Exception:
        pass
    rgb = safe_get(lambda: run.font.color.rgb)
    if rgb is not None:
        return ColorSpec(value=f"#{str(rgb)}")
    return None


def _xml_color(r_pr: Any, theme: dict[str, Any]) -> "ColorSpec | None":
    from percy.bridge.elements import ColorSpec
    if r_pr is None:
        return None
    _XML_MAP = {
        "dk1": "DARK_1", "dk2": "DARK_2", "lt1": "LIGHT_1", "lt2": "LIGHT_2",
        "bg1": "LIGHT_1", "bg2": "LIGHT_2", "tx1": "DARK_1", "tx2": "DARK_2",
        "accent1": "ACCENT_1", "accent2": "ACCENT_2", "accent3": "ACCENT_3",
        "accent4": "ACCENT_4", "accent5": "ACCENT_5", "accent6": "ACCENT_6",
        "hlink": "HYPERLINK", "folHlink": "FOLLOWED_HYPERLINK",
    }
    try:
        solid = find_solidFill(r_pr)
        if solid is None:
            solid = _first(_xpath(r_pr, "./a:solidFill"))
        if solid is None:
            return None
        srgb = find_srgbClr(solid)
        if srgb is not None:
            val = srgb.get("val", "")
            return ColorSpec(value=f"#{val.upper()}") if val else None
        sys_el = find_sysClr(solid)
        if sys_el is not None:
            last_clr = sys_el.get("lastClr", "")
            return ColorSpec(value=f"#{last_clr.upper()}") if last_clr else None
        scheme = find_schemeClr(solid)
        if scheme is None:
            return None
        xml_name = scheme.get("val", "")
        normalized = _XML_MAP.get(xml_name, xml_name.upper())

        def _int_val(tag: str) -> int | None:
            el = scheme.find(qa(tag))
            if el is not None:
                try:
                    return int(el.get("val", ""))
                except (ValueError, TypeError):
                    pass
            return None

        lum_mod = _int_val("lumMod")
        lum_off = _int_val("lumOff")
        shade   = _int_val("shade")
        tint    = _int_val("tint")
        alpha   = _int_val("alpha")

        # Try to resolve to a concrete hex value immediately so bridge objects
        # are self-contained and don't rely on theme_colors at serialization time.
        raw_hex = theme.get("colors", {}).get(xml_name, "")
        if raw_hex:
            base = raw_hex if raw_hex.startswith("#") else f"#{raw_hex}"
            temp = ColorSpec(value=base, lum_mod=lum_mod, lum_off=lum_off,
                             shade=shade, tint=tint)
            resolved = temp.resolve()
            if len(resolved.lstrip("#")) == 6 and resolved.upper() != "#888888":
                return ColorSpec(value=resolved, alpha=alpha)

        return ColorSpec(
            value=f"scheme:{normalized}",
            lum_mod=lum_mod, lum_off=lum_off,
            shade=shade, tint=tint, alpha=alpha,
        )
    except Exception:
        pass
    srgb = _attr(r_pr, "./a:solidFill/a:srgbClr", "val")
    if srgb:
        return ColorSpec(value=f"#{srgb.upper()}")
    scheme_name = _attr(r_pr, "./a:solidFill/a:schemeClr", "val")
    if scheme_name:
        normalized = _XML_MAP.get(scheme_name, scheme_name.upper())
        resolved = theme.get("colors", {}).get(scheme_name)
        if resolved:
            return ColorSpec(value=resolved if resolved.startswith("#") else f"#{resolved}")
        return ColorSpec(value=f"scheme:{normalized}")
    return None


def _resolve_run_str_attr(attr_name: str, run: Any | None, sources: list[dict[str, Any]]) -> ResolvedValue:
    """Resolve a string rPr attribute from run level falling back through the source chain."""
    if run is not None:
        try:
            rPr = find_rPr(run._r)
            if rPr is not None:
                val = rPr.get(attr_name)
                if val is not None:
                    return ResolvedValue(val, "slide:run:rPr")
        except Exception:
            pass
    for source in sources:
        rPr = source.get("rPr")
        if rPr is not None:
            val = rPr.get(attr_name)
            if val is not None:
                return ResolvedValue(val, source["name"])
    return ResolvedValue(None, "unresolved")


def _resolve_run_int_attr(attr_name: str, run: Any | None, sources: list[dict[str, Any]]) -> ResolvedValue:
    """Resolve an int rPr attribute from run level falling back through the source chain."""
    if run is not None:
        try:
            rPr = find_rPr(run._r)
            if rPr is not None:
                val = rPr.get(attr_name)
                if val is not None:
                    try:
                        return ResolvedValue(int(val), "slide:run:rPr")
                    except (ValueError, TypeError):
                        pass
        except Exception:
            pass
    for source in sources:
        rPr = source.get("rPr")
        if rPr is not None:
            val = rPr.get(attr_name)
            if val is not None:
                try:
                    return ResolvedValue(int(val), source["name"])
                except (ValueError, TypeError):
                    pass
    return ResolvedValue(None, "unresolved")


def _resolve_para_line_spacing(paragraph: Any | None, sources: list[dict[str, Any]]) -> ResolvedValue:
    """Resolve line spacing from paragraph pPr, falling back to lstStyle sources."""
    if paragraph is not None:
        try:
            ls = paragraph.line_spacing
            if ls is not None:
                # pptx returns Pt(x) (a Length/int subclass) for fixed spcPts,
                # and a plain float for proportional spcPct.
                # Must check Length BEFORE (int, float) because Length inherits from int.
                from pptx.util import Length
                if isinstance(ls, Length):
                    return ResolvedValue(float(ls.pt), "slide:paragraph")
                if isinstance(ls, (int, float)):
                    return ResolvedValue(float(ls), "slide:paragraph")
        except Exception:
            pass
    for source in sources:
        pPr = source.get("pPr")
        if pPr is None:
            continue
        try:
            lnSpc = pPr.find(qa("lnSpc"))
            if lnSpc is None:
                continue
            spcPct = lnSpc.find(qa("spcPct"))
            if spcPct is not None:
                val = spcPct.get("val")
                if val is not None:
                    return ResolvedValue(int(val) / 100000.0, source["name"])
            spcPts = lnSpc.find(qa("spcPts"))
            if spcPts is not None:
                val = spcPts.get("val")
                if val is not None:
                    return ResolvedValue(int(val) / 100.0, source["name"])
        except Exception:
            pass
    return ResolvedValue(None, "unresolved")


def _resolve_para_spacing(tag_name: str, paragraph: Any | None, sources: list[dict[str, Any]]) -> ResolvedValue:
    """Resolve spcBef or spcAft from slide paragraph pPr then lstStyle sources. Returns pt value."""
    def _extract_spc(pPr: Any) -> float | None:
        if pPr is None:
            return None
        try:
            spcEl = pPr.find(qa(tag_name))
            if spcEl is None:
                return None
            spcPts = spcEl.find(qa("spcPts"))
            if spcPts is not None:
                val = spcPts.get("val")
                if val is not None:
                    return int(val) / 100.0  # hundredths of a point → points
            spcPct = spcEl.find(qa("spcPct"))
            if spcPct is not None:
                val = spcPct.get("val")
                if val is not None:
                    return float(val) / 100000.0  # store as fraction of line height
        except Exception:
            pass
        return None

    if paragraph is not None:
        try:
            pPr = find_pPr(paragraph._p)
            result = _extract_spc(pPr)
            if result is not None:
                return ResolvedValue(result, "slide:paragraph:pPr")
        except Exception:
            pass
    for source in sources:
        result = _extract_spc(source.get("pPr"))
        if result is not None:
            return ResolvedValue(result, source["name"])
    return ResolvedValue(None, "unresolved")


def _resolve_para_emu_attr(attr_name: str, paragraph: Any | None, sources: list[dict[str, Any]]) -> ResolvedValue:
    """Resolve a paragraph pPr EMU attribute (marL, indent) from slide then lstStyle sources.

    The ':paragraph' sources contain para[0]'s pPr from layout/master shapes. Only skip them
    when para[0] has an explicit buNone — meaning it is a non-bullet heading whose marL=0
    should NOT override lstStyle hanging-indent values for body/bullet paragraphs.
    """
    if paragraph is not None:
        try:
            pPr = find_pPr(paragraph._p)
            if pPr is not None:
                val = pPr.get(attr_name)
                if val is not None:
                    return ResolvedValue(int(val) / 914400.0, "slide:paragraph:pPr")
        except Exception:
            pass
    for source in sources:
        # Skip para[0] pPr from layout/master ONLY when it has buNone — a heading paragraph
        # with buNone+marL=0 must not shadow lstStyle bullet hanging-indent for body paragraphs.
        if source.get("name", "").endswith(":paragraph"):
            pPr_check = source.get("pPr")
            if pPr_check is not None and pPr_check.find(qa("buNone")) is not None:
                continue
        pPr = source.get("pPr")
        if pPr is not None:
            val = pPr.get(attr_name)
            if val is not None:
                try:
                    return ResolvedValue(int(val) / 914400.0, source["name"])
                except (ValueError, TypeError):
                    pass
    return ResolvedValue(None, "unresolved")


def _resolve_para_bullet_type(paragraph: Any | None, sources: list[dict[str, Any]]) -> ResolvedValue:
    """Resolve bullet type from slide paragraph pPr then lstStyle sources.

    The 'slide:paragraph' source is para[0]'s pPr — it's per-paragraph content,
    not a shape-wide template. Skip it for bullet resolution so that para[0]'s
    explicit buNone (e.g. a heading) doesn't shadow inherited bullets on para[1+].
    """
    from percy.oxml import bullet_type_from_pPr

    if paragraph is not None:
        try:
            pPr = find_pPr(paragraph._p)
            result = bullet_type_from_pPr(pPr)
            if result is not None:
                return ResolvedValue(result, "slide:paragraph:pPr")
        except Exception:
            pass
    for source in sources:
        # Skip "slide:paragraph" — it is para[0]'s per-paragraph pPr, not a
        # shape-wide default, and its buNone would incorrectly mask inherited bullets.
        if source.get("name") == "slide:paragraph":
            continue
        result = bullet_type_from_pPr(source.get("pPr"))
        if result is not None:
            return ResolvedValue(result, source["name"])
    return ResolvedValue("none", "office-default")


def _resolve_para_bullet_char(paragraph: Any | None, sources: list[dict[str, Any]]) -> ResolvedValue:
    """Resolve bullet character from slide paragraph pPr then lstStyle sources."""
    from percy.oxml import bullet_char_from_pPr

    if paragraph is not None:
        try:
            pPr = find_pPr(paragraph._p)
            result = bullet_char_from_pPr(pPr)
            if result is not None:
                return ResolvedValue(result, "slide:paragraph:pPr")
        except Exception:
            pass
    for source in sources:
        if source.get("name") == "slide:paragraph":
            continue  # skip para[0]'s pPr — not a shape-wide bullet template
        result = bullet_char_from_pPr(source.get("pPr"))
        if result is not None:
            return ResolvedValue(result, source["name"])
    return ResolvedValue(None, "unresolved")


def _resolve_para_bullet_font(paragraph: Any | None, sources: list[dict[str, Any]]) -> ResolvedValue:
    """Resolve bullet font typeface from slide paragraph pPr then lstStyle sources."""
    from percy.oxml import bullet_font_from_pPr

    if paragraph is not None:
        try:
            pPr = find_pPr(paragraph._p)
            result = bullet_font_from_pPr(pPr)
            if result is not None:
                return ResolvedValue(result, "slide:paragraph:pPr")
        except Exception:
            pass
    for source in sources:
        if source.get("name") == "slide:paragraph":
            continue  # skip para[0]'s pPr — not a shape-wide bullet template
        result = bullet_font_from_pPr(source.get("pPr"))
        if result is not None:
            return ResolvedValue(result, source["name"])
    return ResolvedValue(None, "unresolved")


def _def_rpr(element: Any) -> Any | None:
    if element is None:
        return None
    return _first(_xpath(element, "./a:defRPr"))


def _attr(element: Any, xpath: str, attr_name: str) -> str | None:
    if element is None:
        return None
    node = _first(_xpath(element, xpath))
    if node is None:
        return None
    return node.get(attr_name)


def _bool_attr(element: Any, attr_name: str) -> bool | None:
    if element is None or element.get(attr_name) is None:
        return None
    return element.get(attr_name) in {"1", "true", "True"}


def _first(items: list[Any]) -> Any | None:
    return items[0] if items else None


def _xpath(element: Any, xpath: str) -> list[Any]:
    if element is None:
        return []
    try:
        return element.xpath(xpath, namespaces=NS)
    except TypeError:
        return element.xpath(xpath)
