from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches

from percy.diagnostics.corpus import analyze_corpus


def test_analyze_corpus_summarizes_pptx_and_pdf(tmp_path) -> None:
    input_dir = tmp_path / "dump"
    input_dir.mkdir()
    _write_sample_deck(input_dir / "sample.pptx")
    (input_dir / "sample.pdf").write_bytes(b"%PDF-1.4\n")

    report = analyze_corpus(input_dir, tmp_path / "analysis")

    assert report["pptx_count"] == 1
    assert report["pdf_count"] == 1
    assert report["feature_counts"]["text_frame"] == 1
    assert (tmp_path / "analysis" / "corpus.json").exists()


def _write_sample_deck(path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    textbox.text = "Corpus sample"
    presentation.save(path)
