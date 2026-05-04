from __future__ import annotations

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx.util import Pt

from percy.bridge import BridgeShape, BridgeText, PercyDocument, load_percy
from percy.diagnostics.inspect import inspect_pptx
from percy.diagnostics.workflow import roundtrip_pptx


def test_inspect_pptx_reports_shapes(tmp_path) -> None:
    pptx_path = tmp_path / "sample.pptx"
    _write_sample_deck(pptx_path)

    report = inspect_pptx(pptx_path, tmp_path / "inspect")

    assert report["slide_count"] == 1
    assert report["slides"][0]["shape_count"] == 3
    assert (tmp_path / "inspect" / "inspection.json").exists()


def test_roundtrip_writes_percy_and_rebuilt_pptx(tmp_path) -> None:
    pptx_path = tmp_path / "sample.pptx"
    _write_sample_deck(pptx_path)

    report = roundtrip_pptx(pptx_path, tmp_path / "run", render=False)

    document = load_percy(report["percy_path"], PercyDocument)
    assert len(document.slides) == 1
    assert document.metadata.slide_width == 10.0
    assert document.metadata.slide_height == 7.5
    assert any(isinstance(element, BridgeText) for element in document.slides[0].elements)
    assert any(isinstance(element, BridgeShape) for element in document.slides[0].elements)
    assert (tmp_path / "run" / "sample.rebuilt.pptx").exists()
    assert (tmp_path / "run" / "roundtrip.json").exists()

    rebuilt = Presentation(tmp_path / "run" / "sample.rebuilt.pptx")
    rebuilt_shape = rebuilt.slides[0].shapes[1]
    assert rebuilt_shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
    assert rebuilt_shape.text == "Shape text"
    assert str(rebuilt_shape.fill.fore_color.rgb) == "CC3300"
    assert str(rebuilt_shape.line.color.rgb) == "003399"


def _write_sample_deck(path) -> None:
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank)

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    textbox.text = "Hello Percy"

    rounded = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(1), Inches(3), Inches(1))
    rounded.text = "Shape text"
    rounded.fill.solid()
    rounded.fill.fore_color.rgb = RGBColor(0xCC, 0x33, 0x00)
    rounded.line.color.rgb = RGBColor(0x00, 0x33, 0x99)
    rounded.line.width = Pt(2)

    table = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "2"

    presentation.save(path)
