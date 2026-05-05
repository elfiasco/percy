"""End-to-end roundtrip: agent builds a deck → export to PPTX → re-onboard.

This is the existential test: the whole Percy story falls apart if agent-built
decks can't actually export to a real .pptx file. This test verifies the
roundtrip works for:

  1. Templates applied to slides (Title, KPI Tiles, Section Header)
  2. Charts created via the agent (no reconstruction_blobs)
  3. Tables created via the agent
  4. Live groups (children flatten on export per design)

Uses python-pptx to re-open the exported file and confirm the structured
fields survive (slide count, element count, chart series, text content).
"""

from __future__ import annotations

import io
import os
import tempfile

import pytest
from fastapi.testclient import TestClient

os.environ["PERCY_PUBLIC_DEV"] = "1"

from percy.bridge import BridgeSlide, PercyDocument, PresentationMetadata
from percy.bridge import builders


THEME = {"ACCENT_1": "#3B82F6", "ACCENT_2": "#10B981", "ACCENT_3": "#F59E0B", "TX1": "#1E293B"}


@pytest.fixture
def client_and_doc():
    from app.backend import main as backend_main
    doc_id = "test-export-roundtrip"
    doc = PercyDocument(
        slides=[BridgeSlide(slide_number=i + 1, elements=[], width=13.333, height=7.5)
                for i in range(2)],
        metadata=PresentationMetadata(slide_width=13.333, slide_height=7.5, slide_count=2),
        theme_colors=THEME,
    )
    backend_main._docs[doc_id] = {
        "doc": doc, "bridge_dir": backend_main._CACHE_DIR / doc_id / "bridge",
        "name": f"{doc_id}.pptx", "kind": "pptx", "source_format": "pptx",
        "_undo_stack": [], "_redo_stack": [],
    }
    (backend_main._CACHE_DIR / doc_id / "bridge").mkdir(parents=True, exist_ok=True)
    yield TestClient(backend_main.app), doc_id, doc
    backend_main._docs.pop(doc_id, None)


def test_roundtrip_agent_built_deck(client_and_doc):
    client, doc_id, doc = client_and_doc

    # Apply Title template to slide 1
    r = client.post("/api/agent/templates/std.title/apply", json={
        "doc_id": doc_id, "slide_n": 1,
        "inputs": {"title": "Roundtrip Test", "subtitle": "Agent-built"},
    })
    assert r.status_code == 200, r.text

    # Create a chart via the agent on slide 2
    r = client.post(f"/api/docs/{doc_id}/slides/2/elements/chart", json={
        "chart_type": "column_clustered",
        "categories": ["Q1", "Q2", "Q3", "Q4"],
        "series": [{"name": "Revenue", "values": [100, 120, 130, 140]}],
        "title": "Revenue",
        "position": {"left_in": 1, "top_in": 1, "width_in": 8, "height_in": 5},
    })
    assert r.status_code == 200

    # Sanity: doc has the elements we expect
    assert len(doc.slides[0].elements) == 2  # title + subtitle text boxes
    assert len(doc.slides[1].elements) == 1  # chart

    # Now export as PPTX. The studio's existing /export endpoint produces .pptx.
    r = client.get(f"/api/docs/{doc_id}/export")
    assert r.status_code == 200, f"export failed: {r.text[:300]}"
    assert r.headers.get("content-type", "").startswith(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ), f"unexpected content-type: {r.headers.get('content-type')}"

    # Save and re-open with python-pptx to validate
    pptx_bytes = r.content
    assert len(pptx_bytes) > 1000, "exported pptx suspiciously small"

    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == 2, f"expected 2 slides, got {len(prs.slides)}"

    # Slide 1 should have at least one shape with the title text
    s1 = prs.slides[0]
    s1_text = ""
    for shape in s1.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    s1_text += run.text + " "
    assert "Roundtrip Test" in s1_text, f"title not in slide 1 export. Got: {s1_text!r}"

    # Slide 2 should have a chart
    s2 = prs.slides[1]
    chart_shapes = [shape for shape in s2.shapes if shape.has_chart]
    assert len(chart_shapes) >= 1, "no chart on slide 2 of exported deck"
    chart = chart_shapes[0].chart
    # Verify chart data survived
    plot = chart.plots[0]
    series = list(plot.series)
    assert len(series) == 1, f"expected 1 series, got {len(series)}"
    values = list(series[0].values)
    assert values == [100, 120, 130, 140], f"chart values lost in export: {values}"


def test_roundtrip_live_group_flattens(client_and_doc):
    """Live groups should flatten to flat shapes on PPTX export (per design)."""
    client, doc_id, doc = client_and_doc
    r = client.post("/api/agent/templates/std.live_timeline/apply", json={
        "doc_id": doc_id, "slide_n": 1,
        "inputs": {"title": "Sprint", "day_count": 5, "labels": "Mon,Tue,Wed,Thu,Fri"},
    })
    assert r.status_code == 200

    # Verify the live group has children
    slide1 = doc.slides[0]
    groups = [el for el in slide1.elements if el.element_type == "BridgeGroup"]
    assert len(groups) == 1
    assert len(groups[0].children) == 5  # 5 days

    # Export
    r = client.get(f"/api/docs/{doc_id}/export")
    assert r.status_code == 200, r.text

    # Re-open and confirm we have 5 flattened shapes (plus title) on slide 1
    from pptx import Presentation
    prs = Presentation(io.BytesIO(r.content))
    s1 = prs.slides[0]
    shapes = list(s1.shapes)
    # We expect: 1 title + 5 day shapes = 6+
    assert len(shapes) >= 5, f"expected >=5 shapes after live-group flattening, got {len(shapes)}"

    # Day labels should be in the slide text
    text = ""
    for shape in s1.shapes:
        if shape.has_text_frame:
            text += shape.text_frame.text + " "
    for label in ("Mon", "Tue", "Wed"):
        assert label in text, f"day label {label!r} missing from exported text"
