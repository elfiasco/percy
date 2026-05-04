"""
Percy Presentation Analyzer
============================
Quality-checks every file in dump_pptx/ to confirm it is a genuine
corporate-branded deck from the target company — not a template, tutorial,
or random low-quality file.

For each file:
  1. Renders the first slide/page as an image
        PDF   → PyMuPDF (page 1)
        PPTX  → embedded thumbnail (docProps/thumbnail.*) or fallback text check
  2. Sends the image to Gemma-3-27b (vision) for branding verification
  3. Extracts structural metadata with python-pptx (slide count, theme colors, etc.)
  4. Writes results into metadata.json under each file's record
  5. Moves files that fail the check to dump_pptx/rejected/

Usage:
    python analyzer.py                  # analyze all unanalyzed files
    python analyzer.py --recheck        # re-analyze everything (overwrites results)
    python analyzer.py --dry-run        # show verdicts without moving any files
    python analyzer.py --file foo.pptx  # analyze a single file by filename
"""

from __future__ import annotations

import argparse
import base64
import json
import logging
import re
import shutil
import zipfile
from datetime import datetime
from io import BytesIO
from pathlib import Path

import fitz  # PyMuPDF
from openai import OpenAI
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

BASE_DIR = Path(__file__).parent
DUMP_DIR = BASE_DIR / "dump_pptx"
REJECTED_DIR = DUMP_DIR / "rejected"
METADATA_PATH = BASE_DIR / "metadata.json"

REJECTED_DIR.mkdir(exist_ok=True)

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[logging.StreamHandler(), logging.FileHandler(BASE_DIR / "analyzer.log")],
)
log = logging.getLogger(__name__)


# ===========================================================================
# Gemma quality-check prompt
# ===========================================================================

BRANDING_PROMPT = """You are curating a training dataset of corporate presentation slides for an AI tool that replaces PowerPoint. You are reviewing the FIRST SLIDE of a presentation purportedly from {company_name}.

Your job: decide whether this is a genuine, high-quality {company_name} corporate presentation with their branded design system.

ACCEPT if: company logo/wordmark is visible, uses a consistent branded color palette (not default Office colors), professional layout, corporate content (investor metrics, product info, strategy, earnings, conference keynote).

REJECT if: looks like a blank or default Office template, is clearly not from {company_name} (wrong branding), is a tutorial/sample/placeholder deck, has very low design quality, or the content is unrelated to corporate/business materials.

Respond with a JSON object ONLY — no markdown fences, no extra text:
{{
  "verdict": "ACCEPT" or "REJECT",
  "confidence": 0.0-1.0,
  "has_visible_logo": true or false or null,
  "color_scheme": "brief description, e.g. 'navy blue + orange brand colors'",
  "design_quality": "high" or "medium" or "low",
  "content_type": "investor presentation" or "conference deck" or "earnings slides" or "product demo" or "template/sample" or "other",
  "rejection_reason": "reason if REJECT, empty string if ACCEPT",
  "observations": "one sentence describing what you see on the slide"
}}"""


# ===========================================================================
# Local LLM (Gemma via LM Studio)
# ===========================================================================

class LocalLLM:
    MODEL = "google/gemma-4-e4b"
    BASE_URL = "http://localhost:1234/v1"

    def __init__(self):
        self.client = OpenAI(base_url=self.BASE_URL, api_key="lm-studio")

    def assess_branding(self, img_b64: str, company_name: str) -> dict:
        """Send first-slide image to Gemma and get a structured quality verdict."""
        prompt = BRANDING_PROMPT.format(company_name=company_name)
        try:
            resp = self.client.chat.completions.create(
                model=self.MODEL,
                messages=[{"role": "user", "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_b64}"}},
                ]}],
                max_tokens=700,
                temperature=0.0,
            )
            raw = resp.choices[0].message.content.strip()
        except Exception as exc:
            log.warning("LLM vision call failed: %s", exc)
            return _unknown_result("LLM call failed")

        # Strip markdown code fences if present
        raw = re.sub(r"^```(?:json)?\s*", "", raw, flags=re.M)
        raw = re.sub(r"\s*```$", "", raw, flags=re.M)

        try:
            return json.loads(raw)
        except json.JSONDecodeError:
            # Try to extract verdict at minimum
            verdict = "ACCEPT" if '"ACCEPT"' in raw else ("REJECT" if '"REJECT"' in raw else "UNKNOWN")
            return _unknown_result(f"JSON parse failed — raw: {raw[:120]}", verdict=verdict)

    def assess_text_only(self, company_name: str, first_slide_text: str, slide_count: int) -> dict:
        """Fallback text-based check when no image is available (PPTX with no thumbnail)."""
        prompt = (
            f"A presentation file from '{company_name}' has {slide_count} slides. "
            f"The text on the first slide is:\n\n{first_slide_text[:600]}\n\n"
            "Based only on this text, does this look like an official corporate presentation "
            f"from {company_name} (investor deck, conference slides, earnings presentation)? "
            "Or does it look like a template, tutorial, or unrelated document?\n\n"
            "Reply with JSON ONLY:\n"
            '{"verdict": "ACCEPT" or "REJECT", "confidence": 0.0-1.0, '
            '"rejection_reason": "reason or empty", "observations": "one sentence"}'
        )
        try:
            resp = self.client.chat.completions.create(
                model=self.MODEL,
                messages=[{"role": "user", "content": prompt}],
                max_tokens=400,
                temperature=0.0,
            )
            raw = resp.choices[0].message.content.strip()
            raw = re.sub(r"^```(?:json)?\s*", "", raw, flags=re.M)
            raw = re.sub(r"\s*```$", "", raw, flags=re.M)
            return json.loads(raw)
        except Exception as exc:
            log.warning("Text-only LLM check failed: %s", exc)
            return _unknown_result("text LLM failed")


def _unknown_result(reason: str, verdict: str = "UNKNOWN") -> dict:
    return {
        "verdict": verdict,
        "confidence": 0.0,
        "has_visible_logo": None,
        "color_scheme": "",
        "design_quality": "low",
        "content_type": "other",
        "rejection_reason": reason,
        "observations": "",
    }


# ===========================================================================
# Rendering utilities
# ===========================================================================

def render_pdf_first_page(path: Path, max_width: int = 1280) -> bytes | None:
    """Render PDF page 1 to PNG bytes using PyMuPDF."""
    try:
        doc = fitz.open(str(path))
        page = doc[0]
        # Scale so longest side ~= max_width
        scale = min(max_width / max(page.rect.width, page.rect.height), 2.0)
        mat = fitz.Matrix(scale, scale)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        return pix.tobytes("png")
    except Exception as exc:
        log.warning("PDF render failed for %s: %s", path.name, exc)
        return None


def extract_pptx_thumbnail(path: Path) -> bytes | None:
    """Extract the embedded thumbnail from a PPTX file (stored in docProps/)."""
    try:
        with zipfile.ZipFile(path, "r") as z:
            for candidate in [
                "docProps/thumbnail.jpeg",
                "docProps/thumbnail.png",
                "docProps/thumbnail.jpg",
                "docProps/thumbnail.wmf",
            ]:
                if candidate in z.namelist():
                    data = z.read(candidate)
                    # WMF thumbnails aren't usable for vision — skip them
                    if candidate.endswith(".wmf"):
                        continue
                    return data
    except Exception as exc:
        log.debug("Thumbnail extract failed for %s: %s", path.name, exc)
    return None


def get_pptx_structural_info(path: Path) -> dict:
    """
    Extract structural metadata from a PPTX:
    slide count, theme colors, whether custom theme is present, first-slide text.
    """
    info = {
        "slide_count": 0,
        "master_slide_count": 0,
        "has_custom_theme": False,
        "theme_colors": [],
        "first_slide_text": "",
        "has_animations": False,
        "has_smartart": False,
        "has_embedded_excel": False,
    }
    try:
        prs = Presentation(str(path))
        info["slide_count"] = len(prs.slides)
        info["master_slide_count"] = len(prs.slide_masters)

        # Theme colors from the first slide master
        if prs.slide_masters:
            master = prs.slide_masters[0]
            try:
                theme = master.theme_color_map  # may not exist in all versions
                info["has_custom_theme"] = True
            except Exception:
                pass
            # Try to extract accent/brand colors from the theme XML
            try:
                theme_xml = master._element.xml
                color_matches = re.findall(r'val="([0-9A-Fa-f]{6})"', theme_xml)
                info["theme_colors"] = list(dict.fromkeys(color_matches))[:12]
                # If no default Office grays/blues, it's likely a custom theme
                default_office = {"4472C4", "ED7D31", "A9D18E", "FF0000"}
                info["has_custom_theme"] = not any(
                    c.upper() in default_office for c in info["theme_colors"]
                )
            except Exception:
                pass

        # First-slide text
        if prs.slides:
            slide = prs.slides[0]
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
                # Check for embedded OLE (Excel charts)
                if hasattr(shape, "shape_type") and shape.shape_type == 3:  # MSO_SHAPE_TYPE.OLE_OBJECT
                    info["has_embedded_excel"] = True

            info["first_slide_text"] = "\n".join(texts)

            # Animations (check slide XML for animClr/animEffect elements)
            try:
                slide_xml = slide._element.xml
                info["has_animations"] = "<p:animation" in slide_xml or "<p:anim" in slide_xml
            except Exception:
                pass

        # SmartArt — look for graphicData with SmartArt URI across all slides
        for slide in prs.slides:
            try:
                if "dgm:" in slide._element.xml or "SmartArt" in slide._element.xml:
                    info["has_smartart"] = True
                    break
            except Exception:
                pass

        # Animations across all slides
        if not info["has_animations"]:
            for slide in prs.slides:
                try:
                    if "<p:animation" in slide._element.xml or "<p:anim " in slide._element.xml:
                        info["has_animations"] = True
                        break
                except Exception:
                    pass

    except Exception as exc:
        log.warning("python-pptx parse failed for %s: %s", path.name, exc)

    return info


def to_b64_png(img_bytes: bytes, max_width: int = 1280) -> str:
    """Convert raw image bytes (any format) to a resized base64 PNG."""
    try:
        img = Image.open(BytesIO(img_bytes)).convert("RGB")
        if img.width > max_width:
            ratio = max_width / img.width
            img = img.resize((max_width, int(img.height * ratio)), Image.LANCZOS)
        buf = BytesIO()
        img.save(buf, format="PNG", optimize=True)
        return base64.b64encode(buf.getvalue()).decode()
    except Exception as exc:
        raise RuntimeError(f"Image conversion failed: {exc}") from exc


# ===========================================================================
# Per-file analysis
# ===========================================================================

def analyze_file(
    record: dict,
    llm: LocalLLM,
    dry_run: bool = False,
) -> dict:
    """
    Run the full quality check on a single downloaded file.
    Mutates `record` in place and returns it.
    """
    filename = record["filename"]
    company = record["company"]
    file_type = record.get("file_type", Path(filename).suffix.lstrip(".").lower())
    path = DUMP_DIR / filename

    log.info("Analyzing: %s  [%s]", filename, company)

    if not path.exists():
        log.warning("  File missing: %s", path)
        record["quality_check"] = {
            "passed": False,
            "checked_at": datetime.utcnow().isoformat(),
            "verdict": "MISSING",
            "rejection_reason": "file not found on disk",
        }
        return record

    # ----- Step 1: get image representation --------------------------------
    img_b64: str | None = None
    structural: dict = {}
    used_vision = False

    if file_type == "pdf":
        png_bytes = render_pdf_first_page(path)
        if png_bytes:
            img_b64 = to_b64_png(png_bytes)
            used_vision = True
        else:
            log.warning("  Could not render PDF: %s", filename)

    elif file_type in ("pptx", "ppt"):
        structural = get_pptx_structural_info(path)

        # Update formatting_analysis with structural data
        fa = record.setdefault("formatting_analysis", {})
        fa["slide_count"] = structural.get("slide_count")
        fa["master_slide_count"] = structural.get("master_slide_count")
        fa["has_animations"] = structural.get("has_animations")
        fa["has_smartart"] = structural.get("has_smartart")
        fa["has_embedded_excel"] = structural.get("has_embedded_excel")
        fa["has_custom_fonts"] = structural.get("has_custom_theme")  # proxy
        fa["color_palette"] = structural.get("theme_colors", [])

        thumb_bytes = extract_pptx_thumbnail(path)
        if thumb_bytes:
            try:
                img_b64 = to_b64_png(thumb_bytes)
                used_vision = True
            except RuntimeError as exc:
                log.debug("  Thumbnail unusable: %s", exc)

    # ----- Step 2: Gemma quality check ------------------------------------
    if img_b64 and used_vision:
        log.info("  Sending to Gemma (vision)...")
        gemma_result = llm.assess_branding(img_b64, company)
    elif file_type in ("pptx", "ppt") and structural.get("first_slide_text"):
        log.info("  No thumbnail — using text-based Gemma check...")
        gemma_result = llm.assess_text_only(
            company,
            structural["first_slide_text"],
            structural.get("slide_count", 0),
        )
    else:
        log.warning("  No image or text available for %s — skipping LLM check", filename)
        gemma_result = _unknown_result("no renderable content")

    verdict = gemma_result.get("verdict", "UNKNOWN")
    confidence = gemma_result.get("confidence", 0.0)
    passed = verdict == "ACCEPT" and confidence >= 0.5

    log.info(
        "  Verdict: %s  (confidence=%.2f)  quality=%s  type=%s",
        verdict,
        confidence,
        gemma_result.get("design_quality", "?"),
        gemma_result.get("content_type", "?"),
    )
    if not passed:
        log.info("  Rejection reason: %s", gemma_result.get("rejection_reason", ""))

    # ----- Step 3: persist results ----------------------------------------
    record["quality_check"] = {
        "passed": passed,
        "checked_at": datetime.utcnow().isoformat(),
        "used_vision": used_vision,
        **gemma_result,
    }
    record["status"] = "verified" if passed else "rejected"

    # Update notes field with Gemma's observations
    if gemma_result.get("observations") and not record.get("notes"):
        record["notes"] = gemma_result["observations"]

    # ----- Step 4: move rejected files ------------------------------------
    if not passed and not dry_run:
        dest = REJECTED_DIR / filename
        try:
            shutil.move(str(path), str(dest))
            log.info("  Moved to rejected/: %s", filename)
        except Exception as exc:
            log.warning("  Could not move %s: %s", filename, exc)

    return record


# ===========================================================================
# Main run loop
# ===========================================================================

def run(recheck: bool = False, dry_run: bool = False, only_file: str | None = None) -> None:
    with open(METADATA_PATH, encoding="utf-8") as f:
        metadata = json.load(f)

    llm = LocalLLM()
    records = metadata.get("scraped_files", [])

    if only_file:
        records = [r for r in records if r["filename"] == only_file]
        if not records:
            log.error("File not found in metadata: %s", only_file)
            return

    to_process = [
        r for r in records
        if recheck or "quality_check" not in r
    ]

    log.info("Files to analyze: %d  (recheck=%s, dry_run=%s)", len(to_process), recheck, dry_run)

    passed = rejected = errors = 0
    for record in to_process:
        try:
            analyze_file(record, llm, dry_run=dry_run)
            if record.get("quality_check", {}).get("passed"):
                passed += 1
            else:
                rejected += 1
        except Exception as exc:
            log.exception("Unhandled error on %s: %s", record.get("filename"), exc)
            errors += 1

    if not dry_run:
        metadata["last_updated"] = datetime.utcnow().strftime("%Y-%m-%d")
        with open(METADATA_PATH, "w", encoding="utf-8") as f:
            json.dump(metadata, f, indent=2)
        log.info("metadata.json updated.")

    log.info(
        "Analysis complete — PASSED: %d | REJECTED: %d | ERRORS: %d | TOTAL: %d",
        passed, rejected, errors, len(to_process),
    )

    # Print a summary table
    print("\n" + "=" * 70)
    print(f"{'FILE':<45} {'VERDICT':<8} {'CONF':>5}  {'TYPE'}")
    print("-" * 70)
    for r in to_process:
        qc = r.get("quality_check", {})
        v = qc.get("verdict", "?")
        c = qc.get("confidence", 0)
        t = qc.get("content_type", "")
        name = r["filename"][:44]
        print(f"{name:<45} {v:<8} {c:>5.2f}  {t}")
    print("=" * 70)


# ===========================================================================
# CLI
# ===========================================================================

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Percy Presentation Quality Analyzer")
    parser.add_argument("--recheck", action="store_true",
                        help="Re-analyze all files, even already-checked ones")
    parser.add_argument("--dry-run", action="store_true",
                        help="Show verdicts but don't move any files or write metadata")
    parser.add_argument("--file", metavar="FILENAME",
                        help="Analyze a single file by filename (must exist in metadata)")
    args = parser.parse_args()
    run(recheck=args.recheck, dry_run=args.dry_run, only_file=args.file)
